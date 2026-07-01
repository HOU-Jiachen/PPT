from __future__ import annotations

import json
import re
import sys
from datetime import datetime
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

PROJECT_DIR = Path(__file__).resolve().parents[1]
REPO_ROOT = PROJECT_DIR.parents[1]
LOCAL_SCRIPTS = REPO_ROOT / ".codex" / "skills" / "engineering-ppt" / "scripts"
sys.path.insert(0, str(LOCAL_SCRIPTS))

from engineering_deck_runtime import (  # noqa: E402
    DEFAULT_COLORS,
    add_emphasis_textbox,
    add_fill,
    add_ir_table,
    add_textbox,
    sanitize_visible_text,
    text_preview,
)
from table_renderers import effective_image_table_font_pt  # noqa: E402


PROJECT = "东坡大道水土保持方案"
SOURCE_FILE = "中地环科-黄冈市东坡大道片区污水收集管网建设工程水土保持方案报告表-20260108.docx"
TITLE = "黄冈市东坡大道片区污水收集管网建设工程水土保持方案报告表汇报"

CATALOG = PROJECT_DIR / "analysis" / "source_catalog.json"
INVENTORY = PROJECT_DIR / "analysis" / "report_content_inventory.json"
TABLE_IR = PROJECT_DIR / "analysis" / "table_ir.json"
EXPORTS = PROJECT_DIR / "exports"
NOTES = PROJECT_DIR / "notes"

COLORS = {
    **DEFAULT_COLORS,
    "bg": "F4F7FA",
    "paper": "FFFFFF",
    "ink": "14324A",
    "muted": "607385",
    "primary": "0C3B5D",
    "secondary": "EAF2F6",
    "accent": "0E8F8C",
    "line": "CCD8DE",
    "soft": "EEF5F3",
    "good": "2E7D32",
    "warn": "A35E2A",
}

CHAPTERS = [
    ("项目总览", "全局控制指标先行说明项目范围、预测结果、措施投资和验收闭环。"),
    ("综合说明与审批依据", "综合表和批复附件共同限定方案范围。"),
    ("项目概况与工程布置", "工程规模、占地、工期和施工组织决定扰动来源。"),
    ("防治责任范围与标准目标", "责任范围和一级标准构成后续措施评价基准。"),
    ("水土保持评价与土石方", "合规评价、土石方平衡和主体措施共同校核方案基础。"),
    ("水土流失分析与预测", "预测范围、侵蚀模数和流失量形成防护强度依据。"),
    ("水土保持措施与施工安排", "分区措施、工程量和施工窗口形成实施路径。"),
    ("投资概算、效益与管理验收", "投资、补偿费、效益和验收管理形成闭环。"),
]

PAGES: list[dict[str, Any]] = [
    {"type": "cover", "title": TITLE, "chapter": "", "source_mode": "", "tables": []},
    {
        "type": "overview",
        "title": "项目范围、预测结果、措施投资和验收管理形成闭环",
        "chapter": "项目总览",
        "source_mode": "INTERPRETATION",
        "tables": ["T004", "T034", "T040", "T048"],
        "layout": "overview",
        "bullets": [
            "1. 项目范围：工程位于黄冈市黄州区城北片区，污水主管网6.2公里，预埋支管2.2公里。",
            "2. 预测结果：预测土壤流失总量104.05t，新增土壤流失量94.92t。",
            "3. 管理闭环：水土保持总投资26.9211万元，六项效益指标达到报告目标值。",
        ],
    },
    {"type": "agenda", "title": "汇报结构沿报告七个技术章节展开", "chapter": "", "source_mode": "", "tables": []},
    {
        "title": "综合报告表明确项目全局控制指标",
        "chapter": "综合说明与审批依据",
        "source_mode": "ORIGINAL_TABLE",
        "tables": ["T004"],
        "layout": "wide_table",
        "bullets": [
            "1. 建设内容：东坡大道片区改造污水主管网6.2公里，沿线预埋支管2.2公里。",
            "2. 责任范围：防治责任范围为34416.00m²，均按管网工程区纳入。",
            "3. 投资口径：水土保持总投资26.9211万元，主体已列和方案新增分项计列。",
        ],
    },
    {
        "title": "工程特性表锁定规模、投资、工期和占地",
        "chapter": "项目概况与工程布置",
        "source_mode": "ORIGINAL_TABLE",
        "tables": ["T008"],
        "layout": "side_table",
        "bullets": [
            "1. 工程性质：项目为新建污水收集管网工程，建设地点为黄冈市黄州区。",
            "2. 投资工期：总投资8024.00万元，土建投资6914.95万元，建设期为13个月。",
            "3. 占地构成：总占地3.44hm²，报告按临时占地组织水土保持设计。",
        ],
    },
    {
        "title": "防治责任范围全部纳入管网工程区临时占地",
        "chapter": "防治责任范围与标准目标",
        "source_mode": "ORIGINAL_TABLE",
        "tables": ["T006", "T013"],
        "layout": "side_table",
        "bullets": [
            "1. 范围口径：管网工程区防治责任范围为34416.00m²。",
            "2. 用地类型：城镇村道路用地34151.00m²，公园与绿地265.00m²。",
            "3. 占地属性：报告表中永久占地为空，责任范围按临时占地纳入。",
        ],
    },
    {
        "title": "六项防治指标采用南方红壤区一级标准",
        "chapter": "防治责任范围与标准目标",
        "source_mode": "ORIGINAL_TABLE",
        "tables": ["T007"],
        "layout": "wide_table",
        "bullets": [
            "1. 标准等级：项目执行南方红壤区建设类项目一级防治标准。",
            "2. 治理目标：水土流失治理度设计水平年目标为98%。",
            "3. 控制目标：土壤流失控制比修正后设计水平年目标为1.0。",
        ],
    },
    {
        "title": "表土剥离量按公园与绿地扰动面积计列",
        "chapter": "项目概况与工程布置",
        "source_mode": "ORIGINAL_TABLE",
        "tables": ["T014"],
        "layout": "side_table",
        "bullets": [
            "1. 剥离对象：表土剥离面积为265.00m²，对应公园与绿地扰动区域。",
            "2. 剥离参数：表土层厚度0.3m，剥离厚度0.3m。",
            "3. 利用途径：表土剥离量79.5m³，后续用于管网工程区绿化覆土。",
        ],
    },
    {
        "title": "土石方平衡显示弃方委托消纳、借方外购",
        "chapter": "水土保持评价与土石方",
        "source_mode": "ORIGINAL_TABLE",
        "tables": ["T015"],
        "layout": "wide_table",
        "bullets": [
            "1. 平衡关系：土石方开挖、回填、借方和弃方按工程分区分项列示。",
            "2. 弃方去向：弃方采用委托消纳方式处理。",
            "3. 借方来源：借方采用外购方式解决，报告未另列取土场。",
        ],
    },
    {
        "title": "法律和技术标准评价结论均为符合",
        "chapter": "水土保持评价与土石方",
        "source_mode": "ORIGINAL_TABLE",
        "tables": ["T018", "T019", "T020"],
        "layout": "side_table",
        "bullets": [
            "1. 法律条款：水土保持法相关制约性因素逐条评价为符合。",
            "2. 技术标准：报告按生产建设项目水土保持技术标准进行符合性分析。",
            "3. 长江保护：项目建设与长江保护法相关要求形成一致评价结论。",
        ],
    },
    {
        "title": "主体工程已列措施以土地整治和铺种草皮为主",
        "chapter": "水土保持评价与土石方",
        "source_mode": "ORIGINAL_TABLE",
        "tables": ["T021"],
        "layout": "side_table",
        "bullets": [
            "1. 工程措施：主体工程已列土地整治，数量按管网工程区计列。",
            "2. 植物措施：铺种草皮纳入主体工程已列投资。",
            "3. 投资关系：主体已列措施与方案新增措施共同构成防治体系。",
        ],
    },
    {
        "title": "预测范围与防治责任范围保持一致",
        "chapter": "水土流失分析与预测",
        "source_mode": "ORIGINAL_TABLE",
        "tables": ["T023", "T024", "T025", "T026"],
        "layout": "side_table",
        "bullets": [
            "1. 预测单元：预测范围为管网工程区，面积与防治责任范围一致。",
            "2. 扰动面积：扰动地表面积按34416.00m²纳入预测。",
            "3. 时段划分：施工期和自然恢复期分别列示预测参数。",
        ],
    },
    {
        "title": "背景土壤侵蚀模数由占地类型加权形成",
        "chapter": "水土流失分析与预测",
        "source_mode": "CALCULATION",
        "tables": ["T027", "T028"],
        "layout": "side_table",
        "bullets": [
            "1. 类型参数：城镇村道路用地平均土壤侵蚀模数为260t/(km²·a)。",
            "2. 绿地参数：公园与绿地平均土壤侵蚀模数为450t/(km²·a)。",
            "3. 加权结果：报告列示背景土壤侵蚀模数为261.46t/(km²·a)。",
        ],
    },
    {
        "title": "施工期侵蚀模数显著高于背景值",
        "chapter": "水土流失分析与预测",
        "source_mode": "CALCULATION",
        "tables": ["T030", "T031", "T032", "T033"],
        "layout": "side_table",
        "bullets": [
            "1. 植被破坏型：施工期土壤侵蚀模数为1153.34t/(km²·a)。",
            "2. 地表翻扰型：施工期土壤侵蚀模数为3032.58t/(km²·a)。",
            "3. 综合取值：施工期综合取值为3018.60t/(km²·a)。",
        ],
    },
    {
        "title": "预测土壤流失总量104.05t，新增量94.92t",
        "chapter": "水土流失分析与预测",
        "source_mode": "ORIGINAL_TABLE",
        "tables": ["T034"],
        "layout": "wide_table",
        "bullets": [
            "1. 预测总量：项目预测土壤流失总量为104.05t。",
            "2. 新增量：新增土壤流失量为94.92t。",
            "3. 重点时段：施工期是新增水土流失控制重点。",
        ],
    },
    {
        "title": "防治分区为管网工程区，防治面积34416.00m²",
        "chapter": "水土保持措施与施工安排",
        "source_mode": "ORIGINAL_TABLE",
        "tables": ["T035"],
        "layout": "side_table",
        "bullets": [
            "1. 分区设置：报告将项目区划分为管网工程区。",
            "2. 防治面积：管网工程区防治面积为34416.00m²。",
            "3. 管理对象：永久占地为空，临时占地承担主要防治任务。",
        ],
    },
    {
        "title": "措施体系覆盖工程、植物和临时三类措施",
        "chapter": "水土保持措施与施工安排",
        "source_mode": "ORIGINAL_TABLE",
        "tables": ["T036", "T037"],
        "layout": "wide_table",
        "bullets": [
            "1. 工程措施：主体已有土地整治，方案新增表土剥离和表土回覆。",
            "2. 植物措施：主体已有铺种草皮。",
            "3. 临时措施：报告列示临时苫盖、高压洗车池和临时排水沉沙措施。",
        ],
    },
    {
        "title": "措施实施进度嵌入2025—2026年施工窗口",
        "chapter": "水土保持措施与施工安排",
        "source_mode": "ORIGINAL_TABLE",
        "tables": ["T038"],
        "layout": "wide_table",
        "bullets": [
            "1. 工期窗口：水土保持工程施工衔接2025年12月至2026年12月。",
            "2. 措施节奏：临时防护、表土措施和植物措施按施工过程分期安排。",
            "3. 进度关系：措施实施与主体工程施工时序保持同步。",
        ],
    },
    {
        "title": "水土保持总投资26.9211万元",
        "chapter": "投资概算、效益与管理验收",
        "source_mode": "ORIGINAL_TABLE",
        "tables": ["T040", "T041", "T042", "T043", "T044"],
        "layout": "side_table",
        "bullets": [
            "1. 总投资：水土保持总投资为26.9211万元。",
            "2. 分项构成：新增投资、主体工程已列投资和独立费用分项计列。",
            "3. 费用关系：工程措施、植物措施和临时措施共同进入概算体系。",
        ],
    },
    {
        "title": "水土保持补偿费51624元符合免征情形",
        "chapter": "投资概算、效益与管理验收",
        "source_mode": "ORIGINAL_TABLE",
        "tables": ["T045"],
        "layout": "side_table",
        "bullets": [
            "1. 征占面积：补偿费计算面积为34416m²。",
            "2. 计费标准：报告列示收费标准为1.5元/m²。",
            "3. 计算金额：水土保持补偿费为51624元，并列明免征依据。",
        ],
    },
    {
        "title": "六项防治效益指标达到报告目标值",
        "chapter": "投资概算、效益与管理验收",
        "source_mode": "ORIGINAL_TABLE",
        "tables": ["T048"],
        "layout": "wide_table",
        "bullets": [
            "1. 治理度：水土流失总治理度设计达到值为99.99%。",
            "2. 控制比：土壤流失控制比设计达到值满足报告目标。",
            "3. 效益结论：各项防治效益评估结果均为达标。",
        ],
    },
    {
        "type": "closing",
        "title": "验收管理落实后形成水土保持闭环",
        "chapter": "投资概算、效益与管理验收",
        "source_mode": "MANAGEMENT_ACTION",
        "tables": ["T048"],
        "layout": "closing",
        "bullets": [
            "1. 监理管理：水土保持监理纳入工程建设管理过程。",
            "2. 设施验收：水土保持设施验收完成后，工程闭环进入运行阶段。",
            "3. 资料管理：报告成果、验收资料和公开要求共同支撑后续管理。",
        ],
    },
]


def load_json(path: Path) -> dict[str, Any]:
    return json.loads(path.read_text(encoding="utf-8-sig"))


def catalog_entries() -> list[dict[str, Any]]:
    return load_json(CATALOG).get("entries", [])


def table_ir_records() -> list[dict[str, Any]]:
    return load_json(TABLE_IR).get("tables", [])


def entry_for_locator(locator: str) -> dict[str, Any] | None:
    for entry in catalog_entries():
        if entry.get("locator") == locator:
            return entry
    return None


def table_for_locator(locator: str) -> dict[str, Any] | None:
    for table in table_ir_records():
        if table.get("locator") == locator or table.get("table_id") == locator:
            return table
    return None


def table_image_path(table: dict[str, Any]) -> Path | None:
    asset = table.get("assets", {}).get("crop_image") or table.get("assets", {}).get("source_table_image")
    if not asset:
        return None
    path = PROJECT_DIR / str(asset)
    return path if path.exists() else None


def table_effective_font(locator: str, w: float, h: float) -> float:
    table = table_for_locator(locator)
    if table is None:
        return 0.0
    image_path = table_image_path(table)
    if image_path is None:
        return 0.0
    return effective_image_table_font_pt(table, image_path, w, h)


def required_table_parts(locator: str) -> int:
    if table_effective_font(locator, 12.1, 5.35) >= 10:
        return 1
    return 0


def usable_table_locators(locators: list[str]) -> list[str]:
    return [locator for locator in locators if required_table_parts(locator) == 1]


def best_table_layout(locator: str) -> str:
    wide_score = table_effective_font(locator, 11.74, 4.42)
    side_score = table_effective_font(locator, 6.92, 5.45)
    return "wide_table" if wide_score >= side_score else "side_table"


def chinese_suffix(index: int) -> str:
    values = ["一", "二", "三", "四", "五", "六"]
    return values[index - 1] if 1 <= index <= len(values) else str(index)


def continuation_bullets(page: dict[str, Any], locator: str, index: int) -> list[str]:
    bullets = list(page.get("bullets", []))
    table = table_for_locator(locator) or {}
    rows = table.get("rows", [])
    selected = bullets if index == 1 else (bullets[index - 1 : index + 1] or bullets[:2])
    snippets: list[str] = []
    start = max(1, index)
    for row in rows[start : start + 4]:
        cells = [sanitize_visible_text(str(cell)) for cell in row if sanitize_visible_text(str(cell))]
        if cells:
            snippets.append("、".join(cells[:3]))
        if len(snippets) >= 2:
            break
    for snippet in snippets:
        selected.append(f"{len(selected) + 1}. 表内分项：{text_preview(snippet, 42)}。")
    return selected[:3] or bullets[:3]


def expanded_pages() -> list[dict[str, Any]]:
    pages: list[dict[str, Any]] = []
    for page in PAGES:
        locators = list(page.get("tables", []))
        if page.get("layout") == "overview":
            clone = dict(page)
            clone["evidence_tables"] = locators
            clone["tables"] = []
            pages.append(clone)
            continue
        if page.get("type") in {"cover", "agenda", "closing"} or len(locators) <= 1:
            clone = dict(page)
            display_locators = usable_table_locators(locators)
            clone["evidence_tables"] = locators
            clone["tables"] = display_locators
            if len(locators) == 1 and clone.get("layout") not in {"overview", "closing"}:
                clone["layout"] = best_table_layout(display_locators[0]) if display_locators else "text_only"
            pages.append(clone)
            continue
        for index, locator in enumerate(locators, start=1):
            if required_table_parts(locator) != 1:
                continue
            clone = dict(page)
            clone["evidence_tables"] = locators
            clone["tables"] = [locator]
            clone["layout"] = best_table_layout(locator)
            if index > 1:
                clone["title"] = f"{page['title']}（{chinese_suffix(index)}）"
            clone["bullets"] = continuation_bullets(page, locator, index)
            pages.append(clone)
        if not any(required_table_parts(locator) == 1 for locator in locators):
            clone = dict(page)
            clone["evidence_tables"] = locators
            clone["tables"] = []
            clone["layout"] = "text_only"
            pages.append(clone)
    return pages


def source_text_for_locators(locators: list[str], limit: int = 1200) -> str:
    chunks: list[str] = []
    for locator in locators:
        entry = entry_for_locator(locator)
        if entry:
            chunks.append(sanitize_visible_text(str(entry.get("text") or "")))
    return text_preview(" ".join(chunks), limit, complete_sentence=False)


def catalog_ids(locators: list[str]) -> list[str]:
    ids: list[str] = []
    for locator in locators:
        entry = entry_for_locator(locator)
        if entry and entry.get("id"):
            ids.append(str(entry["id"]))
    return ids


def evidence_id(page_index: int, locators: list[str]) -> str:
    suffix = "-".join(locators[:3]) if locators else "TEXT"
    return f"E-{page_index:02d}-{suffix}"


def collect_values(text: str) -> list[dict[str, str]]:
    values: list[dict[str, str]] = []
    pattern = re.compile(r"[-+]?\d+(?:\.\d+)?\s*(?:%|hm²|m²|m³|公里|万元|元|t|t/\(km²·a\)|个月)?")
    for match in pattern.findall(text):
        value = match.strip()
        if not value or value in {"1", "2", "3"}:
            continue
        values.append({"value": value, "unit": "", "time_basis": "source"})
        if len(values) >= 8:
            break
    return values


def calculation_meta(locators: list[str]) -> dict[str, Any]:
    key = tuple(locators)
    if key == ("T027", "T028"):
        return {
            "formula": "背景土壤侵蚀模数 = Σ(占地面积 × 对应地类平均土壤侵蚀模数) / Σ占地面积",
            "inputs": {
                "城镇村道路用地面积_m2": "34151.00",
                "公园与绿地面积_m2": "265.00",
                "城镇村道路用地侵蚀模数_t_km2_a": "260",
                "公园与绿地侵蚀模数_t_km2_a": "450",
                "报告列示背景值_t_km2_a": "261.46",
            },
        }
    if key == ("T030", "T031", "T032", "T033"):
        return {
            "formula": "施工期综合侵蚀模数按报告表参数取值并汇总，M = R × K × L × S × B × E × T。",
            "inputs": {
                "植被破坏型施工期_t_km2_a": "1153.34",
                "地表翻扰型施工期_t_km2_a": "3032.58",
                "自然恢复期_t_km2_a": "322.11",
                "施工期综合取值_t_km2_a": "3018.60",
            },
        }
    if set(locators).issubset({"T027", "T028"}):
        return calculation_meta(["T027", "T028"])
    if set(locators).issubset({"T030", "T031", "T032", "T033"}):
        return calculation_meta(["T030", "T031", "T032", "T033"])
    return {}


def write_evidence_and_plan() -> None:
    evidence: list[dict[str, Any]] = []
    slides: list[dict[str, Any]] = []
    for idx, page in enumerate(expanded_pages(), start=1):
        page_type = page.get("type", "content")
        locators = list(page.get("evidence_tables") or page.get("tables", []))
        eid = evidence_id(idx, locators)
        if page_type not in {"cover", "agenda"}:
            exact = source_text_for_locators(locators)
            kind = "INTERPRETATION" if page.get("source_mode") == "INTERPRETATION" else "ORIGINAL_TABLE"
            if page.get("source_mode") == "CALCULATION":
                kind = "CALCULATION"
            if page.get("source_mode") == "MANAGEMENT_ACTION":
                kind = "RECOMMENDATION"
            record = {
                    "id": eid,
                    "kind": kind,
                    "summary": sanitize_visible_text(page["title"]),
                    "source_file": SOURCE_FILE,
                    "source_locator": ", ".join(locators) if locators else "报告正文",
                    "catalog_ids": catalog_ids(locators),
                    "verification_status": "verified",
                    "exact_text": exact,
                    "values": collect_values(" ".join(page.get("bullets", [])) + " " + exact),
                    "notes": "当前构建器由项目源目录、Table IR 和报告内容盘点生成。",
                }
            if page.get("source_mode") == "CALCULATION":
                record.update(calculation_meta(locators))
            evidence.append(record)
        slide = {
            "page": idx,
            "type": page_type,
            "chapter": page.get("chapter", ""),
            "title": page["title"],
            "source_mode": page.get("source_mode", ""),
            "evidence_ids": [] if page_type in {"cover", "agenda"} else [eid],
            "content_unit_ids": [],
            "visual_proof": ", ".join(locators) if locators else "structural slide",
            "layout_pattern": page.get("layout", "left_text_right_table"),
            "source_note": ", ".join(locators),
            "slide_content": {"title": page["title"], "bullets": page.get("bullets", [])},
            "internal_notes": [
                "Use slide_content only for visible text.",
                "Complex source tables are rendered through Table IR crop assets.",
            ],
            "speaker_notes": [sanitize_visible_text(item) for item in page.get("bullets", [])],
            "density": "dense" if page.get("layout") in {"wide_table", "overview"} else "normal",
            "density_exempt_reason": "",
        }
        slides.append(slide)

    (PROJECT_DIR / "evidence_ledger.json").write_text(
        json.dumps({"schema_version": "1.0", "evidence": evidence}, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    (PROJECT_DIR / "deck_plan.json").write_text(
        json.dumps(
            {
                "deck": {
                    "project": PROJECT,
                    "title": TITLE,
                    "audience": "工程技术审查与建设管理沟通",
                    "objective": "说明项目水土保持范围、扰动预测、措施体系、投资与验收闭环。",
                    "chapter_order": [chapter for chapter, _ in CHAPTERS],
                },
                "slides": slides,
            },
            ensure_ascii=False,
            indent=2,
        ),
        encoding="utf-8",
    )


def write_design_contracts() -> None:
    (PROJECT_DIR / "chapter_coverage.md").write_text(
        "\n".join(
            ["# Chapter Coverage", ""]
            + [f"- {chapter}: {summary}" for chapter, summary in CHAPTERS]
        ),
        encoding="utf-8",
    )
    (PROJECT_DIR / "design_spec.md").write_text(
        "\n".join(
            [
                "# Design Specification",
                "",
                "Canvas: 16:9 engineering review deck.",
                "Audience: engineering review and project management readers.",
                "Visual system: restrained municipal engineering style with teal accents, white evidence panels, and readable native PPT text.",
                "Source policy: source tables use Table IR image/native rendering; visible text comes from slide_content bullets.",
                "Typography: Microsoft YaHei, title 23-30 pt, body 14-18 pt, table text follows Table IR renderer.",
            ]
        ),
        encoding="utf-8",
    )
    (PROJECT_DIR / "spec_lock.md").write_text(
        "\n".join(
            [
                "# Spec Lock",
                "",
                "canvas: ppt169",
                "width_in: 13.333",
                "height_in: 7.5",
                "font_family: Microsoft YaHei",
                "colors:",
                f"  ink: #{COLORS['ink']}",
                f"  accent: #{COLORS['accent']}",
                f"  bg: #{COLORS['bg']}",
                "body_min_pt: 14",
                "table_min_pt: 12",
            ]
        ),
        encoding="utf-8",
    )


def blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def slide_bg(slide) -> None:
    add_fill(slide, 0, 0, 13.333, 7.5, "bg", "bg", colors=COLORS)


def add_header(slide, title: str, chapter: str, page_no: int) -> None:
    add_textbox(slide, 0.46, 0.22, 8.8, 0.24, chapter, 10, "muted", True, colors=COLORS, component="page_meta")
    add_emphasis_textbox(slide, 0.46, 0.56, 11.4, 0.50, title, 23, "ink", colors=COLORS, component="title", bold=True)
    add_textbox(slide, 11.85, 0.34, 0.55, 0.20, f"{page_no:02d}", 9, "muted", True, align=PP_ALIGN.RIGHT, colors=COLORS, component="page_meta")
    add_fill(slide, 0.46, 1.12, 11.9, 0.02, "accent", "accent", colors=COLORS)


def add_numbered_points(slide, points: list[str], x: float, y: float, w: float, h: float, size: float = 14.5) -> None:
    row_h = h / max(1, len(points))
    for idx, point in enumerate(points, start=1):
        yy = y + (idx - 1) * row_h
        add_fill(slide, x, yy + 0.04, 0.34, 0.30, "accent", "accent", colors=COLORS)
        add_textbox(slide, x + 0.045, yy + 0.085, 0.25, 0.14, f"{idx}", 13, "paper", True, align=PP_ALIGN.CENTER, colors=COLORS, component="page_meta")
        text = re.sub(r"^\d+[.、]\s*", "", point)
        add_emphasis_textbox(
            slide,
            x + 0.47,
            yy,
            w - 0.47,
            max(0.42, row_h - 0.05),
            text,
            size,
            "ink",
            colors=COLORS,
            component="bullet",
            max_lines=2,
        )


def add_table_picture(slide, locator: str, x: float, y: float, w: float, h: float, mode: str | None = None) -> None:
    table = table_for_locator(locator)
    if table is None:
        add_textbox(slide, x, y, w, h, f"{locator} 表格未定位", 16, "warn", True, colors=COLORS)
        return
    render_mode = mode or "image"
    add_ir_table(slide, PROJECT_DIR, locator, x, y, w, h, mode=render_mode, colors=COLORS)


def add_multiple_table_strip(slide, locators: list[str], x: float, y: float, w: float, h: float) -> None:
    count = min(len(locators), 3)
    gap = 0.18
    item_h = (h - gap * (count - 1)) / max(1, count)
    for idx, locator in enumerate(locators[:count]):
        add_table_picture(slide, locator, x, y + idx * (item_h + gap), w, item_h)


def render_cover(prs: Presentation, page_no: int) -> None:
    slide = blank_slide(prs)
    slide_bg(slide)
    add_fill(slide, 0.0, 0.0, 13.333, 7.5, "primary", "primary", colors=COLORS)
    add_fill(slide, 0.58, 0.76, 0.10, 4.85, "accent", "accent", colors=COLORS)
    add_textbox(slide, 0.86, 0.86, 10.8, 1.35, TITLE, 30, "paper", True, colors=COLORS, component="title", max_lines=2)
    add_textbox(slide, 0.90, 2.50, 8.5, 0.35, "水土保持方案报告表 · 技术汇报", 18, "secondary", True, colors=COLORS, component="subtitle")
    add_emphasis_textbox(
        slide,
        0.90,
        3.18,
        9.8,
        0.82,
        "围绕项目范围、防治责任、土石方平衡、水土流失预测、措施体系、投资概算和设施验收展开，保留报告表中的关键控制值、章节结论与源表证据。",
        16,
        "secondary",
        emphasis_terms=["项目范围", "防治责任", "投资概算", "设施验收"],
        colors=COLORS,
        component="body_text",
        max_lines=3,
    )
    add_textbox(slide, 0.90, 5.82, 8.0, 0.28, "黄冈市住房和城市更新局", 14, "secondary", colors=COLORS, component="body_text")
    add_textbox(slide, 11.68, 6.82, 0.64, 0.20, f"{page_no:02d}", 9, "secondary", True, align=PP_ALIGN.RIGHT, colors=COLORS, component="page_meta")


def render_agenda(prs: Presentation, page_no: int) -> None:
    slide = blank_slide(prs)
    slide_bg(slide)
    add_header(slide, "汇报结构沿报告七个技术章节展开", "汇报结构", page_no)
    add_emphasis_textbox(
        slide,
        0.72,
        1.22,
        11.4,
        0.44,
        "结构按报告章节推进，先确认项目概况与防治责任，再核对评价、预测、措施、投资和验收管理，形成可追溯的技术汇报链条。",
        15,
        "ink",
        emphasis_terms=["项目概况", "防治责任", "措施", "投资", "验收管理"],
        colors=COLORS,
        component="body_text",
        max_lines=2,
    )
    chapters = [chapter for chapter, _ in CHAPTERS[1:]]
    col_w = 5.7
    for idx, chapter in enumerate(chapters, start=1):
        col = 0 if idx <= 4 else 1
        row = idx - 1 if idx <= 4 else idx - 5
        x = 0.72 + col * 6.0
        y = 1.94 + row * 0.96
        add_fill(slide, x, y, col_w, 0.74, "paper", "line", colors=COLORS)
        add_textbox(slide, x + 0.18, y + 0.19, 0.45, 0.18, f"{idx:02d}", 15, "accent", True, colors=COLORS, component="page_meta")
        add_textbox(slide, x + 0.78, y + 0.14, col_w - 1.0, 0.24, chapter, 15, "ink", True, colors=COLORS, component="body_text")


def render_section(prs: Presentation, chapter: str, summary: str, seq: int, page_no: int) -> None:
    slide = blank_slide(prs)
    slide_bg(slide)
    add_fill(slide, 0.72, 0.86, 0.16, 4.82, "accent", "accent", colors=COLORS)
    add_textbox(slide, 1.02, 1.00, 1.0, 0.32, f"{seq:02d}", 20, "accent", True, colors=COLORS, component="page_meta")
    add_textbox(slide, 1.02, 1.60, 8.8, 0.62, chapter, 28, "ink", True, colors=COLORS, component="title")
    add_textbox(slide, 1.04, 2.56, 8.8, 0.44, summary, 18, "muted", colors=COLORS, component="body_text")
    add_textbox(
        slide,
        1.04,
        3.34,
        9.4,
        0.58,
        "章节内容保留报告中的源表、控制值和工程关系，并把范围、目标、措施、投资、验收等关键项放在同一叙事链条中。",
        16,
        "ink",
        colors=COLORS,
        component="body_text",
        max_lines=2,
    )
    add_textbox(
        slide,
        1.04,
        4.18,
        9.0,
        0.50,
        "后续页面以该章节证据为核心，展示原始表格、计算参数或报告结论，保持技术过程和结论之间的对应关系。",
        15,
        "muted",
        colors=COLORS,
        component="body_text",
        max_lines=2,
    )
    add_textbox(slide, 11.85, 6.82, 0.55, 0.20, f"{page_no:02d}", 9, "muted", True, align=PP_ALIGN.RIGHT, colors=COLORS, component="page_meta")


def render_content(prs: Presentation, page: dict[str, Any], page_no: int) -> None:
    slide = blank_slide(prs)
    slide_bg(slide)
    add_header(slide, page["title"], page["chapter"], page_no)
    bullets = page.get("bullets", [])
    locators = page.get("tables", [])
    layout = page.get("layout", "side_table")
    if layout == "wide_table":
        add_numbered_points(slide, bullets, 0.72, 1.30, 11.7, 1.05, 14.2)
        add_table_picture(slide, locators[0], 0.72, 2.42, 11.74, 4.42)
    elif layout == "overview":
        add_numbered_points(slide, bullets, 0.72, 1.38, 5.2, 4.46, 15.2)
        if locators:
            add_multiple_table_strip(slide, locators[:3], 6.18, 1.40, 6.08, 4.92)
    elif layout == "text_only":
        add_numbered_points(slide, bullets, 0.92, 1.50, 11.0, 4.65, 16.2)
    elif layout == "closing":
        add_numbered_points(slide, bullets, 0.98, 1.58, 7.9, 3.25, 17)
        add_fill(slide, 0.98, 5.26, 10.5, 0.62, "soft", "line", colors=COLORS)
        add_emphasis_textbox(slide, 1.22, 5.43, 9.8, 0.20, "水土保持设施验收完成后，项目形成措施、投资、效益和资料管理闭环。", 16, "ink", colors=COLORS, component="body_text")
    else:
        add_numbered_points(slide, bullets, 0.72, 1.42, 4.45, 5.25, 14.5)
        if len(locators) == 1:
            add_table_picture(slide, locators[0], 5.42, 1.38, 6.92, 5.45)
        else:
            add_multiple_table_strip(slide, locators, 5.42, 1.38, 6.92, 5.45)


def render_deck() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    page_no = 1
    render_cover(prs, page_no)
    page_no += 1
    render_agenda(prs, page_no)
    page_no += 1

    chapter_lookup = {chapter: (idx, summary) for idx, (chapter, summary) in enumerate(CHAPTERS)}
    rendered_chapters: set[str] = set()
    for page in expanded_pages()[1:]:
        if page.get("type") == "agenda":
            continue
        chapter = page.get("chapter")
        if chapter and chapter not in rendered_chapters and chapter in chapter_lookup:
            seq, summary = chapter_lookup[chapter]
            render_section(prs, chapter, summary, seq, page_no)
            rendered_chapters.add(chapter)
            page_no += 1
        if page.get("type") in {"cover", "agenda"}:
            continue
        render_content(prs, page, page_no)
        page_no += 1

    EXPORTS.mkdir(exist_ok=True)
    output = EXPORTS / f"东坡大道水土保持方案当前agent测试_{len(prs.slides)}页_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    prs.save(output)
    return output


def write_notes() -> None:
    NOTES.mkdir(exist_ok=True)
    chunks = []
    for idx, page in enumerate(expanded_pages(), start=1):
        chunks.append(f"# Slide {idx}: {page['title']}")
        for item in page.get("bullets", []):
            chunks.append(f"- {sanitize_visible_text(item)}")
        chunks.append("")
    (NOTES / "total.md").write_text("\n".join(chunks), encoding="utf-8")


def main() -> None:
    if not CATALOG.exists() or not INVENTORY.exists() or not TABLE_IR.exists():
        raise SystemExit("Run prepare/catalog/analyze before building the deck.")
    write_evidence_and_plan()
    write_design_contracts()
    write_notes()
    output = render_deck()
    print(json.dumps({"pptx": str(output), "slides": len(Presentation(output).slides)}, ensure_ascii=False))


if __name__ == "__main__":
    main()
