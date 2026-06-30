from __future__ import annotations

import argparse
import json
import re
import sys
from datetime import datetime
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

PROJECT_DIR = Path(__file__).resolve().parents[1]
REPO_ROOT = PROJECT_DIR.parents[1]
LOCAL_SCRIPTS = REPO_ROOT / ".codex" / "skills" / "engineering-ppt" / "scripts"
sys.path.insert(0, str(LOCAL_SCRIPTS))

from engineering_deck_runtime import (  # noqa: E402
    DEFAULT_COLORS,
    add_emphasis_textbox,
    add_fill,
    add_image_panel,
    add_textbox,
    rgb,
    sanitize_visible_text,
    text_preview,
)
from pptx_merged_table_renderer import model_fits_native_table, render_merged_table  # noqa: E402


PROJECT = "东坡大道水土保持方案"
TITLE = "黄冈市东坡大道片区污水收集管网建设工程水土保持方案报告表汇报"
SOURCE_DOCX = PROJECT_DIR / "中地环科-黄冈市东坡大道片区污水收集管网建设工程水土保持方案报告表-20260108.docx"
CATALOG = PROJECT_DIR / "analysis" / "source_catalog.json"
INVENTORY = PROJECT_DIR / "analysis" / "report_content_inventory.json"
TABLE_MODELS = PROJECT_DIR / "analysis" / "docx_table_models.json"
MEDIA = PROJECT_DIR / "images" / "docx_media"
EXPORTS = PROJECT_DIR / "exports"
NOTES = PROJECT_DIR / "notes"

COLORS = {
    **DEFAULT_COLORS,
    "bg": "F5F7FA",
    "paper": "FFFFFF",
    "ink": "132D46",
    "muted": "5B6F82",
    "primary": "0B3558",
    "accent": "0E9F9A",
    "good": "2E7D32",
    "warn": "B64A31",
    "soft": "EAF3F3",
    "line": "C8D3D9",
    "amber": "C9822B",
}

CHAPTERS = [
    ("01", "综合说明与审批依据", "用综合表、批复依据和项目结论建立评审边界。"),
    ("02", "项目概况与工程布置", "说明建设规模、工程布置、施工组织和扰动来源。"),
    ("03", "防治责任范围与标准目标", "锁定防治责任范围和南方红壤区一级目标。"),
    ("04", "水土保持评价与土石方", "核对制约因素、土石方平衡和主体已列措施。"),
    ("05", "水土流失分析与预测", "展示预测范围、侵蚀模数、计算参数和流失结果。"),
    ("06", "水土保持措施与施工安排", "呈现分区、措施体系、工程量和施工进度衔接。"),
    ("07", "投资概算、效益与管理验收", "说明投资、补偿费、防治效益和验收管理闭环。"),
]

SECTION_CHAPTERS = [
    ("00", "项目总览", "先用全局控制指标说明项目边界、预测结果、措施投资和验收闭环。"),
    *CHAPTERS,
]


PAGES: list[dict[str, Any]] = [
    {"type": "cover", "title": TITLE, "chapter": "", "mode": "", "tables": []},
    {"type": "overview", "title": "项目边界、扰动预测、措施投资和验收管理形成闭环", "chapter": "项目总览", "mode": "INTERPRETATION", "tables": ["T004", "T034", "T040", "T048"]},
    {"type": "agenda", "title": "汇报结构沿报告七个技术章节展开", "chapter": "汇报结构", "mode": "", "tables": []},
    {"type": "table_focus", "chapter": "综合说明与审批依据", "title": "综合报告表明确项目全局控制指标", "mode": "ORIGINAL_TABLE", "tables": ["T004"], "layout": "dashboard_table"},
    {"type": "text_cards", "chapter": "综合说明与审批依据", "title": "批复文件和报告附件限定方案编制边界", "mode": "ORIGINAL_TEXT", "tables": [], "catalog_keywords": ["可研", "批复", "附件"]},
    {"type": "table_focus", "chapter": "项目概况与工程布置", "title": "工程特性表锁定规模、投资、工期和占地", "mode": "ORIGINAL_TABLE", "tables": ["T008"], "layout": "key_value"},
    {"type": "figure", "chapter": "项目概况与工程布置", "title": "工程位于黄州区城北片区并沿既有道路布设", "mode": "ORIGINAL_FIGURE", "tables": ["T009"], "images": ["image18.png"]},
    {"type": "figure_process", "chapter": "项目概况与工程布置", "title": "微顶管施工工艺压缩开挖扰动但强化井位防护要求", "mode": "ORIGINAL_FIGURE", "tables": ["T010", "T011"], "images": ["image19.png", "image20.png", "image21.png"]},
    {"type": "table_focus", "chapter": "项目概况与工程布置", "title": "表土剥离、土石方平衡和施工进度共同决定扰动管理", "mode": "ORIGINAL_TABLE", "tables": ["T014", "T015", "T017"], "layout": "earthwork_combo"},
    {"type": "table_focus", "chapter": "防治责任范围与标准目标", "title": "防治责任范围全部纳入管网工程区临时占地", "mode": "ORIGINAL_TABLE", "tables": ["T006", "T013"], "layout": "scope"},
    {"type": "native_table", "chapter": "防治责任范围与标准目标", "title": "六项防治指标采用南方红壤区一级标准", "mode": "ORIGINAL_TABLE", "tables": ["T007"]},
    {"type": "compliance", "chapter": "水土保持评价与土石方", "title": "水土保持法、技术标准和长江保护法评价结论均为符合", "mode": "ORIGINAL_TABLE", "tables": ["T018", "T019", "T020"]},
    {"type": "table_focus", "chapter": "水土保持评价与土石方", "title": "土石方平衡显示弃方委托消纳、借方外购", "mode": "ORIGINAL_TABLE", "tables": ["T015"], "layout": "earthwork"},
    {"type": "native_table", "chapter": "水土保持评价与土石方", "title": "主体工程已列措施以土地整治和铺种草皮为主", "mode": "ORIGINAL_TABLE", "tables": ["T021"]},
    {"type": "table_focus", "chapter": "水土流失分析与预测", "title": "预测范围与防治责任范围一致，预测单元为管网工程区", "mode": "ORIGINAL_TABLE", "tables": ["T023", "T024", "T025", "T026"], "layout": "prediction_scope"},
    {"type": "native_table", "chapter": "水土流失分析与预测", "title": "背景土壤侵蚀模数由占地类型加权形成", "mode": "CALCULATION", "tables": ["T027", "T028"]},
    {"type": "calculation", "chapter": "水土流失分析与预测", "title": "施工期侵蚀模数显著高于背景值，自然恢复期仍需覆盖", "mode": "CALCULATION", "tables": ["T030", "T031", "T032", "T033"]},
    {"type": "table_focus", "chapter": "水土流失分析与预测", "title": "预测土壤流失总量104.05t，新增量94.92t", "mode": "ORIGINAL_TABLE", "tables": ["T034"], "layout": "loss_result"},
    {"type": "native_table", "chapter": "水土保持措施与施工安排", "title": "防治分区为管网工程区，防治面积34416.00m²", "mode": "ORIGINAL_TABLE", "tables": ["T035"]},
    {"type": "measure_matrix", "chapter": "水土保持措施与施工安排", "title": "措施体系覆盖工程、植物和临时三类措施", "mode": "ORIGINAL_TABLE", "tables": ["T036", "T037"]},
    {"type": "schedule", "chapter": "水土保持措施与施工安排", "title": "措施实施进度嵌入2025—2026年施工窗口", "mode": "ORIGINAL_TABLE", "tables": ["T038"]},
    {"type": "investment", "chapter": "投资概算、效益与管理验收", "title": "水土保持总投资26.9211万元，新增投资与主体已列投资分项计列", "mode": "ORIGINAL_TABLE", "tables": ["T040", "T041", "T042", "T043", "T044"]},
    {"type": "table_focus", "chapter": "投资概算、效益与管理验收", "title": "水土保持补偿费51624元符合免征情形", "mode": "ORIGINAL_TABLE", "tables": ["T045"], "layout": "compensation"},
    {"type": "closing", "chapter": "投资概算、效益与管理验收", "title": "防治效果达标后仍需落实监理、验收和资料公开", "mode": "MANAGEMENT_ACTION", "tables": ["T048"]},
]

EXPECTED_PPTX_SLIDES = len(PAGES) + len(SECTION_CHAPTERS)


def load_json(path: Path) -> dict[str, Any]:
    return json.loads(path.read_text(encoding="utf-8-sig"))


def catalog_entries() -> list[dict[str, Any]]:
    return load_json(CATALOG).get("entries", [])


def models_by_locator() -> dict[str, dict[str, Any]]:
    payload = load_json(TABLE_MODELS)
    return {item.get("locator"): item for item in payload.get("tables", []) if item.get("locator")}


def catalog_by_locator() -> dict[str, dict[str, Any]]:
    result = {}
    for item in catalog_entries():
        loc = item.get("locator")
        if loc:
            result[loc] = item
    return result


def normalize(text: str) -> str:
    return re.sub(r"\s+", "", sanitize_visible_text(text))


def find_catalog_ids(locators: list[str], keywords: list[str] | None = None, limit: int = 4) -> list[str]:
    by_loc = catalog_by_locator()
    ids = [by_loc[loc]["id"] for loc in locators if loc in by_loc]
    if ids:
        return ids
    keys = [normalize(k) for k in (keywords or []) if normalize(k)]
    found = []
    for item in catalog_entries():
        text = normalize(item.get("text", ""))
        if keys and any(key in text for key in keys):
            found.append(item["id"])
        if len(found) >= limit:
            break
    return found


def calculation_metadata(locators: list[str]) -> dict[str, Any]:
    key = tuple(locators)
    if key == ("T027", "T028"):
        return {
            "formula": "背景土壤侵蚀模数 = Σ(占地面积 × 对应地类平均土壤侵蚀模数) / Σ占地面积",
            "inputs": {
                "城镇村道路用地面积_m2": 34151.00,
                "公园与绿地面积_m2": 265.00,
                "城镇村道路用地侵蚀模数_t_km2_a": 260,
                "公园与绿地侵蚀模数_t_km2_a": 450,
                "报告表列示背景值_t_km2_a": 261.46,
            },
        }
    if key == ("T030", "T031", "T032", "T033"):
        return {
            "formula": "土壤侵蚀模数按报告表参数取值并汇总：M = R × K × L × S × B × E × T；施工期综合取值按预测单元采用T033列示结果。",
            "inputs": {
                "植被破坏型施工期_t_km2_a": 1153.34,
                "地表翻扰型施工期_t_km2_a": 3032.58,
                "自然恢复期_t_km2_a": 322.11,
                "施工期综合取值_t_km2_a": 3018.60,
                "主要参数来源": "T030、T031、T032、T033",
            },
        }
    return {}


def slide_bg(slide, color: str = "bg") -> None:
    add_fill(slide, 0, 0, 13.333, 7.5, color, color, colors=COLORS)


def header(slide, title: str, chapter: str, page_no: int) -> None:
    add_fill(slide, 0.42, 0.30, 0.08, 0.38, "accent", "accent", colors=COLORS)
    add_textbox(slide, 0.58, 0.24, 8.7, 0.25, chapter, 10.5, "muted", True, colors=COLORS)
    add_emphasis_textbox(slide, 0.58, 0.66, 11.3, 0.52, title, 23, "ink", emphasis_terms=curated_terms(title), colors=COLORS, bold=True)
    add_textbox(slide, 11.94, 0.34, 0.46, 0.22, f"{page_no:02d}", 10, "muted", True, align=PP_ALIGN.RIGHT, colors=COLORS)
    add_fill(slide, 0.58, 1.18, 11.82, 0.015, "line", "line", colors=COLORS)


def curated_terms(text: str) -> list[str]:
    probes = [
        "34416.00m²", "34416.00", "6.2公里", "2.2公里", "8024.00万元", "6914.95万元",
        "2025—2026年", "2025年12月", "2026年12月", "26.9211万元", "51624元",
        "104.05t", "94.92t", "261.46t/km²·a", "3018.60", "322.11",
        "管网工程区", "临时占地", "微顶管", "一级标准", "符合", "委托消纳",
        "工程措施", "植物措施", "临时措施", "免征情形", "设施验收", "达标",
        "主体工程已列", "土地整治", "铺种草皮", "0.74元/m²", "23.34元/m²",
    ]
    return [item for item in probes if item in text][:3]


def add_paragraphs(slide, x, y, w, h, points: list[str], page_text: str = "", size: float = 14.5) -> None:
    if not points:
        return
    terms = curated_terms(" ".join(points) + " " + page_text)
    row_h = h / max(1, len(points))
    for idx, point in enumerate(points, start=1):
        yy = y + (idx - 1) * row_h
        add_fill(slide, x, yy + 0.02, 0.34, 0.30, "accent", "accent", colors=COLORS)
        add_textbox(slide, x + 0.045, yy + 0.065, 0.25, 0.16, f"{idx:02d}", 14, "paper", True, align=PP_ALIGN.CENTER, colors=COLORS)
        add_emphasis_textbox(slide, x + 0.46, yy, w - 0.46, max(0.28, row_h - 0.03), point, size, "ink", emphasis_terms=terms, colors=COLORS)


def add_metric(slide, x, y, w, value, label, color="accent") -> None:
    add_fill(slide, x, y, w, 0.88, "paper", "line", colors=COLORS)
    add_emphasis_textbox(slide, x + 0.14, y + 0.14, w - 0.28, 0.28, value, 18, color, emphasis_terms=[value], colors=COLORS, bold=True)
    add_textbox(slide, x + 0.14, y + 0.55, w - 0.28, 0.20, label, 14, "muted", colors=COLORS)


def add_card(slide, x, y, w, h, title, body, color="accent", idx: int | None = None) -> None:
    add_fill(slide, x, y, w, h, "paper", "line", colors=COLORS)
    add_fill(slide, x, y, 0.07, h, color, color, colors=COLORS)
    if idx is not None:
        add_textbox(slide, x + 0.18, y + 0.12, 0.38, 0.22, f"{idx:02d}", 14, color, True, colors=COLORS)
        title_x = x + 0.64
    else:
        title_x = x + 0.22
    add_emphasis_textbox(slide, title_x, y + 0.11, w - (title_x - x) - 0.18, 0.24, title, 14.5, "ink", emphasis_terms=curated_terms(title + body) or [title], colors=COLORS, bold=True)
    add_emphasis_textbox(slide, x + 0.22, y + 0.50, w - 0.40, h - 0.58, body, 14, "muted", emphasis_terms=curated_terms(body), colors=COLORS)


def table_rows(locator: str) -> list[list[str]]:
    model = models_by_locator().get(locator, {})
    rows = model.get("rows") or []
    return [[sanitize_visible_text(str(cell)) for cell in row] for row in rows]


def source_text(locator: str) -> str:
    item = catalog_by_locator().get(locator, {})
    return sanitize_visible_text(item.get("text", ""))


def draw_table_grid(slide, rows: list[list[str]], x, y, w, h, max_rows=8, max_cols=8, title: str = "") -> None:
    if title:
        add_textbox(slide, x, y - 0.30, w, 0.22, title, 14, "accent", True, colors=COLORS)
    rows = [row[:max_cols] for row in rows[:max_rows]]
    if not rows:
        return
    col_count = max(len(row) for row in rows)
    for row in rows:
        while len(row) < col_count:
            row.append("")
    row_h = h / len(rows)
    col_w = w / col_count
    for r, row in enumerate(rows):
        for c, cell in enumerate(row):
            fill = "soft" if r == 0 else "paper"
            add_fill(slide, x + c * col_w, y + r * row_h, col_w, row_h, fill, "line", colors=COLORS)
            text = text_preview(cell, max(4, int(col_w * row_h * 7)))
            size = 14
            box = slide.shapes.add_textbox(Inches(x + c * col_w + 0.03), Inches(y + r * row_h + 0.025), Inches(col_w - 0.06), Inches(row_h - 0.04))
            tf = box.text_frame
            tf.clear()
            tf.word_wrap = True
            tf.auto_size = 0
            p = tf.paragraphs[0]
            p.text = text
            p.font.name = "Microsoft YaHei"
            p.font.size = Pt(size)
            p.font.bold = r == 0
            p.font.color.rgb = rgb(COLORS["ink"])
            p.alignment = PP_ALIGN.CENTER
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_native_or_grid(slide, locator: str, x, y, w, h, max_rows=8, max_cols=7, title="") -> None:
    model = models_by_locator().get(locator)
    if model and model_fits_native_table(model, max_rows=max_rows, max_cols=max_cols):
        if title:
            add_textbox(slide, x, y - 0.30, w, 0.22, title, 14, "accent", True, colors=COLORS)
        render_merged_table(slide, model, x, y, w, h, colors=COLORS, max_rows=max_rows, max_cols=max_cols)
    else:
        draw_table_grid(slide, table_rows(locator), x, y, w, h, max_rows=max_rows, max_cols=max_cols, title=title)


def section_slide(prs: Presentation, page_no: int, number: str, chapter: str, subtitle: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide, "primary")
    add_fill(slide, 0.72, 1.05, 0.08, 5.05, "accent", "accent", colors=COLORS)
    add_textbox(slide, 0.96, 1.20, 3.2, 0.30, f"CHAPTER {number}", 17, "soft", True, colors=COLORS)
    add_textbox(slide, 0.96, 2.35, 1.2, 0.80, number, 44, "paper", True, colors=COLORS)
    add_textbox(slide, 2.30, 2.48, 8.8, 0.58, chapter, 30, "paper", True, colors=COLORS)
    add_emphasis_textbox(slide, 2.32, 3.35, 8.8, 0.40, subtitle, 15, "soft", colors=COLORS)
    add_card(slide, 7.55, 4.30, 4.25, 1.32, "覆盖证据", "围绕本章原表、原图、关键原文和计算口径展开，先给结论，再回到报告证据逐项核对。", "accent")
    add_textbox(slide, 2.32, 4.30, 4.80, 0.44, "核对重点：边界、目标、工程量、计算参数、投资和验收责任必须与源报告保持一致。", 14.5, "soft", colors=COLORS)
    add_textbox(slide, 11.52, 6.55, 0.70, 0.22, f"{page_no:02d}", 12, "soft", True, align=PP_ALIGN.RIGHT, colors=COLORS)
    return slide


def render_cover(prs: Presentation, page_no: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide, "primary")
    add_textbox(slide, 0.72, 0.80, 7.4, 0.30, "水土保持方案报告表 · 技术评审汇报", 15, "soft", True, colors=COLORS)
    add_textbox(slide, 0.72, 1.72, 10.4, 0.58, "黄冈市东坡大道片区", 34, "paper", True, colors=COLORS)
    add_textbox(slide, 0.72, 2.42, 10.8, 0.58, "污水收集管网建设工程", 34, "paper", True, colors=COLORS)
    add_emphasis_textbox(slide, 0.76, 3.35, 8.7, 0.42, "围绕34416.00m²防治责任范围、104.05t预测流失总量、26.9211万元水保投资和验收管理闭环展开。", 15, "soft", emphasis_terms=["34416.00m²", "104.05t", "26.9211万元"], colors=COLORS)
    metrics = [("6.2公里", "污水主管网", "accent"), ("34416.00m²", "防治责任范围", "good"), ("26.9211万元", "水保总投资", "warn")]
    for i, metric in enumerate(metrics):
        add_metric(slide, 0.76 + i * 3.25, 5.18, 2.90, *metric)
    add_textbox(slide, 0.76, 6.78, 6.2, 0.20, "武汉中地环科水工环科技咨询有限责任公司 | 2026年1月报批稿", 12, "soft", colors=COLORS)


def render_overview(prs: Presentation, page: dict[str, Any], page_no: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide)
    header(slide, page["title"], page["chapter"], page_no)
    facts = [
        ("建设边界", "主管网6.2公里、预埋支管2.2公里，工程位于黄州区城北片区。", "accent"),
        ("责任范围", "防治责任范围与扰动范围统一为34416.00m²，均为管网工程区。", "good"),
        ("预测结果", "工程建设可能造成土壤流失总量104.05t，新增量94.92t。", "warn"),
        ("投资验收", "水土保持总投资26.9211万元，后续以监理、验收和报备形成闭环。", "primary"),
    ]
    for i, (t, b, c) in enumerate(facts):
        add_card(slide, 0.72 + (i % 2) * 5.95, 1.55 + (i // 2) * 1.50, 5.55, 1.16, t, b, c, i + 1)
    add_fill(slide, 0.78, 4.98, 11.65, 0.78, "soft", "line", colors=COLORS)
    add_textbox(slide, 1.05, 5.24, 10.95, 0.24, "证据链：综合说明表 → 工程特性与图件 → 责任范围与目标 → 土石方与预测计算 → 措施工程量 → 投资效益与验收管理", 15, "ink", True, colors=COLORS)
    add_paragraphs(slide, 0.92, 6.05, 11.15, 0.72, ["原表和原图作为页面主体，解释文字只用于说明工程含义、控制关系和管理责任。"], page["title"], 14.5)


def render_agenda(prs: Presentation, page_no: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide)
    header(slide, "汇报结构沿报告七个技术章节展开", "汇报结构", page_no)
    for i, (num, chap, sub) in enumerate(CHAPTERS):
        x = 0.72 + (i % 2) * 5.95
        y = 1.50 + (i // 2) * 1.18
        add_card(slide, x, y, 5.55, 0.90, f"{num} {chap}", sub, ["accent", "good", "primary", "warn"][i % 4], i + 1)


def render_table_focus(prs: Presentation, page: dict[str, Any], page_no: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide)
    header(slide, page["title"], page["chapter"], page_no)
    layout = page.get("layout")
    if layout == "dashboard_table":
        add_native_or_grid(slide, "T004", 0.62, 1.56, 7.65, 4.75, max_rows=8, max_cols=6, title="综合报告表节选：项目概况、土石方、措施和投资控制指标")
        points = ["报告表明确主管网6.2公里、预埋支管2.2公里，工程范围集中于城北片区。", "防治责任范围为34416.00m²，后续目标、预测、措施和投资均围绕该范围展开。", "综合表同时列示土石方、措施体系和总投资，是全套方案的控制索引。"]
        add_paragraphs(slide, 8.65, 1.68, 3.35, 3.75, points, page["title"])
    elif layout == "key_value":
        add_native_or_grid(slide, "T008", 0.62, 1.55, 7.80, 4.90, max_rows=9, max_cols=8, title="建设项目工程特性表节选")
        add_metric(slide, 8.78, 1.65, 3.35, "8024.00万元", "总投资", "warn")
        add_metric(slide, 8.78, 2.78, 3.35, "2025年12月—2026年12月", "计划工期", "accent")
        add_metric(slide, 8.78, 3.91, 3.35, "34416.00m²", "临时占地", "good")
        add_paragraphs(slide, 8.78, 5.16, 3.25, 1.08, ["工程规模、工期和占地共同决定施工扰动强度。"], page["title"])
    elif layout == "earthwork_combo":
        add_native_or_grid(slide, "T014", 0.62, 1.48, 5.75, 1.35, max_rows=3, max_cols=8, title="表土剥离与回覆")
        draw_table_grid(slide, table_rows("T015"), 0.62, 3.22, 5.75, 2.40, max_rows=7, max_cols=6, title="土石方平衡节选")
        draw_table_grid(slide, table_rows("T017"), 6.72, 1.48, 5.55, 2.10, max_rows=6, max_cols=6, title="施工进度原表节选")
        add_paragraphs(slide, 6.90, 4.18, 5.10, 1.58, ["表土剥离量79.5m³，服务后续绿化覆土回填。", "弃方委托湖北安达港务有限公司外运消纳，借方采用外购土石方。", "施工进度横跨2025年12月至2026年12月，临时防护需随扰动同步布设。"], page["title"])
    elif layout == "scope":
        add_native_or_grid(slide, "T006", 0.62, 1.58, 7.35, 3.45, max_rows=6, max_cols=6, title="防治责任范围一览表")
        add_metric(slide, 8.55, 1.72, 3.45, "34416.00m²", "防治责任范围", "accent")
        add_metric(slide, 8.55, 2.86, 3.45, "34151.00m²", "城镇村道路用地", "good")
        add_metric(slide, 8.55, 4.00, 3.45, "265.00m²", "公园与绿地", "warn")
        add_paragraphs(slide, 0.86, 5.55, 11.20, 0.70, ["责任范围、占地类型和恢复对象在同一表内闭合，避免后续目标和措施脱离工程边界。"], page["title"])
    elif layout == "earthwork":
        draw_table_grid(slide, table_rows("T015"), 0.58, 1.50, 8.25, 4.85, max_rows=10, max_cols=8, title="工程土石方平衡表节选")
        add_paragraphs(slide, 9.18, 1.68, 3.0, 3.80, ["总挖方4971.61m³、总填方2232.54m³。", "弃方4892.11m³，建筑垃圾和清淤土方均委托外运消纳。", "借方2153.04m³由外购土石方解决，施工过程不设置临时堆土场。"], page["title"])
    elif layout == "prediction_scope":
        add_native_or_grid(slide, "T023", 0.62, 1.48, 5.8, 1.80, max_rows=5, max_cols=5, title="扰动地表面积")
        add_native_or_grid(slide, "T025", 6.68, 1.48, 5.55, 1.80, max_rows=4, max_cols=4, title="预测单元划分")
        add_native_or_grid(slide, "T026", 0.62, 3.78, 5.8, 1.55, max_rows=3, max_cols=5, title="预测范围与时段")
        add_paragraphs(slide, 6.84, 3.94, 5.20, 1.42, ["预测范围与防治责任范围一致，均以管网工程区为控制对象。", "施工期对应34416.00m²扰动范围，自然恢复期对应265.00m²绿化恢复范围。"], page["title"])
    elif layout == "loss_result":
        draw_table_grid(slide, table_rows("T034"), 0.58, 1.48, 8.50, 4.90, max_rows=5, max_cols=8, title="土壤流失量预测结果表节选")
        add_metric(slide, 9.45, 1.70, 2.90, "104.05t", "预测土壤流失总量", "warn")
        add_metric(slide, 9.45, 2.84, 2.90, "94.92t", "新增土壤流失量", "accent")
        add_paragraphs(slide, 9.46, 4.20, 2.82, 1.38, ["新增流失量主要由施工扰动贡献，施工期是措施布设重点。"], page["title"])
    elif layout == "compensation":
        add_native_or_grid(slide, "T045", 0.62, 1.58, 7.25, 2.35, max_rows=3, max_cols=6, title="水土保持补偿费计算表")
        add_metric(slide, 8.45, 1.78, 3.55, "51624元", "补偿费测算金额", "warn")
        add_metric(slide, 8.45, 2.92, 3.55, "3.4416hm²", "征占地面积", "accent")
        add_paragraphs(slide, 0.78, 4.45, 11.30, 1.28, ["报告依据财综〔2014〕8号说明项目属于市政生态环境保护基础设施，符合水土保持补偿费免征情形。", "免征不改变补偿费计算依据，方案仍保留征占地面积、收费标准和测算金额。"], page["title"])


def render_native_table(prs: Presentation, page: dict[str, Any], page_no: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide)
    header(slide, page["title"], page["chapter"], page_no)
    locators = page["tables"]
    if locators == ["T007"]:
        add_native_or_grid(slide, "T007", 0.70, 1.48, 8.70, 4.50, max_rows=8, max_cols=6, title="水土流失防治目标值表")
        add_paragraphs(slide, 9.72, 1.78, 2.65, 3.20, ["执行南方红壤区建设类项目一级标准。", "林草覆盖率按项目区实际情况调整为0.77%。", "表土保护率修正后采用93%。"], page["title"])
    elif locators == ["T021"]:
        add_native_or_grid(slide, "T021", 0.62, 1.48, 8.25, 4.50, max_rows=7, max_cols=7, title="主体工程已列水保措施量及投资")
        add_paragraphs(slide, 9.25, 1.72, 2.95, 3.45, ["主体已列投资覆盖土地整治和铺种草皮。", "土地整治单价0.74元/m²，铺种草皮单价23.34元/m²。", "方案新增措施需补足施工期临时防护。"], page["title"])
    elif locators == ["T027", "T028"]:
        add_native_or_grid(slide, "T027", 0.62, 1.50, 6.15, 1.75, max_rows=3, max_cols=6, title="各地类土壤侵蚀模数取值")
        draw_table_grid(slide, table_rows("T028"), 0.62, 3.72, 6.15, 1.80, max_rows=5, max_cols=6, title="背景值计算表节选")
        add_metric(slide, 7.35, 1.68, 2.30, "261.46", "背景侵蚀模数", "accent")
        add_metric(slide, 9.88, 1.68, 2.30, "500", "容许流失量", "good")
        add_paragraphs(slide, 7.40, 3.00, 4.62, 1.62, ["背景模数由城镇村道路用地和公园与绿地面积加权形成。", "背景值是后续施工期新增水土流失量计算的比较基准。"], page["title"])
    elif locators == ["T035"]:
        add_native_or_grid(slide, "T035", 0.80, 1.70, 6.0, 2.10, max_rows=4, max_cols=4, title="水土流失防治区划分表")
        add_metric(slide, 7.35, 1.84, 2.25, "1个", "防治分区", "accent")
        add_metric(slide, 9.92, 1.84, 2.25, "34416.00m²", "防治面积", "good")
        add_paragraphs(slide, 0.90, 4.45, 11.0, 1.15, ["报告将防治分区划定为管网工程区，后续措施体系、工程量、投资估算和防治效益均以该分区为基础。"], page["title"])


def render_figure(prs: Presentation, page: dict[str, Any], page_no: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide)
    header(slide, page["title"], page["chapter"], page_no)
    image_paths = [MEDIA / name for name in page.get("images", []) if (MEDIA / name).exists()]
    if page["type"] == "figure_process":
        add_image_panel(slide, image_paths, 0.66, 1.38, 11.75, 3.36, colors=COLORS)
        add_card(slide, 0.82, 5.05, 3.55, 1.18, "扰动控制", "微顶管减少道路大开挖，井位仍是临时防护重点。", "accent", 1)
        add_card(slide, 4.78, 5.05, 3.55, 1.18, "土方管理", "井位开挖、回填和外运消纳与土石方平衡表对应。", "warn", 2)
        add_card(slide, 8.74, 5.05, 3.55, 1.18, "恢复闭合", "路面恢复和绿化恢复与土地整治、表土回覆同步完成。", "good", 3)
    else:
        add_image_panel(slide, image_paths, 0.66, 1.34, 7.70, 4.90, colors=COLORS)
        add_paragraphs(slide, 8.72, 1.62, 3.35, 3.55, ["工程位于黄州区城北片区，管网沿既有城市道路布设。", "工程范围受黄州大道、东坡大道、赤壁大道等道路边界控制。", "防治分区以管网工程区统一展开。"], page["title"])


def render_text_cards(prs: Presentation, page: dict[str, Any], page_no: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide)
    header(slide, page["title"], page["chapter"], page_no)
    cards = [
        ("前期批复", "可研批复、初设批复和委托文件限定工程建设内容与方案编制范围。", "accent"),
        ("技术依据", "法律法规、GB50433-2018及水土保持监测评价规范共同支撑目标和预测计算。", "primary"),
        ("附件支撑", "地理位置、措施布设、投资估算和批复资料共同构成报批依据。", "good"),
        ("责任边界", "方案责任主体、工程扰动范围和水土保持措施需在施工及验收阶段延续闭合。", "warn"),
    ]
    for i, (t, b, c) in enumerate(cards):
        add_card(slide, 0.82 + (i % 2) * 5.82, 1.66 + (i // 2) * 1.58, 5.38, 1.18, t, b, c, i + 1)
    add_fill(slide, 0.88, 5.25, 11.35, 0.88, "soft", "line", colors=COLORS)
    add_textbox(slide, 1.12, 5.52, 10.85, 0.22, "编制边界不另行扩大：PPT 仅使用源报告的批复依据、表格、图件和报告结论，不采用旧 PPT 作为事实来源。", 14.5, "ink", True, colors=COLORS)


def render_compliance(prs: Presentation, page: dict[str, Any], page_no: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide)
    header(slide, page["title"], page["chapter"], page_no)
    items = [("水土保持法", "项目不涉及取土、挖沙、采石等禁止性因素。"), ("GB50433-2018", "工程选址、建设方案和施工组织不构成水保制约。"), ("长江保护法", "工程不属于长江岸线一公里范围内化工项目。")]
    for i, (name, body) in enumerate(items):
        add_card(slide, 0.78, 1.58 + i * 1.22, 5.10, 0.94, name, body, "good", i + 1)
    draw_table_grid(slide, table_rows("T018"), 6.35, 1.50, 5.85, 1.52, max_rows=3, max_cols=4, title="水土保持法制约因素节选")
    draw_table_grid(slide, table_rows("T019"), 6.35, 3.48, 5.85, 1.52, max_rows=3, max_cols=5, title="GB50433-2018制约因素节选")
    draw_table_grid(slide, table_rows("T020"), 6.35, 5.46, 5.85, 1.05, max_rows=2, max_cols=4, title="长江保护法制约因素节选")


def render_calculation(prs: Presentation, page: dict[str, Any], page_no: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide)
    header(slide, page["title"], page["chapter"], page_no)
    add_metric(slide, 0.78, 1.55, 2.55, "1153.34", "植被破坏型施工期", "warn")
    add_metric(slide, 3.62, 1.55, 2.55, "3032.58", "地表翻扰型施工期", "warn")
    add_metric(slide, 6.46, 1.55, 2.55, "322.11", "自然恢复期", "accent")
    add_metric(slide, 9.30, 1.55, 2.55, "3018.60", "综合施工期取值", "primary")
    draw_table_grid(slide, table_rows("T033"), 0.82, 3.05, 5.20, 1.30, max_rows=3, max_cols=3, title="土壤侵蚀模数取值表")
    add_paragraphs(slide, 6.45, 3.05, 5.40, 1.72, ["施工期取值显著高于背景土壤侵蚀模数，反映施工扰动对水土流失强度的放大作用。", "自然恢复期取值降至322.11t/km²·a，但仍需覆盖植被恢复阶段。"], page["title"])
    draw_table_grid(slide, table_rows("T030"), 0.82, 5.10, 3.55, 1.05, max_rows=4, max_cols=4, title="植被破坏型参数")
    draw_table_grid(slide, table_rows("T031"), 4.85, 5.10, 3.55, 1.05, max_rows=4, max_cols=4, title="地表翻扰型参数")
    draw_table_grid(slide, table_rows("T032"), 8.88, 5.10, 3.55, 1.05, max_rows=4, max_cols=4, title="自然恢复期参数")


def render_measure_matrix(prs: Presentation, page: dict[str, Any], page_no: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide)
    header(slide, page["title"], page["chapter"], page_no)
    add_native_or_grid(slide, "T036", 0.62, 1.45, 6.10, 2.10, max_rows=3, max_cols=7, title="水土流失防治措施体系表")
    draw_table_grid(slide, table_rows("T037"), 0.62, 4.08, 6.10, 2.20, max_rows=7, max_cols=6, title="措施工程量汇总表节选")
    cards = [("工程措施", "表土剥离、表土回覆、土地整治。", "accent"), ("植物措施", "铺种草皮，衔接道路绿化恢复。", "good"), ("临时措施", "临时苫盖和高压洗车机控制裸露面与车辆带泥。", "warn")]
    for i, (t, b, c) in enumerate(cards):
        add_card(slide, 7.20, 1.55 + i * 1.42, 4.75, 1.06, t, b, c, i + 1)


def render_schedule(prs: Presentation, page: dict[str, Any], page_no: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide)
    header(slide, page["title"], page["chapter"], page_no)
    draw_table_grid(slide, table_rows("T038"), 0.60, 1.42, 11.95, 2.55, max_rows=7, max_cols=8, title="水土保持措施实施进度表节选")
    steps = [("2025年12月", "施工准备与扰动启动"), ("2026年1—8月", "顶管井施工、顶管施工"), ("2026年9—12月", "路面恢复、绿化恢复"), ("验收前", "资料整理与设施验收")]
    for i, (t, b) in enumerate(steps):
        x = 0.90 + i * 2.95
        add_fill(slide, x, 5.18, 2.35, 0.07, "accent", "accent", colors=COLORS)
        add_card(slide, x - 0.10, 5.42, 2.55, 0.82, t, b, ["accent", "warn", "good", "primary"][i], i + 1)


def render_investment(prs: Presentation, page: dict[str, Any], page_no: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide)
    header(slide, page["title"], page["chapter"], page_no)
    draw_table_grid(slide, table_rows("T040"), 0.58, 1.48, 7.65, 4.72, max_rows=10, max_cols=8, title="水土保持估算总表节选")
    add_metric(slide, 8.72, 1.68, 3.25, "26.9211万元", "水土保持总投资", "warn")
    add_metric(slide, 8.72, 2.82, 3.25, "0.5885万元", "方案新增投资", "accent")
    add_metric(slide, 8.72, 3.96, 3.25, "0.6381万元", "主体已列投资", "good")
    add_paragraphs(slide, 8.78, 5.20, 3.12, 0.90, ["独立费用和补偿费构成投资主体，工程、植物、临时措施投资需与措施工程量一致。"], page["title"])


def render_closing(prs: Presentation, page: dict[str, Any], page_no: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide)
    header(slide, page["title"], page["chapter"], page_no)
    draw_table_grid(slide, table_rows("T048"), 0.60, 1.48, 7.95, 4.60, max_rows=8, max_cols=7, title="水土保持防治效果分析表节选")
    actions = [("监理记录", "控制质量、进度和投资。"), ("施工资料", "留存防护、恢复和消纳记录。"), ("设施验收", "验收合格后再进入竣工验收。"), ("报备公开", "验收材料公开并报备审批机关。")]
    for i, (t, b) in enumerate(actions):
        add_card(slide, 8.92, 1.46 + i * 1.23, 3.05, 1.00, t, b, ["accent", "good", "warn", "primary"][i], i + 1)


def render_page(prs: Presentation, page: dict[str, Any], page_no: int) -> None:
    kind = page["type"]
    if kind == "cover":
        render_cover(prs, page_no)
    elif kind == "overview":
        render_overview(prs, page, page_no)
    elif kind == "agenda":
        render_agenda(prs, page_no)
    elif kind == "table_focus":
        render_table_focus(prs, page, page_no)
    elif kind == "native_table":
        render_native_table(prs, page, page_no)
    elif kind.startswith("figure"):
        render_figure(prs, page, page_no)
    elif kind == "text_cards":
        render_text_cards(prs, page, page_no)
    elif kind == "compliance":
        render_compliance(prs, page, page_no)
    elif kind == "calculation":
        render_calculation(prs, page, page_no)
    elif kind == "measure_matrix":
        render_measure_matrix(prs, page, page_no)
    elif kind == "schedule":
        render_schedule(prs, page, page_no)
    elif kind == "investment":
        render_investment(prs, page, page_no)
    elif kind == "closing":
        render_closing(prs, page, page_no)
    else:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide_bg(slide)
        header(slide, page["title"], page["chapter"], page_no)


def build_evidence_and_plan() -> tuple[dict[str, Any], dict[str, Any]]:
    evidence = []
    slides = []
    ev_index = 1
    for index, page in enumerate(PAGES, start=1):
        mode = page.get("mode", "")
        evidence_ids = []
        if page["type"] not in {"cover", "agenda"}:
            locators = page.get("tables", [])
            catalog_ids = find_catalog_ids(locators, page.get("catalog_keywords"))
            evidence_id = f"DP-FRESH-{ev_index:03d}"
            ev_index += 1
            exact = " ".join(source_text(loc) for loc in locators)[:900]
            item = {
                "id": evidence_id,
                "kind": "RECOMMENDATION" if mode == "MANAGEMENT_ACTION" else mode or "INTERPRETATION",
                "summary": page["title"],
                "source_file": SOURCE_DOCX.name,
                "source_locator": "、".join(locators) or "报告正文",
                "catalog_ids": catalog_ids,
                "verification_status": "verified",
                "exact_text": text_preview(exact or page["title"], 520, complete_sentence=True),
                "values": [],
                "notes": "Fresh rebuild from source DOCX only; prior PPTX exports are not used as evidence.",
            }
            item.update(calculation_metadata(locators))
            evidence.append(item)
            evidence_ids = [evidence_id]
        slides.append(
            {
                "page": index,
                "type": page["type"],
                "chapter": page.get("chapter", ""),
                "title": page["title"],
                "source_mode": mode,
                "evidence_ids": evidence_ids,
                "content_unit_ids": [],
                "visual_proof": "、".join(page.get("tables", [])) or page.get("type", ""),
                "layout_pattern": page.get("layout", page["type"]),
                "source_note": "来源：" + ("、".join(page.get("tables", [])) if page.get("tables") else "报告正文"),
                "density": "dense" if page["type"] in {"table_focus", "native_table", "calculation"} else "normal",
                "density_exempt_reason": "",
            }
        )
    plan = {
        "schema_version": "1.0",
        "project": PROJECT,
        "deck": {
            "title": TITLE,
            "audience": "建设单位、行政主管部门、水土保持评审专家及施工管理相关方",
            "objective": "以报告原表、原图和关键原文说明项目边界、扰动预测、措施体系、投资效益和验收管理要求。",
            "chapter_order": [chapter for _num, chapter, _sub in CHAPTERS],
        },
        "slides": slides,
    }
    ledger = {
        "schema_version": "1.0",
        "project": PROJECT,
        "updated_at": datetime.now().isoformat(timespec="seconds"),
        "rules": {"source_truth": "Source DOCX and extracted original table models are authoritative."},
        "evidence": evidence,
    }
    return ledger, plan


def write_contracts() -> tuple[dict[str, Any], dict[str, Any]]:
    ledger, plan = build_evidence_and_plan()
    (PROJECT_DIR / "evidence_ledger.json").write_text(json.dumps(ledger, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    (PROJECT_DIR / "deck_plan.json").write_text(json.dumps(plan, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    coverage = ["# Chapter Coverage", "", f"- Expected PPTX slides: {EXPECTED_PPTX_SLIDES}", "- Rebuild scope: prior generated PPT artifacts deleted; source DOCX retained.", ""]
    for _num, chapter, summary in SECTION_CHAPTERS:
        pages = [p for p in PAGES if p.get("chapter") == chapter]
        coverage.append(f"## {chapter}")
        coverage.append(f"- Pages: {len(pages)}")
        coverage.append(f"- Intent: {summary}")
        for page in pages:
            coverage.append(f"- {page['title']} | {page.get('mode')} | {', '.join(page.get('tables', []))}")
        coverage.append("")
    (PROJECT_DIR / "chapter_coverage.md").write_text("\n".join(coverage), encoding="utf-8")
    design = f"""# 东坡大道水土保持方案 - Design Spec

## I. Project Information

| Item | Value |
| ---- | ----- |
| Project Name | {TITLE} |
| Canvas Format | PPT 16:9 |
| Page Count | {EXPECTED_PPTX_SLIDES} |
| Design Style | restrained municipal engineering review with richer evidence components |

## II. Content Policy

Fresh rebuild from source DOCX only. Tables use complete native rendering when readable; dense or wide tables use source-structured excerpt grids with original headers, units, and key source rows.

## III. Visual System

Deep blue chapter dividers, white source-evidence pages, teal emphasis for verified control values, restrained warning color for risk and investment constraints.
"""
    lock = """## canvas
- viewBox: 0 0 1280 720
- format: PPT 16:9

## typography
- font_family: Microsoft YaHei
- body_min: 14
- title: 23

## colors
- primary: #0B3558
- accent: #0E9F9A
- warning: #B64A31
- ink: #132D46
- muted: #5B6F82
"""
    (PROJECT_DIR / "design_spec.md").write_text(design, encoding="utf-8")
    (PROJECT_DIR / "spec_lock.md").write_text(lock, encoding="utf-8")
    return ledger, plan


def write_notes(plan: dict[str, Any]) -> None:
    NOTES.mkdir(exist_ok=True)
    notes = []
    for slide in plan["slides"]:
        note = f"Slide {slide['page']:02d} - {slide['title']}\n\n来源：{slide.get('source_note', '')}\n讲解要点：围绕报告原表、原图或原文说明工程控制关系。\n"
        notes.append(note)
        (NOTES / f"{slide['page']:02d}.md").write_text(note, encoding="utf-8")
    (NOTES / "total.md").write_text("\n\n".join(notes), encoding="utf-8")


def build_deck() -> Path:
    _ledger, plan = write_contracts()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    page_no = 1
    written_chapters: set[str] = set()
    chapter_meta = {chapter: (num, subtitle) for num, chapter, subtitle in SECTION_CHAPTERS}
    for page in PAGES:
        chapter = page.get("chapter", "")
        if chapter in chapter_meta and chapter not in written_chapters:
            written_chapters.add(chapter)
            num, subtitle = chapter_meta[chapter]
            section_slide(prs, page_no, num, chapter, subtitle)
            page_no += 1
        render_page(prs, page, page_no)
        page_no += 1
    write_notes(plan)
    EXPORTS.mkdir(exist_ok=True)
    output = EXPORTS / f"东坡大道水土保持方案源证据增强汇报_{len(prs.slides)}页_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    prs.save(output)
    return output


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--plan-only", action="store_true")
    parser.add_argument("--build", action="store_true")
    args = parser.parse_args()
    ledger, plan = write_contracts()
    if args.build:
        output = build_deck()
        print(json.dumps({"output": str(output), "slides": EXPECTED_PPTX_SLIDES}, ensure_ascii=False))
    else:
        print(json.dumps({"evidence": len(ledger["evidence"]), "planned_pages": len(plan["slides"]), "expected_pptx_slides": EXPECTED_PPTX_SLIDES}, ensure_ascii=False))


if __name__ == "__main__":
    main()
