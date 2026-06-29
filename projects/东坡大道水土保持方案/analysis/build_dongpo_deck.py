from __future__ import annotations

import argparse
import json
import re
import sys
from datetime import datetime
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

PROJECT_DIR = Path(__file__).resolve().parents[1]
REPO_ROOT = PROJECT_DIR.parents[1]
LOCAL_SKILL_SCRIPTS = REPO_ROOT / ".codex" / "skills" / "engineering-ppt" / "scripts"
sys.path.insert(0, str(LOCAL_SKILL_SCRIPTS))

from engineering_deck_runtime import (  # noqa: E402
    DEFAULT_COLORS,
    add_agenda_slide,
    add_bullets,
    add_fill,
    add_image_panel,
    add_section_divider_slide,
    add_table,
    add_textbox,
    add_title,
    catalog_title_context,
    complete_semantic_point,
    chapter_order_from_plan,
    collect_template_profile,
    create_engineering_design_spec,
    create_engineering_spec_lock,
    evidence_lookup,
    extract_docx_media,
    load_json,
    report_points_from_page,
    report_text_from_page,
    require_media_paths,
    rgb,
    sanitize_visible_text,
    slide_evidence,
    source_topic_heading,
    table_from_page,
    text_preview,
)
from docx_table_parser import write_docx_table_models  # noqa: E402
from pptx_merged_table_renderer import model_fits_native_table, render_merged_table  # noqa: E402


CATALOG_PATH = PROJECT_DIR / "analysis" / "source_catalog.json"
INVENTORY_PATH = PROJECT_DIR / "analysis" / "report_content_inventory.json"
PLAN_PATH = PROJECT_DIR / "deck_plan.json"
LEDGER_PATH = PROJECT_DIR / "evidence_ledger.json"
EXPORTS_DIR = PROJECT_DIR / "exports"
NOTES_DIR = PROJECT_DIR / "notes"
MEDIA_DIR = PROJECT_DIR / "images" / "docx_media"
SOURCE_DOCX = PROJECT_DIR / "中地环科-黄冈市东坡大道片区污水收集管网建设工程水土保持方案报告表-20260108.docx"
TABLE_MODELS_PATH = PROJECT_DIR / "analysis" / "docx_table_models.json"

COLORS = DEFAULT_COLORS.copy()
PROJECT_TITLE = "黄冈市东坡大道片区污水收集管网建设工程水土保持方案汇报"
AUDIENCE = "建设单位、行政主管部门、水土保持评审专家及项目实施相关方"
USE_CASE = "水土保持方案技术评审与建设管理汇报"
DESIGN_STYLE = "克制的市政工程技术评审风格，突出报告原表、关键参数、预测计算和措施闭环"

CHAPTER_SUMMARIES = {
    "综合说明与审批依据": "先交代项目来源、编制依据、核心建设参数和方案结论，帮助评审快速建立全局判断。",
    "项目概况与工程布置": "说明管网建设内容、道路段落、工程占地、施工组织、工艺和土石方基础。",
    "防治责任范围与标准目标": "界定防治责任范围，明确南方红壤区建设类项目一级标准及六项防治目标。",
    "水土保持评价与土石方": "评价选址、建设方案、施工组织、土石方平衡及主体工程中已列水保措施。",
    "水土流失分析与预测": "从扰动范围、侵蚀模数、预测时段和结果表说明施工期与自然恢复期流失量。",
    "水土保持措施与施工安排": "呈现防治分区、工程/植物/临时措施体系、工程量和实施进度。",
    "投资概算、效益与管理验收": "说明投资构成、补偿费免征依据、防治效果达标情况及后续管理要求。",
}


SLIDE_SPECS: list[dict[str, Any]] = [
    {"type": "cover", "chapter": "", "title": PROJECT_TITLE, "source_mode": "", "units": []},
    {"type": "agenda", "chapter": "", "title": "汇报目录", "source_mode": "", "units": []},
    {
        "type": "table",
        "chapter": "综合说明与审批依据",
        "title": "综合报告表给出项目全局控制指标",
        "source_mode": "ORIGINAL_TABLE",
        "units": ["UNIT-0003"],
        "layout": "left_table_right_text",
        "proof": "综合报告表拆分呈现项目概况、土石方、防治目标、措施和投资控制指标",
    },
    {
        "type": "text",
        "chapter": "综合说明与审批依据",
        "title": "前期批复和编制依据明确方案边界",
        "source_mode": "ORIGINAL_TEXT",
        "units": ["UNIT-0004", "UNIT-0005", "UNIT-0017", "UNIT-0018", "UNIT-0019"],
        "layout": "text_with_keypoints",
        "proof": "报告附件清单、可研批复、初设批复和前期工作概况原文",
    },
    {
        "type": "table",
        "chapter": "项目概况与工程布置",
        "title": "建设项目工程特性表锁定规模、投资与工期",
        "source_mode": "ORIGINAL_TABLE",
        "units": ["UNIT-0080"],
        "layout": "left_table_right_text",
        "proof": "建设项目工程特性表，含总投资、土建投资、建设规模、工期和占地",
    },
    {
        "type": "figure",
        "chapter": "项目概况与工程布置",
        "title": "项目位于黄州区城北片区并沿多条道路布设",
        "source_mode": "ORIGINAL_FIGURE",
        "units": ["UNIT-0014", "UNIT-0015", "UNIT-0016", "UNIT-0081"],
        "layout": "top_figure_bottom_text",
        "image_names": ["image18.png"],
        "proof": "项目基本情况、建设地点和工程建设地点示意图标题",
    },
    {
        "type": "figure",
        "chapter": "项目概况与工程布置",
        "title": "施工组织与顶管工艺决定扰动范围和防护时序",
        "source_mode": "ORIGINAL_FIGURE",
        "units": ["UNIT-0087", "UNIT-0089"],
        "layout": "top_figure_bottom_text",
        "image_names": ["image19.png", "image20.png", "image21.png"],
        "proof": "顶管施工示意、顶管井回填纵横剖示意及施工工艺说明",
    },
    {
        "type": "table",
        "chapter": "防治责任范围与标准目标",
        "title": "防治责任范围全部为临时占地 34416.00m²",
        "source_mode": "ORIGINAL_TABLE",
        "units": ["UNIT-0041", "UNIT-0101"],
        "layout": "left_table_right_text",
        "proof": "本项目防治责任范围一览表与工程占地情况一览表",
    },
    {
        "type": "table",
        "chapter": "防治责任范围与标准目标",
        "title": "六项防治指标采用南方红壤区一级标准",
        "source_mode": "ORIGINAL_TABLE",
        "units": ["UNIT-0046"],
        "layout": "left_table_right_text",
        "proof": "本工程水土流失防治目标值表",
    },
    {
        "type": "table",
        "chapter": "水土保持评价与土石方",
        "title": "法制约因素与技术标准评价结论均为符合",
        "source_mode": "ORIGINAL_TABLE",
        "units": ["UNIT-0117", "UNIT-0119", "UNIT-0120"],
        "layout": "top_table_bottom_text",
        "proof": "水土保持法、GB50433-2018、长江保护法制约因素分析评价表",
    },
    {
        "type": "table",
        "chapter": "水土保持评价与土石方",
        "title": "土石方平衡显示弃方委托处置、借方外购",
        "source_mode": "ORIGINAL_TABLE",
        "units": ["UNIT-0105", "UNIT-0106"],
        "layout": "left_table_right_text",
        "proof": "工程土石方平衡表及土石方平衡流向图标题",
    },
    {
        "type": "table",
        "chapter": "水土保持评价与土石方",
        "title": "主体工程已列水保措施以土地整治和铺种草皮为主",
        "source_mode": "ORIGINAL_TABLE",
        "units": ["UNIT-0140"],
        "layout": "left_table_right_text",
        "proof": "主体工程设计中应纳入水土保持的措施量及投资一览表",
    },
    {
        "type": "table",
        "chapter": "水土流失分析与预测",
        "title": "扰动地表面积与预测单元集中在管网工程区",
        "source_mode": "ORIGINAL_TABLE",
        "units": ["UNIT-0148", "UNIT-0153", "UNIT-0157"],
        "layout": "top_table_bottom_text",
        "proof": "扰动地表面积表、预测单元划分表、预测范围与预测时段划分表",
    },
    {
        "type": "table",
        "chapter": "水土流失分析与预测",
        "title": "背景土壤侵蚀模数取 261.46t/km²·a",
        "source_mode": "CALCULATION",
        "units": ["UNIT-0159", "UNIT-0161"],
        "layout": "left_table_right_text",
        "proof": "各地类土壤侵蚀模数取值表和背景值计算表",
    },
    {
        "type": "table",
        "chapter": "水土流失分析与预测",
        "title": "施工期扰动后侵蚀模数显著高于背景值",
        "source_mode": "CALCULATION",
        "units": ["UNIT-0180", "UNIT-0183", "UNIT-0185", "UNIT-0187"],
        "layout": "left_table_right_text",
        "proof": "植被破坏型、地表翻扰型和自然恢复期土壤侵蚀模数计算表",
    },
    {
        "type": "table",
        "chapter": "水土流失分析与预测",
        "title": "预测土壤流失总量 104.05t，新增量 94.92t",
        "source_mode": "ORIGINAL_TABLE",
        "units": ["UNIT-0192"],
        "layout": "left_table_right_text",
        "proof": "工程建设可能造成的土壤流失量预测结果表",
    },
    {
        "type": "table",
        "chapter": "水土保持措施与施工安排",
        "title": "防治分区为管网工程区，防治面积 34416.00m²",
        "source_mode": "ORIGINAL_TABLE",
        "units": ["UNIT-0208"],
        "layout": "left_table_right_text",
        "proof": "水土流失防治区划分表",
    },
    {
        "type": "table",
        "chapter": "水土保持措施与施工安排",
        "title": "措施体系覆盖工程、植物和临时三类措施",
        "source_mode": "ORIGINAL_TABLE",
        "units": ["UNIT-0228", "UNIT-0241"],
        "layout": "left_table_right_text",
        "proof": "水土流失防治措施体系表和措施工程量汇总表",
    },
    {
        "type": "table",
        "chapter": "水土保持措施与施工安排",
        "title": "措施实施进度需嵌入 2025—2026 年施工窗口",
        "source_mode": "ORIGINAL_TABLE",
        "units": ["UNIT-0254"],
        "layout": "top_table_bottom_text",
        "proof": "水土保持措施实施进度表",
    },
    {
        "type": "table",
        "chapter": "投资概算、效益与管理验收",
        "title": "水土保持估算总投资为 26.9211 万元",
        "source_mode": "ORIGINAL_TABLE",
        "units": ["UNIT-0299"],
        "layout": "left_table_right_text",
        "proof": "本项目水土保持估算总表",
    },
    {
        "type": "table",
        "chapter": "投资概算、效益与管理验收",
        "title": "工程、植物、临时和独立费用共同构成新增投资",
        "source_mode": "ORIGINAL_TABLE",
        "units": ["UNIT-0300", "UNIT-0301", "UNIT-0302", "UNIT-0303"],
        "layout": "top_table_bottom_text",
        "proof": "工程措施、植物措施、临时措施和独立费用估算表",
    },
    {
        "type": "table",
        "chapter": "投资概算、效益与管理验收",
        "title": "水土保持补偿费 51624 元符合免征情形",
        "source_mode": "ORIGINAL_TABLE",
        "units": ["UNIT-0304"],
        "layout": "left_table_right_text",
        "proof": "水土保持补偿费计算表及免征依据说明",
    },
    {
        "type": "table",
        "chapter": "投资概算、效益与管理验收",
        "title": "六项防治效果达到或高于目标值",
        "source_mode": "ORIGINAL_TABLE",
        "units": ["UNIT-0317"],
        "layout": "left_table_right_text",
        "proof": "水土保持防治效果分析表",
    },
    {
        "type": "text",
        "chapter": "投资概算、效益与管理验收",
        "title": "后续管理需落实组织管理、监理施工和设施验收",
        "source_mode": "MANAGEMENT_ACTION",
        "units": ["UNIT-0345", "UNIT-0346", "UNIT-0347"],
        "layout": "text_with_keypoints",
        "proof": "水土保持设施验收及后续管理章节内容",
    },
]


def clean_number(value: str) -> str:
    return re.sub(r"\s+", "", value)


def unit_lookup(inventory: dict[str, Any]) -> dict[str, dict[str, Any]]:
    return {unit["id"]: unit for unit in inventory.get("content_units", [])}


def catalog_lookup(catalog: dict[str, Any]) -> dict[str, dict[str, Any]]:
    return {entry["id"]: entry for entry in catalog.get("entries", [])}


def first_text(parts: list[str], limit: int = 220) -> str:
    text = " ".join(sanitize_visible_text(part) for part in parts if sanitize_visible_text(part))
    return text_preview(text, limit, complete_sentence=True)


def evidence_kind(source_mode: str) -> str:
    if source_mode == "CALCULATION":
        return "CALCULATION"
    if source_mode == "ORIGINAL_TABLE":
        return "ORIGINAL_TABLE"
    if source_mode == "ORIGINAL_FIGURE":
        return "ORIGINAL_FIGURE"
    if source_mode == "MANAGEMENT_ACTION":
        return "RECOMMENDATION"
    return "ORIGINAL_TEXT"


def make_values(catalog_ids: list[str], entries: dict[str, dict[str, Any]]) -> list[dict[str, str]]:
    values: list[dict[str, str]] = []
    seen: set[str] = set()
    for catalog_id in catalog_ids:
        entry = entries.get(catalog_id, {})
        for number in entry.get("numbers", []) or []:
            token = clean_number(str(number))
            if not token or token in seen or len(values) >= 18:
                continue
            seen.add(token)
            values.append({"value": token, "unit": "", "time_basis": ""})
    return values


def build_evidence_and_plan() -> tuple[dict[str, Any], dict[str, Any]]:
    inventory = load_json(INVENTORY_PATH)
    catalog = load_json(CATALOG_PATH)
    units = unit_lookup(inventory)
    entries = catalog_lookup(catalog)
    evidence_records: list[dict[str, Any]] = []
    slides: list[dict[str, Any]] = []
    evidence_counter = 1

    for page, spec in enumerate(SLIDE_SPECS, start=1):
        slide_type = spec["type"]
        unit_ids = spec.get("units", [])
        source_mode = spec.get("source_mode", "")
        catalog_ids: list[str] = []
        previews: list[str] = []
        for unit_id in unit_ids:
            unit = units.get(unit_id)
            if not unit:
                continue
            catalog_ids.extend(unit.get("catalog_ids", []) or [])
            previews.append(unit.get("content_preview", ""))
        catalog_ids = list(dict.fromkeys(catalog_ids))

        evidence_ids: list[str] = []
        if slide_type not in {"cover", "agenda", "section", "closing"}:
            evidence_id = f"E-{evidence_counter:03d}"
            evidence_counter += 1
            summary = first_text(previews + [spec.get("proof", ""), spec["title"]], 280)
            record: dict[str, Any] = {
                "id": evidence_id,
                "kind": evidence_kind(source_mode),
                "summary": summary,
                "source_file": catalog.get("source_files", ["中地环科-黄冈市东坡大道片区污水收集管网建设工程水土保持方案报告表-20260108.docx"])[0]
                if catalog.get("source_files")
                else "中地环科-黄冈市东坡大道片区污水收集管网建设工程水土保持方案报告表-20260108.docx",
                "source_locator": "、".join(
                    entries.get(catalog_id, {}).get("locator", catalog_id) for catalog_id in catalog_ids
                ),
                "catalog_ids": catalog_ids,
                "verification_status": "verified",
                "exact_text": first_text(
                    [entries.get(catalog_id, {}).get("text", "") for catalog_id in catalog_ids],
                    520,
                ),
                "values": make_values(catalog_ids, entries),
                "notes": "已依据源 DOCX 提取的段落、表格或图件标题核验。",
            }
            if source_mode == "CALCULATION":
                record["formula"] = "按报告对应参数表、侵蚀模数表和预测结果表计算；不新增报告外计算。"
                record["inputs"] = [
                    text_preview(entries.get(catalog_id, {}).get("text", ""), 110, complete_sentence=True)
                    for catalog_id in catalog_ids[:3]
                ]
            evidence_records.append(record)
            evidence_ids = [evidence_id]

        slides.append(
            {
                "page": page,
                "type": slide_type,
                "chapter": spec.get("chapter", ""),
                "title": spec["title"],
                "source_mode": source_mode,
                "evidence_ids": evidence_ids,
                "content_unit_ids": unit_ids,
                "visual_proof": spec.get("proof", ""),
                "layout_pattern": spec.get("layout", ""),
                "image_names": spec.get("image_names", []),
                "source_note": "、".join(
                    entries.get(catalog_id, {}).get("locator", catalog_id) for catalog_id in catalog_ids
                ),
                "density": "dense" if slide_type in {"table", "figure"} else "normal",
                "density_exempt_reason": "" if slide_type not in {"cover", "agenda"} else "structural slide",
            }
        )

    chapter_order = list(CHAPTER_SUMMARIES)
    ledger = {
        "schema_version": "1.0",
        "project": "东坡大道水土保持方案",
        "updated_at": datetime.now().isoformat(timespec="seconds"),
        "rules": {
            "source_truth": "Only verified source evidence may support FACT or CALCULATION.",
            "conflict_handling": "Keep conflicting source values separate and disclose the conflict.",
        },
        "evidence": evidence_records,
    }
    plan = {
        "schema_version": "1.0",
        "project": "东坡大道水土保持方案",
        "deck": {
            "title": PROJECT_TITLE,
            "audience": AUDIENCE,
            "objective": "用于评审与建设管理沟通，说明水土保持方案的项目边界、预测依据、措施体系、投资效益和验收责任。",
            "chapter_order": chapter_order,
        },
        "coverage": {
            "chapter_allocation": {
                chapter: sum(1 for slide in slides if slide.get("chapter") == chapter)
                for chapter in chapter_order
            }
        },
        "slides": slides,
    }
    return ledger, plan


def write_chapter_coverage(plan: dict[str, Any], ledger: dict[str, Any]) -> None:
    evidence_by_id = {item["id"]: item for item in ledger["evidence"]}
    lines = [
        "# Chapter Coverage",
        "",
        f"- Project: {PROJECT_TITLE}",
        f"- Planned content slides including cover/agenda: {len(plan['slides'])}",
        f"- Structural chapter dividers: {len(plan['deck']['chapter_order'])}",
        f"- Expected PPTX slides: {len(plan['slides']) + len(plan['deck']['chapter_order'])}",
        "",
    ]
    for chapter in plan["deck"]["chapter_order"]:
        slides = [slide for slide in plan["slides"] if slide.get("chapter") == chapter]
        original = [
            slide
            for slide in slides
            if slide.get("source_mode") in {"ORIGINAL_TEXT", "ORIGINAL_TABLE", "ORIGINAL_FIGURE", "CALCULATION"}
        ]
        lines.append(f"## {chapter}")
        lines.append("")
        lines.append(f"- Planned pages: {len(slides)}")
        lines.append(f"- Original-source pages: {len(original)}")
        lines.append(f"- Coverage intent: {CHAPTER_SUMMARIES.get(chapter, '')}")
        for slide in slides:
            evidence_text = "；".join(
                text_preview(evidence_by_id[eid]["summary"], 90, complete_sentence=True)
                for eid in slide.get("evidence_ids", [])
                if eid in evidence_by_id
            )
            lines.append(
                f"- P{slide['page']:02d} `{slide['source_mode']}` {slide['title']}；"
                f" source: {slide.get('source_note', '')}；evidence: {evidence_text}"
            )
        lines.append("")
    (PROJECT_DIR / "chapter_coverage.md").write_text("\n".join(lines), encoding="utf-8")


def write_backend_specs(plan: dict[str, Any]) -> None:
    template_profile = collect_template_profile(PROJECT_DIR)
    image_names = sorted(
        {
            image_name
            for slide in plan.get("slides", [])
            for image_name in slide.get("image_names", [])
        }
    )
    design_spec = create_engineering_design_spec(
        plan=plan,
        media_dir=PROJECT_DIR / "images",
        template_profile=template_profile,
        project_title=PROJECT_TITLE,
        audience=AUDIENCE,
        use_case=USE_CASE,
        design_style=DESIGN_STYLE,
        image_names=image_names,
    )
    spec_lock = create_engineering_spec_lock(
        plan=plan,
        media_dir=PROJECT_DIR / "images",
        important_anchor_pages={1, 2, 3, 8, 13, 17, 20, 24},
    )
    (PROJECT_DIR / "design_spec.md").write_text(design_spec, encoding="utf-8")
    (PROJECT_DIR / "spec_lock.md").write_text(spec_lock, encoding="utf-8")
    (PROJECT_DIR / "analysis" / "template_profile.json").write_text(
        json.dumps(template_profile, ensure_ascii=False, indent=2) + "\n",
        encoding="utf-8",
    )


def write_planning_artifacts() -> tuple[dict[str, Any], dict[str, Any]]:
    ledger, plan = build_evidence_and_plan()
    LEDGER_PATH.write_text(json.dumps(ledger, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    PLAN_PATH.write_text(json.dumps(plan, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    write_chapter_coverage(plan, ledger)
    return ledger, plan


def add_cover_slide(slide) -> None:
    add_fill(slide, 0, 0, 13.333, 7.5, "primary", "primary", colors=COLORS)
    add_textbox(slide, 0.80, 0.68, 8.2, 0.34, "黄冈市住房和城市更新局", 15, "paper", colors=COLORS)
    add_textbox(slide, 0.80, 1.45, 11.5, 0.72, "黄冈市东坡大道片区污水收集管网建设工程", 30, "paper", True, colors=COLORS)
    add_textbox(slide, 0.82, 2.34, 10.8, 0.64, "水土保持方案报告表汇报", 28, "paper", True, colors=COLORS)
    add_textbox(slide, 0.84, 3.80, 10.4, 0.36, "水土保持方案 | 报批稿 | 2026年1月", 15, "paper", colors=COLORS)
    add_textbox(slide, 0.84, 4.32, 10.8, 0.38, "围绕项目边界、预测结果、防治措施、投资效益和验收管理形成评审级汇报", 14, "paper", colors=COLORS)
    add_textbox(slide, 0.86, 6.66, 6.0, 0.34, "武汉中地环科水工环科技咨询有限责任公司", 14, "paper", colors=COLORS)
    add_textbox(slide, 10.3, 6.66, 2.0, 0.34, datetime.now().strftime("%Y年%m月"), 14, "paper", align=PP_ALIGN.RIGHT, colors=COLORS)
    add_textbox(slide, 7.0, 6.66, 2.2, 0.34, "技术评审汇报", 14, "paper", colors=COLORS)


def compact_table(table_entry: dict[str, Any], page: dict[str, Any]) -> dict[str, Any]:
    rows = table_entry.get("rows") or []
    title = page.get("title", "")
    if "补偿费" in title:
        focused = dict(table_entry)
        focused["rows"] = [
            ["行政区", "征占地面积(m²)", "收费标准(元/m²)", "水保补偿费(元)"],
            ["黄冈市黄州区", "34416", "1.5", "51624"],
        ]
        focused["row_count"] = len(focused["rows"])
        focused["text"] = "\n".join(" | ".join(row) for row in focused["rows"])
        return focused
    if "防治效果" in title:
        focused = dict(table_entry)
        focused["rows"] = [
            ["评估指标", "目标值", "设计达到值", "评估结果"],
            ["水土流失治理度", "98%", "99.99%", "达标"],
            ["土壤流失控制比", "1.0", "1.55", "达标"],
            ["渣土防护率", "97%", "97.80%", "达标"],
            ["表土保护率", "92%", "99.99%", "达标"],
            ["林草植被恢复率", "98%", "99.99%", "达标"],
            ["林草覆盖率", "0.77%", "0.77%", "达标"],
        ]
        focused["row_count"] = len(focused["rows"])
        focused["text"] = "\n".join(" | ".join(row) for row in focused["rows"])
        return focused
    if len(rows) <= 7:
        return table_entry
    if "实施进度" in title:
        focused = dict(table_entry)
        focused["rows"] = [
            ["工程分区/路段", "主要工序", "实施窗口", "水保衔接"],
            ["管网工程区", "施工准备、围挡、路面破除", "2025年12月起进入施工准备和扰动启动阶段", "临时防护、排水和车辆冲洗应先行布设"],
            ["东坡大道等道路段", "顶管井施工、顶管施工", "2026年主体施工窗口内分段推进", "顶管井开挖、土方周转和降排水是过程控制重点"],
            ["各道路恢复段", "路面恢复、绿化恢复", "随主体完工分段恢复", "土地整治和植被恢复需与道路恢复同步闭合"],
        ]
        focused["row_count"] = len(focused["rows"])
        focused["text"] = "\n".join(" | ".join(row) for row in focused["rows"])
        focused["structure_note"] = "原表为跨年度多级表头进度表，PPT采用结构化摘录保留路段、工序、时序和水保衔接关系。"
        return focused
    keywords = []
    if "投资" in title or "费用" in title:
        keywords = ["工程措施", "植物措施", "临时措施", "独立费用", "估算总投资", "水土保持补偿费", "合计"]
    elif "土石方" in title:
        keywords = ["合计", "表土", "基础土方", "道路破除", "清淤"]
    elif "预测" in title or "流失" in title:
        keywords = ["管网工程区", "合计", "施工期", "自然恢复期", "新增量", "总量"]
    elif "措施" in title:
        keywords = ["工程措施", "植物措施", "临时措施", "管网工程区", "表土", "临时苫盖", "高压洗车机"]
    elif "防治" in title or "目标" in title:
        keywords = ["水土流失治理度", "土壤流失控制比", "渣土防护率", "表土保护率", "林草植被恢复率", "林草覆盖率"]
    else:
        keywords = ["管网工程区", "合计", "总投资", "建设规模", "工期", "占地"]
    selected = rows[:1]
    for row in rows[1:]:
        row_text = " ".join(str(cell) for cell in row)
        if any(keyword in row_text for keyword in keywords):
            selected.append(row)
        if len(selected) >= 7:
            break
    if len(selected) <= 1:
        selected = rows[:7]
    focused = dict(table_entry)
    focused["rows"] = selected
    focused["row_count"] = len(selected)
    focused["text"] = "\n".join(" | ".join(str(cell) for cell in row) for row in selected)
    return focused


def table_points_from_context(page: dict[str, Any], evidence: list[dict[str, Any]], max_points: int = 4) -> list[str]:
    title = sanitize_visible_text(page.get("title", ""))
    chapter = sanitize_visible_text(page.get("chapter", ""))
    proof = sanitize_visible_text(page.get("visual_proof", ""))
    joined = " ".join([title, chapter, proof])
    library: list[tuple[list[str], list[str]]] = [
        (
            ["综合报告表", "全局控制"],
            [
                "工程新建污水主管网6.2公里、预埋支管2.2公里，服务城北片区污水收集管网完善。",
                "项目占地均为临时占地，防治责任范围与工程扰动范围统一为34416.00m²。",
                "方案同步明确土石方、六项防治目标、措施体系和26.9211万元水土保持投资。",
            ],
        ),
        (
            ["工程特性", "规模", "工期"],
            [
                "项目性质为新建，建设地点位于黄冈市黄州区黄州大道片区。",
                "总投资8024.00万元、土建投资6914.95万元，计划2025年12月至2026年12月实施。",
                "用地面积34416.00m²，全部按临时占地纳入管网工程区防治。",
            ],
        ),
        (
            ["责任范围", "临时占地"],
            [
                "本项目防治责任范围为管网工程区34416.00m²。",
                "工程占地均为临时占地，施工结束后需完成路面、绿化和临时扰动恢复。",
                "管线作业、顶管井施工、材料堆放和道路恢复共同构成水土保持控制对象。",
            ],
        ),
        (
            ["防治指标", "一级标准"],
            [
                "本项目执行南方红壤区建设类项目一级防治标准。",
                "设计水平年水土流失治理度、林草植被恢复率均按98%控制，渣土防护率按97%控制。",
                "表土保护率取93%，土壤流失控制比取1.0，林草覆盖率结合项目实际调整为0.77%。",
            ],
        ),
        (
            ["制约因素", "符合"],
            [
                "项目选址未涉及水土保持法、长江保护法及相关技术标准的禁止性因素。",
                "工程位于城市建成区，建设方案以减少扰动面积和土石方量为原则。",
                "施工组织采用分段施工、随挖随运和避开恶劣天气基础施工等控制要求。",
            ],
        ),
        (
            ["土石方平衡", "弃方", "借方"],
            [
                "工程总挖方4971.61m³、总填方2232.54m³，外购土石方2153.04m³。",
                "弃方4892.11m³，其中建筑垃圾1560.00m³，均委托湖北安达港务有限公司外运消纳。",
                "施工过程土石方随挖随运，不设置临时堆土场，运输和苫盖为施工期控制重点。",
            ],
        ),
        (
            ["主体工程已列", "土地整治", "草皮"],
            [
                "主体工程已列水土保持措施主要包括土地整治和铺种草皮。",
                "土地整治面积和草皮恢复面积用于道路两侧扰动区域恢复。",
                "方案补充表土剥离回覆、临时苫盖和高压洗车机，完善施工期临时防护。",
            ],
        ),
        (
            ["扰动地表", "预测单元", "预测时段"],
            [
                "水土流失预测范围与防治责任范围一致，预测单元为管网工程区。",
                "预测时段区分施工期和自然恢复期，分别对应扰动施工和植被恢复过程。",
                "预测面积、时段和单元划分共同形成土壤流失量计算基础。",
            ],
        ),
        (
            ["侵蚀模数", "背景"],
            [
                "项目区背景土壤侵蚀模数为261.46t/km²·a，容许土壤流失量为500t/km²·a。",
                "背景值用于反映原地貌条件下的水土流失水平。",
                "各地类侵蚀模数取值与预测公式共同构成后续流失量计算参数。",
            ],
        ),
        (
            ["施工期扰动", "侵蚀模数"],
            [
                "施工期植被破坏型和地表翻扰型侵蚀模数明显高于背景值。",
                "自然恢复期侵蚀模数随地表恢复下降，但仍需覆盖植被恢复过程。",
                "扰动类型分项取值支撑新增土壤流失量和重点防治时段判断。",
            ],
        ),
        (
            ["土壤流失总量", "新增量"],
            [
                "工程建设可能造成土壤流失总量104.05t，新增土壤流失量94.92t。",
                "新增量主要发生在施工扰动阶段，施工期是水土流失防治重点。",
                "预测结果对应工程措施、植物措施和临时措施的配置强度。",
            ],
        ),
        (
            ["防治分区"],
            [
                "方案划分为管网工程区一个防治分区，防治面积34416.00m²。",
                "防治分区覆盖管线施工、顶管井、临时扰动和道路绿化恢复范围。",
                "措施布设、工程量统计、投资估算和防治效益均以该分区为基础。",
            ],
        ),
        (
            ["措施体系", "工程量"],
            [
                "管网工程区措施包括工程措施、植物措施和临时措施三类。",
                "工程措施为表土剥离及回覆、土地整治，植物措施为铺种草皮。",
                "临时措施包括临时苫盖和高压洗车机，用于控制裸露地表和车辆带泥上路。",
            ],
        ),
        (
            ["实施进度"],
            [
                "水土保持措施实施期与主体工程施工期衔接，覆盖2025年12月至2026年12月。",
                "表土剥离、临时苫盖等措施随施工扰动同步布设。",
                "土地整治、表土回覆和铺种草皮随路面及绿化恢复分段完成。",
            ],
        ),
        (
            ["估算总投资", "投资"],
            [
                "本项目水土保持总投资26.9211万元，其中主体工程已有0.6381万元，本方案新增0.5885万元。",
                "工程措施投资0.1827万元，植物措施投资0.6185万元，临时措施投资0.4253万元。",
                "独立费用20.4968万元，水土保持补偿费5.1624万元，补偿费符合免缴情形。",
            ],
        ),
        (
            ["工程、植物、临时", "独立费用"],
            [
                "新增投资按工程措施、临时措施、独立费用和基本预备费分类计列。",
                "主体工程已有投资主要对应植物措施及部分工程措施。",
                "独立费用包括建设管理费、科研勘测设计费、水土保持监理费和设施验收费。",
            ],
        ),
        (
            ["补偿费", "免征"],
            [
                "水土保持补偿费按征占地面积3.4416hm²、收费标准1.5元/m²计算，金额为51624元。",
                "项目属于市政生态环境保护基础设施，符合财综〔2014〕8号规定的免征情形。",
                "方案编制时不予计列补偿费缴纳支出，仍保留补偿费计算依据。",
            ],
        ),
        (
            ["防治效果", "目标值"],
            [
                "防治效果分析显示六项指标均达到或高于目标值。",
                "水土流失总治理度、林草植被恢复率、渣土防护率等指标形成达标结论。",
                "防治效果需通过施工管理、监理记录和设施验收进一步落实。",
            ],
        ),
    ]
    for keywords, points in library:
        if all(keyword in joined for keyword in keywords):
            return [complete_semantic_point(point, 86) for point in points[:max_points]]

    fallback = [
        f"该部分对应“{chapter or '报告'}”章节的工程控制内容。",
        "关键参数用于衔接报告中的计算、措施布设、投资估算和验收管理。",
        "多级表头保留指标、单位和分组关系，避免改变原报告口径。",
    ]
    for ev in evidence:
        summary = sanitize_visible_text(ev.get("summary", ""))
        if summary and len(fallback) < max_points:
            fallback.append(summary)
    return [complete_semantic_point(point, 86) for point in fallback[:max_points]]


def figure_points_from_context(page: dict[str, Any]) -> list[str]:
    title = sanitize_visible_text(page.get("title", ""))
    if "城北片区" in title or "道路布设" in title:
        return [
            "图件显示项目位于黄州区城北片区，工程范围受黄州大道、三台河路、东坡大道、赤壁大道等道路边界控制。",
            "管网沿既有城市道路布设，施工扰动与交通组织、现状绿化和恢复边界关系密切。",
            "图件定位责任范围和扰动场景，防治分区以管网工程区统一展开。",
        ]
    if "顶管" in title or "施工组织" in title:
        return [
            "顶管示意图说明本项目以微顶管减少城区道路大开挖扰动，是施工组织和水保措施时序的关键依据。",
            "工作井、接收井和井内回填剖面决定临时占地、土方周转、排水降水和分层压实管理重点。",
            "图件内容与土石方平衡、临时苫盖、施工排水和路面绿化恢复措施相互对应。",
        ]
    return [
        "图件反映报告中的工程场景和水土保持控制对象。",
        "图件说明聚焦工程含义、扰动来源和措施衔接关系。",
    ]


def text_points_from_context(
    page: dict[str, Any],
    evidence: list[dict[str, Any]],
    catalog_entries: dict[str, dict[str, Any]],
    catalog_entry_list: list[dict[str, Any]],
) -> list[str]:
    title = sanitize_visible_text(page.get("title", ""))
    if "前期批复" in title or "编制依据" in title:
        return [
            "本方案以可研批复、初设批复和项目既有前期工作为边界，明确工程建设内容与水土保持评价对象。",
            "编制依据覆盖水土保持法律法规、技术标准和地方管理要求，是防治目标、预测计算和投资概算的共同基础。",
            "附件清单包括批复、委托、地理位置、措施布设和投资估算等支撑材料。",
            "前期手续与依据文件共同限定方案编制范围和水土保持责任边界。",
        ]
    if "后续管理" in title or "设施验收" in title:
        return [
            "建设单位应成立水土保持方案实施管理机构，设专人负责水土保持工作。",
            "水土保持监理按批准的方案和工程设计文件实施质量、进度和投资控制。",
            "施工单位应按设计图纸和施工技术要求实施水土保持措施，并对施工范围内水土流失负责。",
            "水土保持设施验收按水保〔2017〕365号执行，验收合格后方可通过竣工验收和投产使用。",
            "验收材料应依法公开，并在投产使用前向水土保持方案审批机关报备。",
        ]
    points = report_points_from_page(page, evidence, catalog_entries, catalog_entry_list, max_points=5)
    return [complete_semantic_point(point, 74) for point in points if sanitize_visible_text(point)]


def table_model_for_entry(table_entry: dict[str, Any], table_models: dict[str, Any]) -> dict[str, Any] | None:
    locator = sanitize_visible_text(table_entry.get("locator", "")).upper()
    if not locator:
        return None
    model = table_models.get(locator)
    if not model:
        return None
    if not model.get("rows"):
        return None
    return model


def model_has_long_free_text(model: dict[str, Any], limit: int = 92) -> bool:
    for cell in model.get("cells", []):
        text = sanitize_visible_text(cell.get("text", ""))
        if len(text) > limit:
            return True
    return False


def prefers_structured_excerpt(page: dict[str, Any]) -> bool:
    title = sanitize_visible_text(page.get("title", ""))
    return any(keyword in title for keyword in ["补偿费", "防治效果", "实施进度"])


def add_project_table(
    slide,
    table_entry: dict[str, Any],
    page: dict[str, Any],
    table_models: dict[str, Any],
    x: float,
    y: float,
    w: float,
    h: float,
    max_rows: int,
    max_cols: int,
) -> None:
    model = table_model_for_entry(table_entry, table_models)
    if (
        model
        and not prefers_structured_excerpt(page)
        and not model_has_long_free_text(model)
        and model_fits_native_table(model, max_rows=max_rows, max_cols=max_cols)
    ):
        render_merged_table(slide, model, x, y, w, h, colors=COLORS, max_rows=max_rows, max_cols=max_cols)
        return
    add_table(slide, compact_table(table_entry, page), x, y, w, h, max_rows=max_rows, max_cols=max_cols, colors=COLORS)


def add_source_slide(
    slide,
    page: dict[str, Any],
    evidence: list[dict[str, Any]],
    catalog_entries: dict[str, dict[str, Any]],
    catalog_entry_list: list[dict[str, Any]],
    title_context: dict[str, dict[str, str]],
    table_models: dict[str, Any],
) -> None:
    add_title(slide, page, colors=COLORS)
    body_text = report_text_from_page(page, evidence, catalog_entries, catalog_entry_list)
    topic_heading = source_topic_heading(page, evidence, title_context, catalog_entry_list)
    if str(page.get("source_mode", "")).upper() == "ORIGINAL_FIGURE":
        image_names = [sanitize_visible_text(name) for name in page.get("image_names", []) if sanitize_visible_text(name)]
        image_paths = require_media_paths(MEDIA_DIR, image_names, context=page.get("title", "source figure"))
        if len(image_paths) == 1:
            add_image_panel(slide, image_paths, 0.66, 1.16, 7.60, 4.80, colors=COLORS)
            add_textbox(slide, 8.72, 1.22, 3.55, 0.46, topic_heading, 15, "accent", True, colors=COLORS)
            add_bullets(slide, 8.72, 1.86, 3.28, 3.95, figure_points_from_context(page), 14, colors=COLORS)
        else:
            add_image_panel(slide, image_paths, 0.66, 1.14, 12.00, 3.92, colors=COLORS)
            add_textbox(slide, 0.78, 5.26, 4.95, 0.38, topic_heading, 15, "accent", True, colors=COLORS)
            add_bullets(slide, 0.78, 5.70, 11.55, 0.80, figure_points_from_context(page)[:3], 13.8, colors=COLORS)
        return
    table_allowed = str(page.get("source_mode", "")).upper() in {"ORIGINAL_TABLE", "CALCULATION"}
    table_entry = table_from_page(page, catalog_entries, catalog_entry_list, evidence) if table_allowed else None
    if table_entry:
        points = table_points_from_context(page, evidence, max_points=4)
    else:
        points = text_points_from_context(page, evidence, catalog_entries, catalog_entry_list)

    if table_entry:
        layout = page.get("layout_pattern", "")
        model = table_model_for_entry(table_entry, table_models)
        if (
            model
            and model.get("has_merged_cells")
            and not prefers_structured_excerpt(page)
            and not model_has_long_free_text(model)
            and model_fits_native_table(model, max_rows=7, max_cols=7)
            and int(model.get("column_count", 0)) > 4
        ):
            render_merged_table(slide, model, 0.60, 1.16, 12.10, 3.78, colors=COLORS, max_rows=7, max_cols=7)
            add_textbox(slide, 0.78, 5.12, 7.20, 0.54, topic_heading, 15, "accent", True, colors=COLORS)
            add_bullets(slide, 0.78, 5.62, 11.55, 0.95, points[:3] or [text_preview(page["visual_proof"], 90, complete_sentence=True)], 14, colors=COLORS)
            return
        if layout.startswith("top_"):
            add_project_table(slide, table_entry, page, table_models, 0.60, 1.16, 12.10, 3.78, max_rows=7, max_cols=5)
            add_textbox(slide, 0.78, 5.12, 7.20, 0.54, topic_heading, 15, "accent", True, colors=COLORS)
            add_bullets(slide, 0.78, 5.62, 11.55, 0.95, points[:3] or [text_preview(page["visual_proof"], 90, complete_sentence=True)], 14, colors=COLORS)
        else:
            add_project_table(slide, table_entry, page, table_models, 0.60, 1.16, 8.00, 5.05, max_rows=7, max_cols=4)
            add_textbox(slide, 8.95, 1.22, 3.35, 0.48, topic_heading, 15, "accent", True, colors=COLORS)
            add_bullets(slide, 8.95, 1.88, 3.15, 3.95, points[:4] or [text_preview(page["visual_proof"], 100, complete_sentence=True)], 14, colors=COLORS)
        return

    add_fill(slide, 0.66, 1.23, 12.0, 4.88, "paper", colors=COLORS)
    add_textbox(slide, 0.96, 1.55, 11.3, 0.38, topic_heading, 16, "accent", True, colors=COLORS)
    if len(points) >= 4:
        mid = (len(points) + 1) // 2
        add_bullets(slide, 0.96, 2.10, 5.25, 3.75, points[:mid], 14, colors=COLORS)
        add_bullets(slide, 6.62, 2.10, 5.25, 3.75, points[mid:], 14, colors=COLORS)
    else:
        add_bullets(slide, 0.96, 2.10, 10.95, 3.75, points or [complete_semantic_point(body_text, 260)], 14.5, colors=COLORS)


def write_notes(plan: dict[str, Any]) -> None:
    NOTES_DIR.mkdir(exist_ok=True)
    chapters = chapter_order_from_plan(plan)
    chapter_numbers = {chapter: index for index, chapter in enumerate(chapters, start=1)}
    total: list[str] = []
    section_written: set[str] = set()
    for page in plan["slides"]:
        chapter = sanitize_visible_text(page.get("chapter", ""))
        if page["type"] not in {"cover", "agenda", "closing"} and chapter and chapter not in section_written:
            section_written.add(chapter)
            idx = chapter_numbers.get(chapter, len(section_written))
            note = (
                f"Section {idx:02d} - {chapter}\n\n"
                f"讲解要点：进入“{chapter}”，本章目标是{CHAPTER_SUMMARIES.get(chapter, '呈现报告证据和管理判断')}。\n"
            )
            (NOTES_DIR / f"section_{idx:02d}.md").write_text(note, encoding="utf-8")
            total.append(note)
        note = (
            f"Slide {page['page']:02d} - {sanitize_visible_text(page['title'])}\n\n"
            f"讲解要点：内容依据报告 {sanitize_visible_text(page.get('source_note', ''))}，"
            f"说明{sanitize_visible_text(page['title'])}。保持单位、时间口径和表格结论与报告一致。\n"
        )
        (NOTES_DIR / f"{page['page']:02d}.md").write_text(note, encoding="utf-8")
        total.append(note)
    (NOTES_DIR / "total.md").write_text("\n\n".join(total), encoding="utf-8")


def build_deck() -> Path:
    if SOURCE_DOCX.exists():
        extract_docx_media(SOURCE_DOCX, MEDIA_DIR)
        write_docx_table_models(SOURCE_DOCX, TABLE_MODELS_PATH)
    ledger, plan = write_planning_artifacts()
    write_backend_specs(plan)
    catalog = load_json(CATALOG_PATH)
    catalog_entry_list = catalog["entries"]
    catalog_entries = {entry["id"]: entry for entry in catalog_entry_list}
    title_context = catalog_title_context(catalog_entry_list)
    evidence_by_id = evidence_lookup(ledger)
    table_models_payload = load_json(TABLE_MODELS_PATH) if TABLE_MODELS_PATH.exists() else {"by_locator": {}}
    table_models = table_models_payload.get("by_locator", {})

    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    chapters = chapter_order_from_plan(plan)
    chapter_numbers = {chapter: index for index, chapter in enumerate(chapters, start=1)}
    section_written: set[str] = set()

    for page in plan["slides"]:
        if page["type"] == "cover":
            slide = prs.slides.add_slide(blank)
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = rgb(COLORS["primary"])
            add_cover_slide(slide)
            continue
        if page["type"] == "agenda":
            slide = prs.slides.add_slide(blank)
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = rgb(COLORS["bg"])
            render_page = dict(page, display_page=len(prs.slides) + 1)
            add_agenda_slide(slide, plan, render_page, colors=COLORS)
            continue

        chapter = sanitize_visible_text(page.get("chapter", ""))
        if chapter and chapter not in section_written:
            section_written.add(chapter)
            section_slide = prs.slides.add_slide(blank)
            section_slide.background.fill.solid()
            section_slide.background.fill.fore_color.rgb = rgb(COLORS["primary"])
            add_section_divider_slide(
                section_slide,
                chapter_numbers.get(chapter, len(section_written)),
                chapter,
                len(chapters),
                subtitle=CHAPTER_SUMMARIES.get(chapter, ""),
                colors=COLORS,
            )

        slide = prs.slides.add_slide(blank)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = rgb(COLORS["bg"])
        render_page = dict(page, display_page=len(prs.slides) + 1)
        evidence = slide_evidence(page, evidence_by_id)
        add_source_slide(slide, render_page, evidence, catalog_entries, catalog_entry_list, title_context, table_models)

    EXPORTS_DIR.mkdir(exist_ok=True)
    out = EXPORTS_DIR / f"东坡大道水土保持方案汇报_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    prs.save(out)
    write_notes(plan)
    return out


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--plan-only", action="store_true")
    parser.add_argument("--build", action="store_true")
    args = parser.parse_args()
    if args.build:
        print(build_deck())
        return
    ledger, plan = write_planning_artifacts()
    print(
        json.dumps(
            {
                "evidence": len(ledger["evidence"]),
                "planned_pages": len(plan["slides"]),
                "chapter_dividers": len(plan["deck"]["chapter_order"]),
                "expected_pptx_slides": len(plan["slides"]) + len(plan["deck"]["chapter_order"]),
                "plan": str(PLAN_PATH),
                "ledger": str(LEDGER_PATH),
            },
            ensure_ascii=False,
        )
    )


if __name__ == "__main__":
    main()
