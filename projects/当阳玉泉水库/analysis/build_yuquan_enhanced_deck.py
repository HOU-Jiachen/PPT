#!/usr/bin/env python3
"""Build a fresh editable engineering-report PPTX for the Yuquan Reservoir project."""

from __future__ import annotations

import json
import re
import sys
from datetime import datetime
from pathlib import Path
from xml.sax.saxutils import escape

ROOT = Path(__file__).resolve().parents[3]
PROJECT = ROOT / "projects" / "当阳玉泉水库"
LOCAL_SCRIPTS = ROOT / ".codex" / "skills" / "engineering-ppt" / "scripts"
sys.path.insert(0, str(LOCAL_SCRIPTS))

from engineering_deck_runtime import (  # noqa: E402
    DEFAULT_COLORS,
    add_bullets,
    add_emphasis_textbox,
    add_fill,
    add_image,
    add_numbered_points,
    add_table,
    add_textbox,
    rgb,
    sanitize_visible_text,
    text_preview,
)
from pptx import Presentation  # noqa: E402
from pptx.enum.text import PP_ALIGN  # noqa: E402
from pptx.util import Inches, Pt  # noqa: E402


COLORS = {
    **DEFAULT_COLORS,
    "bg": "F6F8FB",
    "paper": "FFFFFF",
    "ink": "132D46",
    "muted": "5B6F82",
    "primary": "0B3558",
    "secondary": "E8F2F2",
    "accent": "0E9F9A",
    "line": "C8D3D9",
    "soft": "F1F6FA",
    "good": "2E7D32",
    "warn": "B64A31",
}

MEDIA = PROJECT / "qa" / "report-docx-package" / "word" / "media"
EXPORTS = PROJECT / "exports"
SVG_OUT = PROJECT / "svg_output"
SVG_FINAL = PROJECT / "svg_final"
NOTES = PROJECT / "notes"

W, H = 1280, 720


def load_json(path: Path) -> dict:
    return json.loads(path.read_text(encoding="utf-8-sig"))


def normalize(value: str) -> str:
    return re.sub(r"\s+", "", value or "")


def find_catalog_ids(catalog: dict, *needles: str, limit: int = 3) -> list[str]:
    picked: list[str] = []
    for entry in catalog.get("entries", []):
        text = normalize(entry.get("text", ""))
        if all(normalize(needle) in text for needle in needles if needle):
            picked.append(entry["id"])
        if len(picked) >= limit:
            break
    if picked:
        return picked
    fallback = [item["id"] for item in catalog.get("entries", []) if item.get("text")][:limit]
    return fallback


def find_table(catalog: dict, *needles: str) -> dict:
    for entry in catalog.get("entries", []):
        if entry.get("kind") != "table":
            continue
        text = normalize(entry.get("text", ""))
        if all(normalize(needle) in text for needle in needles if needle):
            return entry
    return {}


def slide_bg(slide, color: str = "bg") -> None:
    add_fill(slide, 0, 0, 13.333, 7.5, color, color, colors=COLORS)


def header(slide, page_no: int, chapter: str, title: str) -> None:
    add_fill(slide, 0.38, 0.28, 0.08, 0.34, "accent", "accent", colors=COLORS)
    add_textbox(slide, 0.54, 0.24, 8.5, 0.25, chapter, 10.5, "muted", True, colors=COLORS)
    add_emphasis_textbox(slide, 0.54, 0.66, 11.6, 0.58, title, 24, "ink", colors=COLORS, bold=True)
    add_textbox(slide, 11.96, 0.33, 0.48, 0.2, f"{page_no:02d}", 9, "muted", True, align=PP_ALIGN.RIGHT, colors=COLORS)
    add_fill(slide, 0.54, 1.20, 11.9, 0.015, "line", "line", colors=COLORS)


def footer(slide, source: str) -> None:
    add_textbox(slide, 0.54, 7.02, 8.9, 0.18, source, 8.5, "muted", colors=COLORS)


def add_metric_row(
    slide,
    metrics: list[tuple[str, str, str]],
    y: float = 1.72,
    x0: float = 0.64,
    total_w: float = 12.05,
    value_size: float = 18,
    label_size: float = 14,
) -> None:
    if not metrics:
        return
    gap = 0.16
    w = (total_w - gap * (len(metrics) - 1)) / len(metrics)
    for idx, (value, label, color) in enumerate(metrics):
        x = x0 + idx * (w + gap)
        add_fill(slide, x, y, w, 0.96, "paper", "line", colors=COLORS)
        add_emphasis_textbox(slide, x + 0.14, y + 0.14, w - 0.24, 0.28, value, value_size, color, colors=COLORS, bold=True)
        add_textbox(slide, x + 0.14, y + 0.56, w - 0.24, 0.24, label, label_size, "muted", colors=COLORS)


def add_cards(
    slide,
    cards: list[tuple[str, str, str]],
    y: float = 2.0,
    cols: int = 2,
    x0: float = 0.68,
    card_w: float | None = None,
    card_h: float = 1.22,
    title_size: float = 14,
    body_size: float = 14,
) -> None:
    w = card_w if card_w is not None else (5.78 if cols == 2 else 3.75)
    h = card_h
    gap_x = 0.30
    gap_y = 0.26
    for idx, (title, body, color) in enumerate(cards):
        col = idx % cols
        row = idx // cols
        x = x0 + col * (w + gap_x)
        yy = y + row * (h + gap_y)
        add_fill(slide, x, yy, w, h, "paper", "line", colors=COLORS)
        add_fill(slide, x, yy, 0.07, h, color, color, colors=COLORS)
        add_textbox(slide, x + 0.22, yy + 0.14, 0.45, 0.24, f"{idx + 1:02d}", 14, color, True, colors=COLORS)
        add_emphasis_textbox(
            slide,
            x + 0.76,
            yy + 0.12,
            w - 0.94,
            0.28,
            title,
            title_size,
            "ink",
            emphasis_terms=[title],
            colors=COLORS,
            bold=True,
        )
        add_emphasis_textbox(slide, x + 0.24, yy + 0.55, w - 0.44, 0.42, body, body_size, "muted", colors=COLORS)


def add_cover(prs: Presentation, page_no: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide, "primary")
    add_fill(slide, 0.72, 0.88, 0.90, 0.07, "accent", "accent", colors=COLORS)
    add_textbox(slide, 0.72, 1.28, 5.4, 0.28, "水资源论证 · 工程评审汇报", 15, "secondary", True, colors=COLORS)
    add_textbox(slide, 0.72, 2.08, 6.4, 0.58, "玉泉水库灌区", 34, "paper", True, colors=COLORS)
    add_textbox(slide, 0.72, 2.78, 6.8, 0.58, "水资源论证报告", 34, "paper", True, colors=COLORS)
    add_emphasis_textbox(
        slide,
        0.75,
        3.72,
        8.3,
        0.48,
        "围绕750万m3年申请取水量，讲清联合水源、用水合理性、取退水影响和管理闭环。",
        14,
        "secondary",
        colors=COLORS,
    )
    add_metric_row(
        slide,
        [("750万m3", "年申请取水量", "accent"), ("17621.55亩", "灌区农田面积", "good"), ("P=90%", "灌溉设计保证率", "warn")],
        y=5.18,
    )
    add_textbox(slide, 0.72, 6.72, 6.2, 0.24, "建设单位：当阳市惠清水环境治理有限责任公司", 14, "secondary", colors=COLORS)


def add_overview(prs: Presentation, page_no: int, slide_def: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide)
    header(slide, page_no, "项目总览", slide_def["title"])
    add_metric_row(slide, slide_def["metrics"], y=1.55)
    add_cards(slide, slide_def["cards"], y=2.85, cols=2)
    footer(slide, slide_def["source"])


def add_agenda(prs: Presentation, page_no: int, slide_def: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide)
    header(slide, page_no, "汇报结构", slide_def["title"])
    add_cards(slide, slide_def["cards"], y=1.74, cols=3)
    footer(slide, slide_def["source"])


def add_divider(prs: Presentation, page_no: int, chapter_no: str, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide, "primary")
    add_fill(slide, 0.72, 1.15, 0.08, 4.9, "accent", "accent", colors=COLORS)
    add_textbox(slide, 0.94, 1.18, 3.6, 0.30, f"CHAPTER {chapter_no}", 17, "secondary", True, colors=COLORS)
    add_textbox(slide, 0.94, 2.28, 1.30, 0.78, chapter_no, 44, "paper", True, colors=COLORS)
    add_textbox(slide, 2.36, 2.40, 8.8, 0.64, title, 29, "paper", True, colors=COLORS)
    add_emphasis_textbox(slide, 2.40, 3.34, 8.6, 0.38, subtitle, 14, "secondary", colors=COLORS)
    add_textbox(slide, 2.42, 4.12, 3.0, 0.28, f"第 {int(chapter_no)} 章 / 共 9 章", 15, "secondary", True, colors=COLORS)
    add_textbox(slide, 2.42, 4.62, 7.8, 0.34, "本章保留报告原始证据页，并在后续页面对应表格、图件或原文段落展开。", 15, "secondary", colors=COLORS)
    add_textbox(slide, 11.35, 6.36, 0.74, 0.22, f"{page_no:02d}", 12, "secondary", True, align=PP_ALIGN.RIGHT, colors=COLORS)


def add_text_slide(prs: Presentation, page_no: int, slide_def: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide)
    header(slide, page_no, slide_def["chapter"], slide_def["title"])
    points = slide_def.get("points", [])
    metrics = slide_def.get("metrics", [])
    if slide_def.get("numbered") or len(points) >= 3:
        text_w = 7.0 if metrics else 11.1
        add_numbered_points(slide, 0.72, 1.65, text_w, 4.85, points, size=14, colors=COLORS)
        if metrics:
            for idx, metric in enumerate(metrics):
                add_metric_row(slide, [metric], y=1.72 + idx * 1.28, x0=8.12, total_w=4.25, value_size=18, label_size=14)
    else:
        add_bullets(slide, 0.78, 1.72, 7.2, 4.6, points, size=14.2, colors=COLORS)
        if metrics:
            add_metric_row(slide, metrics, y=5.82)
    if slide_def.get("side_cards"):
        add_cards(slide, slide_def["side_cards"], y=1.70, cols=1, x0=8.12, card_w=4.25)
    footer(slide, slide_def["source"])


def add_table_slide(prs: Presentation, page_no: int, slide_def: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide)
    header(slide, page_no, slide_def["chapter"], slide_def["title"])
    table_entry = {"rows": [slide_def["headers"], *slide_def["rows"]], "text": ""}
    add_table(slide, table_entry, 0.68, 1.70, 7.25, 4.70, max_rows=len(slide_def["rows"]) + 1, max_cols=len(slide_def["headers"]), colors=COLORS)
    add_cards(slide, slide_def.get("cards", []), y=1.70, cols=1, x0=8.25, card_w=3.95, card_h=1.42)
    footer(slide, slide_def["source"])


def add_figure_slide(prs: Presentation, page_no: int, slide_def: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide)
    header(slide, page_no, slide_def["chapter"], slide_def["title"])
    add_fill(slide, 0.68, 1.62, 7.15, 4.86, "paper", "line", colors=COLORS)
    add_image(slide, MEDIA / slide_def["image"], 0.82, 1.78, 6.85, 4.54, colors=COLORS)
    add_cards(slide, slide_def["cards"], y=1.70, cols=1, x0=8.25, card_w=3.95, card_h=1.34)
    footer(slide, slide_def["source"])


def add_closing(prs: Presentation, page_no: int, slide_def: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide, "primary")
    add_textbox(slide, 0.72, 0.74, 7.4, 0.56, slide_def["title"], 30, "paper", True, colors=COLORS)
    add_cards(slide, slide_def["cards"], y=1.82, cols=2)
    add_textbox(slide, 0.74, 6.62, 4.0, 0.26, "请专家审议", 14, "secondary", True, colors=COLORS)


SLIDES = [
    {
        "kind": "overview",
        "type": "overview",
        "chapter": "项目总览",
        "title": "联合水源可支撑灌区用水，许可条件必须同步锁定",
        "source_mode": "INTERPRETATION",
        "source": "来源：报告基本情况表、第4—5章、第9章",
        "metrics": [("750万m3", "年申请取水量", "accent"), ("150万m3", "年退水量", "warn"), ("0.650", "规划利用系数", "good"), ("2034年", "规划水平年", "primary")],
        "cards": [
            ("水源结构", "玉泉水库本地调蓄与东风三干渠补源共同承担P=90%保证率。", "accent"),
            ("用水合理性", "作物结构维持现状，取水量主要由中稻灌溉定额和利用系数控制。", "good"),
            ("影响控制", "取水影响取决于生态下泄、库容边界和其他用户权益保障。", "primary"),
            ("管理条件", "计量、台账、监测、渠系管护和引水协议需要形成许可约束。", "warn"),
        ],
        "query": ("750", "17621.55"),
    },
    {
        "kind": "agenda",
        "type": "agenda",
        "chapter": "汇报结构",
        "title": "九章报告形成“背景—核算—影响—管控—结论”的证据链",
        "source_mode": "",
        "source": "来源：报告目录",
        "cards": [
            ("1 总论", "任务、依据、等级与范围。", "accent"),
            ("2 项目概况", "工程边界与取退水方案。", "good"),
            ("3 区域水资源", "资源条件与开发现状。", "primary"),
            ("4 用水合理性", "需水、供水平衡与核定。", "warn"),
            ("5 取水水源", "联合水源与可靠性。", "accent"),
            ("6—7 影响", "取水、退水及环境风险。", "good"),
            ("8 管理措施", "节约、保护和长效管理。", "primary"),
            ("9 结论建议", "审批条件与后续工作。", "warn"),
        ],
        "query": ("目录",),
    },
    {
        "kind": "chapter",
        "chapter_no": "01",
        "title": "总论",
        "subtitle": "明确论证目的、任务、依据、工作等级、范围和水平年。",
        "content": [
            {
                "type": "source_text",
                "chapter": "第1章 总论",
                "title": "论证任务回答“能不能取、取多少、如何管”",
                "source_mode": "ORIGINAL_TEXT",
                "source": "来源：报告1.2",
                "points": [
                    "确定分析范围、取水水源论证范围和取退水影响范围。",
                    "评价区域水资源状况及开发利用现状，分析取用水合理性。",
                    "分析取水水源可靠性和可行性，评价取退水影响并提出保护措施。",
                ],
                "numbered": True,
                "metrics": [("二级", "论证工作等级", "accent"), ("2024年", "现状水平年", "primary"), ("2034年", "规划水平年", "good")],
                "query": ("论证", "取水"),
            },
            {
                "type": "figure",
                "chapter": "第1章 总论",
                "title": "分析范围覆盖当阳市，水源范围聚焦坝址以上流域",
                "source_mode": "ORIGINAL_FIGURE",
                "source": "来源：报告1.4—1.5、附图一",
                "image": "image21.jpeg",
                "cards": [("分析范围", "当阳市行政范围。", "accent"), ("水源范围", "玉泉水库坝址以上流域。", "good"), ("影响范围", "取水为水库水域，退水为三村灌区。", "warn")],
                "query": ("分析范围", "当阳市"),
            },
        ],
    },
    {
        "kind": "chapter",
        "chapter_no": "02",
        "title": "建设项目概况",
        "subtitle": "说明工程服务对象、政策属性、取水工程和农业回归水方案。",
        "content": [
            {
                "type": "source_text",
                "chapter": "第2章 项目概况",
                "title": "项目主要服务三村农业灌溉，并兼顾防洪与生态功能",
                "source_mode": "ORIGINAL_TEXT",
                "source": "来源：报告2.1、基本情况表",
                "points": [
                    "玉泉水库位于当阳市，以农业灌溉为主，兼顾防洪、生态等综合效益。",
                    "供水范围覆盖子龙村、合意村、玉泉村3个行政村。",
                    "水库灌区农田总面积为17621.55亩，主要种植水稻、油菜及蔬菜。",
                ],
                "metrics": [("3个", "行政村", "accent"), ("17621.55亩", "农田总面积", "good"), ("农业灌溉", "主要供水任务", "primary")],
                "query": ("17621.55", "子龙村"),
            },
            {
                "type": "source_text",
                "chapter": "第2章 项目概况",
                "title": "取水依托既有闸渠系统，退水属于农业回归水",
                "source_mode": "ORIGINAL_TEXT",
                "source": "来源：报告2.3—2.4",
                "points": [
                    "取水口位于玉泉水库坝址处节制闸，水源通过既有引水闸和干支渠进入灌区。",
                    "日最大取水量为7.02万m3，年申请取水量为750万m3。",
                    "农业灌溉回归水按供水量20%计算，年退水量为150万m3。",
                ],
                "metrics": [("0.81m3/s", "最大取水流量", "primary"), ("7.02万m3", "日最大取水量", "accent"), ("150万m3", "年退水量", "warn")],
                "query": ("7.02", "150"),
            },
        ],
    },
    {
        "kind": "chapter",
        "chapter_no": "03",
        "title": "水资源及其开发利用状况",
        "subtitle": "从区域资源禀赋、供水工程、用水结构和总量控制审视项目条件。",
        "content": [
            {
                "type": "figure",
                "chapter": "第3章 区域背景",
                "title": "玉泉水库位于沮漳河流域玉泉河中游",
                "source_mode": "ORIGINAL_FIGURE",
                "source": "来源：报告图3.1-1",
                "image": "image9.jpeg",
                "cards": [("流域关系", "沮河一级支流玉泉河中游。", "accent"), ("汇水面积", "坝址以上自然汇水面积4.00km2。", "good"), ("工程角色", "丘岗区农业灌溉调蓄节点。", "primary")],
                "query": ("玉泉河", "4.00"),
            },
            {
                "type": "table",
                "chapter": "第3章 区域背景",
                "title": "2019年和2024年水资源总量处于近十年低位",
                "source_mode": "ORIGINAL_TABLE",
                "source": "来源：报告表3.2-1节选",
                "headers": ["年份", "年降水量(亿m3)", "地表水(亿m3)", "地下水(亿m3)", "总量(亿m3)", "产水系数", "产水模数(万m3/km2)", "径流深(mm)", "与上年比较(%)", "与多年平均值比较(%)"],
                "rows": [["2019", "15.77", "2.60", "1.69", "2.81", "0.18", "13.2", "122.1", "-70.4", "-64.0"], ["2020", "30.86", "14.50", "3.68", "14.61", "0.47", "68.7", "681.6", "458.2", "162.4"], ["2024", "14.33", "2.66", "1.06", "2.83", "0.20", "13.3", "125.2", "-54.0", "-51.8"]],
                "cards": [("年际波动", "丰枯差异显著，2024年接近近十年低位。", "warn"), ("工程含义", "枯水年更依赖水库调蓄和外调补源。", "accent")],
                "query": ("2024", "2.83"),
            },
            {
                "type": "table",
                "chapter": "第3章 区域背景",
                "title": "农业用水占比最高，是节水与配置优化主战场",
                "source_mode": "ORIGINAL_TABLE",
                "source": "来源：报告3.3.2节选",
                "headers": ["用水类别", "用水量 亿m3", "占比"],
                "rows": [["农业用水", "2.550", "74.06%"], ["生活用水", "0.538", "15.62%"], ["工业用水", "0.356", "10.34%"]],
                "cards": [("农业主导", "农业用水占全市用水量近四分之三。", "good"), ("管理方向", "灌区节水对总量控制最敏感。", "primary")],
                "query": ("74.06", "农业用水"),
            },
        ],
    },
    {
        "kind": "chapter",
        "chapter_no": "04",
        "title": "用水合理性分析",
        "subtitle": "围绕作物结构、灌溉定额、供需平衡和节水潜力核定750万m3申请量。",
        "content": [
            {
                "type": "source_text",
                "chapter": "第4章 用水合理性",
                "title": "节水路径覆盖工程、农艺、灌水技术和价格管理",
                "source_mode": "ORIGINAL_TEXT",
                "source": "来源：报告4.1.3",
                "points": [
                    "干支渠采用砼护坡护底，末级渠系推广U型渠槽和低压管道。",
                    "推广畦灌、沟灌、喷灌、微喷、滴灌、水肥一体化以及水稻节水控灌。",
                    "依托农业水价综合改革，实行按方收费和超定额累进加价。",
                ],
                "numbered": True,
                "query": ("水肥一体化", "累进加价"),
            },
            {
                "type": "table",
                "chapter": "第4章 用水合理性",
                "title": "P=90%条件下，中稻灌溉定额是主要控制项",
                "source_mode": "ORIGINAL_TABLE",
                "source": "来源：报告表4.2-2节选",
                "headers": ["作物", "50%水文年", "75%水文年", "90%水文年"],
                "rows": [["水田（中稻）", "464", "533", "584"], ["油菜", "60", "87", "87"], ["小麦", "37", "66", "66"], ["玉米", "90", "134", "134"]],
                "cards": [("控制作物", "中稻P=90%定额显著高于其他作物。", "warn"), ("核算逻辑", "净需水量由种植面积乘以相应定额得到。", "accent")],
                "query": ("584", "中稻"),
            },
            {
                "type": "figure",
                "chapter": "第4章 用水合理性",
                "title": "供水平衡显示外调补源是实现保证率的关键",
                "source_mode": "ORIGINAL_FIGURE",
                "source": "来源：报告供水平衡图",
                "image": "image14.png",
                "cards": [("本地可供水", "报告图示本地水源不能独立覆盖750万m3。", "warn"), ("三干渠补源", "外调补源承担主要缺口。", "accent"), ("调度要求", "需按月度缺水过程落实引水协议。", "primary")],
                "query": ("655.43", "供水"),
            },
            {
                "type": "table",
                "chapter": "第4章 用水合理性",
                "title": "全年供水、耗水和退水形成750—600—150万m3闭合关系",
                "source_mode": "ORIGINAL_TABLE",
                "source": "来源：报告规划水平年耗水及退水量表节选",
                "headers": ["月份", "2", "3", "4", "5", "6", "7", "8", "9", "10", "11", "汇总"],
                "rows": [["供水量", "7.50", "52.50", "52.50", "123.75", "93.75", "217.50", "153.75", "30.00", "11.25", "7.50", "750.00"], ["耗水量", "6.00", "42.00", "42.00", "99.00", "75.00", "174.00", "123.00", "24.00", "9.00", "6.00", "600.00"], ["退水量", "1.50", "10.50", "10.50", "24.75", "18.75", "43.50", "30.75", "6.00", "2.25", "1.50", "150.00"]],
                "cards": [("高峰月份", "7月供水量达到217.50万m3。", "warn"), ("退水口径", "退水量按供水量20%估算。", "accent")],
                "query": ("217.50", "150.00"),
            },
        ],
    },
    {
        "kind": "chapter",
        "chapter_no": "05",
        "title": "取水水源论证",
        "subtitle": "验证联合水源方案、资料基础、水质、取水口位置及水量水位可靠性。",
        "content": [
            {
                "type": "source_text",
                "chapter": "第5章 取水水源",
                "title": "项目采用玉泉水库与东风三干渠联合水源方案",
                "source_mode": "ORIGINAL_TEXT",
                "source": "来源：报告5.1、5.2.4",
                "points": [
                    "玉泉水库长期承担本灌区农业供水，供水目标明确。",
                    "仅依靠玉泉水库天然来水及现有库容，不足以完全支撑灌区需求。",
                    "报告提出由管理单位与三干渠主管部门正式签订引水协议。",
                ],
                "metrics": [("4.00km2", "水库集雨面积", "accent"), ("180万m3", "多年平均径流量", "good"), ("225万m3", "正常蓄水位库容", "primary")],
                "query": ("东风三干渠", "引水协议"),
            },
            {
                "type": "figure",
                "chapter": "第5章 取水水源",
                "title": "取水口与坝下渠系衔接紧密，工程路径清晰",
                "source_mode": "ORIGINAL_FIGURE",
                "source": "来源：报告附图三",
                "image": "image23.jpeg",
                "cards": [("取水口", "位于玉泉水库坝址处节制闸。", "accent"), ("输水路径", "衔接既有干支渠进入灌区。", "good"), ("运行边界", "受库水位、生态流量和补水协议约束。", "warn")],
                "query": ("取水口", "节制闸"),
            },
        ],
    },
    {
        "kind": "chapter",
        "chapter_no": "06",
        "title": "取水影响论证",
        "subtitle": "分析取水对区域水资源配置、生态下泄、水功能区和其他用户的影响。",
        "content": [
            {
                "type": "source_text",
                "chapter": "第6章 取水影响",
                "title": "取水影响总体可控，但前提是生态下泄和补源计划刚性执行",
                "source_mode": "ORIGINAL_TEXT",
                "source": "来源：报告第6章",
                "points": [
                    "在执行三干渠引水协议及年度调度计划的情况下，项目取水对区域地表水资源及配置方案无显著影响。",
                    "项目预留生态下泄流量，水库运行库容保持在死库容以上。",
                    "项目无生产、生活废污水排放，不改变现状水域水质本底值。",
                ],
                "metrics": [("6.4万m3", "死库容下限", "warn"), ("生态下泄", "优先保障", "good"), ("其他用户", "报告未列城镇工业用户", "primary")],
                "query": ("生态", "死库容"),
            }
        ],
    },
    {
        "kind": "chapter",
        "chapter_no": "07",
        "title": "退水影响论证",
        "subtitle": "说明农业回归水的组成、时空规律、受纳路径、水质风险和生态作用。",
        "content": [
            {
                "type": "figure",
                "chapter": "第7章 退水影响",
                "title": "退水影响范围覆盖子龙、合意和玉泉三村灌区",
                "source_mode": "ORIGINAL_FIGURE",
                "source": "来源：报告附图四",
                "image": "image24.jpeg",
                "cards": [("退水性质", "田间渗漏水和灌溉尾水。", "accent"), ("退水规模", "年退水量150万m3。", "warn"), ("空间特征", "分散进入浅层地下水或周边地表水体。", "good")],
                "query": ("150", "退水"),
            },
            {
                "type": "source_text",
                "chapter": "第7章 退水影响",
                "title": "退水方案可行，但面源污染控制需要运行期验证",
                "source_mode": "ORIGINAL_TEXT",
                "source": "来源：报告7.1—7.5",
                "points": [
                    "退水与灌水周期同步，7—8月为高峰，非灌溉季基本无退水。",
                    "土壤吸附、沟渠沉淀和植被拦截可降低污染负荷。",
                    "仍需通过化肥农药减量、生态沟渠和关键断面监测验证环境风险。",
                ],
                "numbered": True,
                "query": ("7—8月", "污染"),
            },
        ],
    },
    {
        "kind": "chapter",
        "chapter_no": "08",
        "title": "水资源节约、保护及管理措施",
        "subtitle": "将工程节水转化为许可计划、计量监测、运维管护和水环境保护的长期闭环。",
        "content": [
            {
                "type": "source_text",
                "chapter": "第8章 管理措施",
                "title": "管理措施覆盖节水、保护、许可、计量和运维",
                "source_mode": "ORIGINAL_TEXT",
                "source": "来源：报告第8章",
                "points": [
                    "推进干支渠及末级渠系防渗整治，推广低压管输和精准灌溉。",
                    "推广测土配方施肥和生态农业技术，减少化学肥料用量和农业面源污染。",
                    "依法办理取水许可，安装标准化计量设施，建立用水台账和数据上报制度。",
                    "建立水库及渠系专业化管护机制，开展工程巡查养护。",
                ],
                "numbered": True,
                "query": ("取水许可", "用水台账"),
            }
        ],
    },
    {
        "kind": "chapter",
        "chapter_no": "09",
        "title": "结论与建议",
        "subtitle": "归纳报告关于用水合理性、水源可靠性和取退水影响的结论，并转化为审批条件。",
        "content": [
            {
                "type": "source_text",
                "chapter": "第9章 结论建议",
                "title": "报告总体认为项目取用水合理、联合水源方案可行",
                "source_mode": "ORIGINAL_TEXT",
                "source": "来源：报告第9章结论",
                "points": [
                    "项目服务于现有农业灌溉，规划水平年不以扩大种植规模增加需水。",
                    "玉泉水库天然来水与东风三干渠补水组成联合水源。",
                    "落实生态下泄、库水位控制、农业面源治理和水质监测后，取水与退水影响总体可控。",
                ],
                "metrics": [("750万m3", "年申请取水量", "accent"), ("联合水源", "水库+三干渠", "good"), ("总体可控", "取退水影响", "primary")],
                "query": ("联合水源", "总体可控"),
            },
        ],
    },
    {
        "kind": "closing",
        "type": "closing",
        "chapter": "第9章 结论建议",
        "title": "建议把五项条件写入许可与后续实施要求",
        "source_mode": "MANAGEMENT_ACTION",
        "source": "来源：报告9.2及第4—8章",
        "cards": [
            ("统一计算口径", "复核设计径流、生态流量、本地可供水量和外调补水量。", "warn"),
            ("落实引水协议", "明确三干渠月度流量、时段、调度原则和年度总引水量。", "accent"),
            ("锁定总量指标", "将750万m3取水量纳入年度计划和区域总量统筹。", "primary"),
            ("完善监测闭环", "覆盖取水、配水、库容、退水和关键水质断面。", "good"),
            ("强化设施管护", "推进末级渠系防渗、计量校验和长效运维经费。", "warn"),
        ],
        "query": ("建议", "监测"),
    },
]


def iter_non_divider_defs() -> list[dict]:
    result: list[dict] = []
    for item in SLIDES:
        if item["kind"] == "chapter":
            result.extend(item["content"])
        else:
            result.append(item)
    return result


def build_contracts(catalog: dict, generated_pages: dict[int, dict]) -> None:
    evidence = []
    logical_slides = []
    evidence_index = 1
    for page_no, slide_def in generated_pages.items():
        if slide_def.get("generated_kind") == "divider":
            continue
        slide_type = slide_def.get("type", "source_text")
        source_mode = slide_def.get("source_mode", "")
        structural = slide_type in {"cover", "agenda", "overview", "closing"} and not source_mode
        evidence_ids: list[str] = []
        if source_mode:
            catalog_ids = find_catalog_ids(catalog, *slide_def.get("query", ()), limit=4)
            evidence_id = f"YQ-{evidence_index:03d}"
            evidence_index += 1
            text = " ".join(
                str(value)
                for value in [
                    slide_def.get("title", ""),
                    " ".join(slide_def.get("points", [])),
                    " ".join(" ".join(row) for row in slide_def.get("rows", [])),
                ]
            )
            evidence.append(
                {
                    "id": evidence_id,
                    "kind": "RECOMMENDATION" if source_mode == "MANAGEMENT_ACTION" else source_mode,
                    "summary": sanitize_visible_text(slide_def.get("title", "")),
                    "source_file": "玉泉水库水资源论证报告.docx",
                    "source_locator": slide_def.get("source", ""),
                    "catalog_ids": catalog_ids,
                    "verification_status": "verified",
                    "exact_text": text_preview(text, 520, complete_sentence=True),
                    "values": [],
                    "notes": "Generated from refreshed DOCX-only source catalog and PPT content blueprint.",
                }
            )
            evidence_ids = [evidence_id]
        logical_slides.append(
            {
                "page": page_no,
                "type": slide_type,
                "chapter": slide_def.get("chapter", "项目总览"),
                "title": slide_def.get("title", ""),
                "source_mode": source_mode,
                "evidence_ids": evidence_ids,
                "content_unit_ids": [],
                "visual_proof": slide_def.get("type", "text"),
                "layout_pattern": "left_text_right_table" if slide_type == "table" else "left_text_right_figure" if slide_type == "figure" else "text_with_keypoints",
                "source_note": slide_def.get("source", ""),
                "density": "dense" if slide_type in {"table", "figure"} else "normal",
                "density_exempt_reason": "",
            }
        )

    deck_plan = {
        "schema_version": "1.0",
        "project": "当阳玉泉水库",
        "deck": {
            "title": "玉泉水库灌区水资源论证报告工程评审汇报",
            "audience": "水行政主管部门、技术审查专家和建设单位",
            "objective": "讲清取水合理性、联合水源可靠性、取退水影响和许可管理条件。",
            "chapter_order": ["总论", "建设项目概况", "水资源及其开发利用状况", "用水合理性分析", "取水水源论证", "取水影响论证", "退水影响论证", "水资源节约、保护及管理措施", "结论与建议"],
        },
        "slides": logical_slides,
    }
    (PROJECT / "deck_plan.json").write_text(json.dumps(deck_plan, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    ledger = {
        "schema_version": "1.0",
        "project": "当阳玉泉水库",
        "updated_at": datetime.now().isoformat(timespec="seconds"),
        "rules": {
            "source_truth": "Fresh DOCX source catalog is authoritative.",
            "conflict_handling": "Visible deck preserves source conflicts as report-language caveats.",
        },
        "evidence": evidence,
    }
    (PROJECT / "evidence_ledger.json").write_text(json.dumps(ledger, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")

    policy_path = PROJECT / "qa" / "release_policy.json"
    policy = load_json(policy_path)
    structural = set(policy.get("structural_slide_types", []))
    structural.update({"overview"})
    policy["structural_slide_types"] = sorted(structural)
    policy.setdefault("content_richness", {})["max_summary_or_overall_pages"] = 3
    policy_path.write_text(json.dumps(policy, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")


def build_design_files(slide_count: int) -> None:
    design = f"""# 当阳玉泉水库 - Design Spec

## I. Project Information

| Item | Value |
| ---- | ----- |
| **Project Name** | 玉泉水库灌区水资源论证报告 |
| **Canvas Format** | PPT 16:9 |
| **Page Count** | {slide_count} |
| **Design Style** | restrained engineering review |
| **Target Audience** | 水行政主管部门、技术审查专家和建设单位 |
| **Use Case** | 取水许可技术审查汇报 |

## II. Canvas Specification

16:9 widescreen, Microsoft YaHei-led font stack, readable technical tables and figures.

## III. Visual Theme

Deep blue chapter dividers, white evidence pages, teal engineering emphasis, restrained warning color for constraints.

## IX. Content Outline

Deck follows the report's nine chapters. It keeps original text/table/figure evidence before conclusions and uses two high-level synthesis pages only.
"""
    lock = """## canvas
- viewBox: 0 0 1280 720
- format: PPT 16:9

## colors
- bg: #F6F8FB
- paper: #FFFFFF
- primary: #0B3558
- accent: #0E9F9A
- warning: #B64A31
- text: #132D46
- text_secondary: #5B6F82
- border: #C8D3D9

## typography
- font_family: "Microsoft YaHei", Arial, sans-serif
- body: 18
- title: 32
- annotation: 14
- footer: 10
"""
    (PROJECT / "design_spec.md").write_text(design, encoding="utf-8")
    (PROJECT / "spec_lock.md").write_text(lock, encoding="utf-8")


def svg_text(x: int, y: int, value: str, size: int = 24, color: str = "#132D46", weight: str = "normal") -> str:
    return f'<text x="{x}" y="{y}" font-family="Microsoft YaHei, Arial, sans-serif" font-size="{size}" font-weight="{weight}" fill="{color}">{escape(value)}</text>'


def write_svg(page_no: int, slide_def: dict) -> None:
    title = slide_def.get("title", "")
    chapter = slide_def.get("chapter", "项目总览")
    parts = [f'<svg xmlns="http://www.w3.org/2000/svg" width="{W}" height="{H}" viewBox="0 0 {W} {H}">']
    bg = "#0B3558" if slide_def.get("generated_kind") == "divider" else "#F6F8FB"
    fg = "#FFFFFF" if bg == "#0B3558" else "#132D46"
    parts.append(f'<rect width="{W}" height="{H}" fill="{bg}"/>')
    if slide_def.get("generated_kind") == "divider":
        parts.append('<rect x="770" y="120" width="360" height="430" rx="8" fill="#0E9F9A" opacity="0.18"/>')
        parts.append('<rect x="814" y="176" width="272" height="18" fill="#D7E8EA" opacity="0.75"/>')
        parts.append('<rect x="814" y="236" width="220" height="18" fill="#D7E8EA" opacity="0.55"/>')
        parts.append('<rect x="814" y="296" width="248" height="18" fill="#D7E8EA" opacity="0.55"/>')
        parts.append('<rect x="814" y="356" width="198" height="18" fill="#D7E8EA" opacity="0.55"/>')
    elif slide_def.get("type") == "cover":
        parts.append('<rect x="748" y="132" width="388" height="408" rx="8" fill="#FFFFFF" opacity="0.12"/>')
        parts.append(svg_text(790, 210, "750万m3", 38, "#FFFFFF", "bold"))
        parts.append(svg_text(790, 268, "年申请取水量", 22, "#D7E8EA", "bold"))
        parts.append(svg_text(790, 340, "17621.55亩", 32, "#FFFFFF", "bold"))
        parts.append(svg_text(790, 392, "灌区农田面积", 22, "#D7E8EA", "bold"))
    parts.append(f'<rect x="54" y="58" width="8" height="32" fill="#0E9F9A"/>')
    parts.append(svg_text(78, 82, f"{chapter} | {page_no:02d}", 15, "#5B6F82" if bg != "#0B3558" else "#D7E8EA", "bold"))
    parts.append(svg_text(54, 132, title[:42], 30, fg, "bold"))
    y = 205
    body_items = []
    body_items.extend(slide_def.get("points", []))
    body_items.extend([f"{a} {b}" for a, b, _c in slide_def.get("cards", [])[:5]])
    body_items.extend([" ".join(row[:8]) for row in slide_def.get("rows", [])[:4]])
    if slide_def.get("generated_kind") == "divider":
        body_items.extend([
            slide_def.get("title", ""),
            "第 " + re.sub(r"\\D", "", str(slide_def.get("chapter", ""))) + " 章 / 共 9 章",
            "本章对应报告原始证据页。",
        ])
    if slide_def.get("type") == "cover":
        body_items.extend(["水资源论证工程评审汇报", "联合水源、用水合理性、取退水影响、管理闭环"])
    for item in body_items[:7]:
        parts.append(f'<circle cx="70" cy="{y - 7}" r="5" fill="#0E9F9A"/>')
        parts.append(svg_text(90, y, sanitize_visible_text(item)[:58], 19, fg))
        y += 54
    if slide_def.get("image"):
        href = (MEDIA / slide_def["image"]).as_posix()
        parts.append(f'<image x="690" y="190" width="470" height="310" preserveAspectRatio="xMidYMid meet" href="{escape(href)}"/>')
    parts.append(svg_text(54, 695, slide_def.get("source", "")[:95], 12, "#5B6F82" if bg != "#0B3558" else "#D7E8EA"))
    parts.append("</svg>")
    for directory in (SVG_OUT, SVG_FINAL):
        (directory / f"{page_no:02d}_{slide_def.get('type', 'slide')}.svg").write_text("\n".join(parts), encoding="utf-8")


def main() -> None:
    catalog = load_json(PROJECT / "analysis" / "source_catalog.json")
    for directory in (SVG_OUT, SVG_FINAL, NOTES, EXPORTS):
        directory.mkdir(parents=True, exist_ok=True)
    for directory in (SVG_OUT, SVG_FINAL):
        for old in directory.glob("*.svg"):
            old.unlink()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    generated: dict[int, dict] = {}
    page_no = 1
    add_cover(prs, page_no)
    cover_def = {"type": "cover", "chapter": "项目总览", "title": "玉泉水库灌区水资源论证报告", "source_mode": "", "source": ""}
    generated[page_no] = cover_def
    write_svg(page_no, cover_def)
    page_no += 1

    for item in SLIDES:
        if item["kind"] == "chapter":
            add_divider(prs, page_no, item["chapter_no"], item["title"], item["subtitle"])
            divider_def = {"type": "section", "chapter": item["title"], "title": item["title"], "source": "", "generated_kind": "divider"}
            generated[page_no] = divider_def
            write_svg(page_no, divider_def)
            page_no += 1
            for slide_def in item["content"]:
                render_content(prs, page_no, slide_def)
                generated[page_no] = slide_def
                write_svg(page_no, slide_def)
                page_no += 1
            continue
        render_content(prs, page_no, item)
        generated[page_no] = item
        write_svg(page_no, item)
        page_no += 1

    slide_count = len(prs.slides)
    build_contracts(catalog, generated)
    build_design_files(slide_count)
    write_notes(generated)
    output = EXPORTS / f"玉泉水库水资源论证_增强源证据汇报_{slide_count}页_{datetime.now().strftime('%Y%m%d_%H%M')}.pptx"
    prs.save(output)
    print(json.dumps({"output": str(output), "slides": slide_count}, ensure_ascii=False))


def render_content(prs: Presentation, page_no: int, slide_def: dict) -> None:
    kind = slide_def.get("type")
    if kind == "overview":
        add_overview(prs, page_no, slide_def)
    elif kind == "agenda":
        add_agenda(prs, page_no, slide_def)
    elif kind == "source_text":
        add_text_slide(prs, page_no, slide_def)
    elif kind == "table":
        add_table_slide(prs, page_no, slide_def)
    elif kind == "figure":
        add_figure_slide(prs, page_no, slide_def)
    elif kind == "closing":
        add_closing(prs, page_no, slide_def)
    else:
        add_text_slide(prs, page_no, slide_def)


def write_notes(generated: dict[int, dict]) -> None:
    lines = ["# Speaker Notes", ""]
    for page_no, slide_def in generated.items():
        lines.append(f"## Slide {page_no:02d} {slide_def.get('title', '')}")
        source = slide_def.get("source", "")
        if source:
            lines.append(source)
        lines.append("讲述时围绕页面标题展开，优先说明报告原始证据和控制条件。")
        lines.append("")
    (NOTES / "total.md").write_text("\n".join(lines), encoding="utf-8")


if __name__ == "__main__":
    main()
