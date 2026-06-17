from __future__ import annotations

import json
import sys
from datetime import datetime
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Inches

PROJECT_DIR = Path(__file__).resolve().parents[1]
REPO_ROOT = PROJECT_DIR.parents[1]
LOCAL_SKILL_SCRIPTS = REPO_ROOT / ".codex" / "skills" / "engineering-ppt" / "scripts"
sys.path.insert(0, str(LOCAL_SKILL_SCRIPTS))

from engineering_deck_runtime import (  # noqa: E402
    DEFAULT_COLORS,
    add_agenda_slide,
    add_bullets,
    add_fill,
    add_image_panel,
    add_table,
    add_textbox,
    add_title,
    catalog_title_context,
    collect_template_profile,
    create_engineering_design_spec,
    create_engineering_spec_lock,
    evidence_lookup,
    figure_image_index,
    image_paths_from_text,
    load_json,
    report_points_from_page,
    report_text_from_page,
    rgb,
    sanitize_visible_text,
    slide_evidence,
    source_topic_heading,
    table_from_catalog,
    text_preview,
)


CATALOG_PATH = PROJECT_DIR / "analysis" / "source_catalog.json"
PLAN_PATH = PROJECT_DIR / "deck_plan.json"
LEDGER_PATH = PROJECT_DIR / "evidence_ledger.json"
MEDIA_DIR = PROJECT_DIR / "sources" / "大旺塘水文地质勘察报告_files"
MARKDOWN_PATH = PROJECT_DIR / "sources" / "大旺塘水文地质勘察报告.md"
EXPORTS_DIR = PROJECT_DIR / "exports"
ANALYSIS_DIR = PROJECT_DIR / "analysis"
NOTES_DIR = PROJECT_DIR / "notes"

COLORS = DEFAULT_COLORS.copy()
ANCHOR_PAGES = {1, 2, 7, 13, 19, 31, 41, 53, 63, 69, 74}
IMPORTANT_IMAGE_NAMES = [
    "image_003.jpeg",
    "image_034.png",
    "image_060.png",
    "image_143.jpeg",
    "image_194.jpeg",
    "image_195.jpeg",
    "image_207.png",
    "image_216.png",
]

PROJECT_TITLE = "大旺塘矿山地下水治理专项水文地质勘察工程评审汇报"
AUDIENCE = "业主单位、矿山安全与水文地质评审专家、设计及治理实施相关方"
USE_CASE = "专项勘察工程评审汇报"
DESIGN_STYLE = "学习完整版评审PPT的深蓝工程评审风格，保持源证据优先和专家评审可读性"


def write_backend_specs(plan: dict[str, Any], template_profile: dict[str, Any]) -> None:
    design_spec = create_engineering_design_spec(
        plan=plan,
        media_dir=MEDIA_DIR,
        template_profile=template_profile,
        project_title=PROJECT_TITLE,
        audience=AUDIENCE,
        use_case=USE_CASE,
        design_style=DESIGN_STYLE,
        image_names=IMPORTANT_IMAGE_NAMES,
    )
    spec_lock = create_engineering_spec_lock(
        plan=plan,
        media_dir=MEDIA_DIR,
        important_anchor_pages=ANCHOR_PAGES,
    )
    (PROJECT_DIR / "design_spec.md").write_text(design_spec, encoding="utf-8")
    (PROJECT_DIR / "spec_lock.md").write_text(spec_lock, encoding="utf-8")


def write_analysis_artifacts(template_profile: dict[str, Any], figure_index: dict[str, list[str]]) -> None:
    ANALYSIS_DIR.mkdir(exist_ok=True)
    (ANALYSIS_DIR / "template_profile.json").write_text(
        json.dumps(template_profile, ensure_ascii=False, indent=2) + "\n",
        encoding="utf-8",
    )
    (ANALYSIS_DIR / "figure_image_index.json").write_text(
        json.dumps(figure_index, ensure_ascii=False, indent=2) + "\n",
        encoding="utf-8",
    )


def add_cover_slide(slide) -> None:
    add_fill(slide, 0, 0, 13.333, 7.5, "primary", "primary", colors=COLORS)
    add_textbox(slide, 0.78, 0.72, 4.7, 0.28, "华润水泥（封开）有限公司", 14, "paper", colors=COLORS)
    add_textbox(slide, 0.80, 1.38, 11.55, 0.78, "大旺塘矿山地下水治理专项", 34, "paper", True, colors=COLORS)
    add_textbox(slide, 0.80, 2.28, 11.45, 0.70, "水文地质勘察工程评审汇报", 28, "paper", True, colors=COLORS)
    add_textbox(slide, 0.84, 3.92, 10.7, 0.36, "专项勘察 | 工程编号 K062601 | 2026年4月报告", 15, "paper", colors=COLORS)
    add_textbox(slide, 0.84, 4.42, 10.8, 0.36, "依据报告章节、原始图表、计算口径和防治水建议形成评审级汇报", 14, "paper", colors=COLORS)
    add_textbox(slide, 0.86, 6.72, 4.2, 0.3, datetime.now().strftime("%Y年%m月"), 14, "paper", colors=COLORS)


def figure_paths_for_page(page: dict[str, Any], figure_index: dict[str, list[str]]) -> list[Path]:
    source_text = " ".join(
        [
            page.get("visual_proof", ""),
            page.get("source_note", ""),
            page.get("title", ""),
        ]
    )
    return image_paths_from_text(MEDIA_DIR, source_text, figure_index)


def add_calculation_slide(
    slide,
    page: dict[str, Any],
    evidence: list[dict[str, Any]],
    body_text: str,
    table_entry: dict[str, Any] | None,
    topic_heading: str,
) -> None:
    add_fill(slide, 0.55, 1.18, 5.5, 4.95, "paper", colors=COLORS)
    formulas = [sanitize_visible_text(ev.get("formula", "")) for ev in evidence if ev.get("formula")]
    inputs: list[str] = []
    for ev in evidence:
        inputs.extend(ev.get("inputs", []) or [])

    add_textbox(slide, 0.85, 1.45, 4.9, 0.3, topic_heading, 16, "accent", True, colors=COLORS)
    add_textbox(
        slide,
        0.85,
        1.9,
        4.9,
        1.25,
        text_preview("；".join(formulas) or body_text or "报告按公式、参数表和结果表列示计算过程。", 240),
        14,
        "ink",
        colors=COLORS,
    )
    add_bullets(
        slide,
        0.85,
        3.35,
        4.9,
        2.0,
        [text_preview(item, 130) for item in inputs[:4]] or [text_preview(body_text, 130)],
        14,
        colors=COLORS,
    )
    if table_entry:
        add_table(slide, table_entry, 6.35, 1.18, 5.9, 4.95, max_rows=7, max_cols=4, colors=COLORS)
    else:
        add_textbox(slide, 6.45, 1.3, 5.5, 3.5, text_preview(body_text, 420), 14, "ink", colors=COLORS)


def add_source_slide(
    slide,
    page: dict[str, Any],
    evidence: list[dict[str, Any]],
    catalog_entries: dict[str, dict[str, Any]],
    catalog_entry_list: list[dict[str, Any]],
    title_context: dict[str, dict[str, str]],
    figure_index: dict[str, list[str]],
) -> None:
    add_title(slide, page, colors=COLORS)
    body_text = report_text_from_page(page, evidence, catalog_entries, catalog_entry_list)
    topic_heading = source_topic_heading(page, evidence, title_context, catalog_entry_list)
    image_paths = figure_paths_for_page(page, figure_index)
    table_entry = table_from_catalog(catalog_entries, evidence)

    if page["source_mode"] == "ORIGINAL_FIGURE" or image_paths:
        add_image_panel(slide, image_paths, 0.55, 1.10, 7.5, 5.55, colors=COLORS)
        add_textbox(slide, 8.35, 1.18, 4.25, 0.32, topic_heading, 15, "accent", True, colors=COLORS)
        add_textbox(slide, 8.35, 1.58, 4.15, 3.75, text_preview(body_text or page["title"], 420), 14, "ink", colors=COLORS)
    elif page["source_mode"] == "ORIGINAL_TABLE" and table_entry:
        add_table(slide, table_entry, 0.55, 1.18, 8.35, 4.95, max_rows=7, max_cols=4, colors=COLORS)
        add_textbox(slide, 9.25, 1.25, 3.2, 0.32, topic_heading, 15, "accent", True, colors=COLORS)
        add_textbox(slide, 9.25, 1.65, 3.1, 3.45, text_preview(body_text or table_entry.get("text", ""), 360), 14, "ink", colors=COLORS)
    elif page["source_mode"] == "CALCULATION":
        add_calculation_slide(slide, page, evidence, body_text, table_entry, topic_heading)
    else:
        add_fill(slide, 0.65, 1.25, 12.0, 4.9, "paper", colors=COLORS)
        add_textbox(slide, 0.95, 1.55, 11.35, 0.35, topic_heading, 16, "accent", True, colors=COLORS)
        points = report_points_from_page(page, evidence, catalog_entries, catalog_entry_list, max_points=5)
        if len(points) >= 2:
            add_bullets(slide, 0.95, 2.08, 10.95, 3.30, points, 14, colors=COLORS)
        else:
            add_textbox(slide, 0.95, 2.08, 11.1, 3.15, text_preview(body_text or page["title"], 760), 15, "ink", colors=COLORS)


def write_notes(plan: dict[str, Any]) -> None:
    NOTES_DIR.mkdir(exist_ok=True)
    total: list[str] = []
    for page in plan["slides"]:
        note = (
            f"Slide {page['page']:02d} - {sanitize_visible_text(page['title'])}\n\n"
            f"讲解要点：围绕“{sanitize_visible_text(page['title'])}”展开，"
            f"按报告对应章节、图表或计算口径说明；资料来源：{sanitize_visible_text(page.get('source_note', ''))}。\n"
        )
        (NOTES_DIR / f"{page['page']:02d}.md").write_text(note, encoding="utf-8")
        total.append(note)
    (NOTES_DIR / "total.md").write_text("\n\n".join(total), encoding="utf-8")


def build_deck() -> Path:
    plan = load_json(PLAN_PATH)
    ledger = load_json(LEDGER_PATH)
    catalog = load_json(CATALOG_PATH)
    catalog_entry_list = catalog["entries"]
    catalog_entries = {entry["id"]: entry for entry in catalog_entry_list}
    title_context = catalog_title_context(catalog_entry_list)
    evidence_by_id = evidence_lookup(ledger)

    template_profile = collect_template_profile(PROJECT_DIR)
    figure_index = figure_image_index(MARKDOWN_PATH)
    write_analysis_artifacts(template_profile, figure_index)
    write_backend_specs(plan, template_profile)

    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for page in plan["slides"]:
        slide = prs.slides.add_slide(blank)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = rgb(COLORS["bg"])

        if page["page"] == 1:
            add_cover_slide(slide)
        elif page["type"] == "agenda":
            add_agenda_slide(slide, plan, page, colors=COLORS)
        else:
            evidence = slide_evidence(page, evidence_by_id)
            add_source_slide(slide, page, evidence, catalog_entries, catalog_entry_list, title_context, figure_index)

    EXPORTS_DIR.mkdir(exist_ok=True)
    out = EXPORTS_DIR / f"大旺塘工程评审汇报_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    prs.save(out)
    write_notes(plan)
    return out


if __name__ == "__main__":
    print(build_deck())
