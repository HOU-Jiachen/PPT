#!/usr/bin/env python3
"""Deterministic PPT table renderers backed by analysis/table_ir.json."""

from __future__ import annotations

import json
from pathlib import Path
from typing import Any

from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from text_fitting import SlideContentSanitizer, TextFittingEngine, TextStyle, component_char_limit


TEXT_SANITIZER = SlideContentSanitizer()
TEXT_FITTER = TextFittingEngine()
TABLE_MIN_PT = 12.0
BODY_MIN_PT = 14.0


def _rgb(hex_value: str):
    from pptx.dml.color import RGBColor

    value = (hex_value or "#FFFFFF").strip().lstrip("#")
    if len(value) != 6:
        value = "FFFFFF"
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def _fit_text(text: str, w: float, h: float, size: float, min_size: float, component: str) -> tuple[str, float, bool]:
    fitted = TEXT_FITTER.fit_text_to_box(
        text,
        w,
        h,
        TextStyle(
            font_size=size,
            min_font_size=min_size,
            max_chars=component_char_limit(component),
            component=component,
        ),
    )
    return fitted.text, fitted.font_size, fitted.overflow


def _fit_rect(x: float, y: float, w: float, h: float, aspect: float) -> tuple[float, float, float, float]:
    if aspect <= 0:
        return x, y, w, h
    box_aspect = w / max(0.01, h)
    if box_aspect > aspect:
        fitted_w = h * aspect
        return x + (w - fitted_w) / 2, y, fitted_w, h
    fitted_h = w / aspect
    return x, y + (h - fitted_h) / 2, w, fitted_h


def _alignment(value: str):
    return {
        "left": PP_ALIGN.LEFT,
        "right": PP_ALIGN.RIGHT,
        "center": PP_ALIGN.CENTER,
        "justify": PP_ALIGN.JUSTIFY,
    }.get(value or "center", PP_ALIGN.CENTER)


def _valign(value: str):
    return {
        "top": MSO_ANCHOR.TOP,
        "bottom": MSO_ANCHOR.BOTTOM,
        "middle": MSO_ANCHOR.MIDDLE,
        "center": MSO_ANCHOR.MIDDLE,
    }.get(value or "middle", MSO_ANCHOR.MIDDLE)


def load_table_ir(project: Path) -> dict[str, Any]:
    path = project / "analysis" / "table_ir.json"
    return json.loads(path.read_text(encoding="utf-8-sig")) if path.exists() else {"tables": [], "by_table_id": {}}


def get_table_ir(project: Path, table_id: str) -> dict[str, Any]:
    payload = load_table_ir(project)
    if table_id in payload.get("by_table_id", {}):
        return payload["by_table_id"][table_id]
    for table in payload.get("tables", []):
        if table.get("table_id") == table_id or table.get("locator") == table_id:
            return table
    raise KeyError(f"Unknown table_id: {table_id}")


class NativeTableRenderer:
    """Render simple Table IR records as editable PowerPoint tables."""

    def __init__(self, colors: dict[str, str] | None = None, font_name: str = "Microsoft YaHei"):
        self.colors = colors or {"ink": "132D46", "secondary": "EAF3F3", "line": "CCD8E2"}
        self.font_name = font_name

    def render(self, slide, table_ir: dict[str, Any], x: float, y: float, w: float, h: float):
        structure = table_ir.get("structure", {})
        rows = int(structure.get("rows") or len(table_ir.get("rows", [])) or 1)
        cols = int(structure.get("cols") or max((len(row) for row in table_ir.get("rows", [])), default=1))
        shape = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h))
        table = shape.table
        self._apply_geometry(table, table_ir, w, h, rows, cols)
        self._apply_cells(table, table_ir, rows, cols)
        self._apply_merges(table, table_ir, rows, cols)
        return shape

    def _apply_geometry(self, ppt_table, table_ir: dict[str, Any], w: float, h: float, rows: int, cols: int) -> None:
        col_widths = [float(value or 0) for value in table_ir.get("geometry", {}).get("col_widths", [])]
        if len(col_widths) != cols or not any(col_widths):
            col_widths = [1.0] * cols
        total_width = sum(col_widths) or 1.0
        for index, column in enumerate(ppt_table.columns):
            column.width = Inches(w * col_widths[index] / total_width)

        row_heights = [float(value or 0) for value in table_ir.get("geometry", {}).get("row_heights", [])]
        if len(row_heights) != rows or not any(row_heights):
            row_heights = [1.0] * rows
        total_height = sum(row_heights) or 1.0
        for index, row in enumerate(ppt_table.rows):
            row.height = Inches(h * row_heights[index] / total_height)

    def _apply_cells(self, ppt_table, table_ir: dict[str, Any], rows: int, cols: int) -> None:
        cell_lookup = {(int(cell["row"]), int(cell["col"])): cell for cell in table_ir.get("cells", [])}
        raw_rows = table_ir.get("rows", [])
        for r_idx in range(rows):
            for c_idx in range(cols):
                ppt_cell = ppt_table.cell(r_idx, c_idx)
                ir_cell = cell_lookup.get((r_idx, c_idx), {})
                text = ir_cell.get("text")
                if text is None and r_idx < len(raw_rows) and c_idx < len(raw_rows[r_idx]):
                    text = raw_rows[r_idx][c_idx]
                style = ir_cell.get("style", {})
                fitted_text, fitted_size, cell_overflow = _fit_text(
                    str(text or ""),
                    1.0,
                    0.35,
                    max(float(style.get("font_size", TABLE_MIN_PT)), TABLE_MIN_PT),
                    TABLE_MIN_PT,
                    "table_cell",
                )
                if cell_overflow:
                    ppt_cell.text = fitted_text
                else:
                    ppt_cell.text = fitted_text
                ppt_cell.vertical_anchor = _valign(style.get("valign", "middle"))
                margin = style.get("margin", {})
                ppt_cell.margin_left = Inches(float(margin.get("left", 0.04)))
                ppt_cell.margin_right = Inches(float(margin.get("right", 0.04)))
                ppt_cell.margin_top = Inches(float(margin.get("top", 0.02)))
                ppt_cell.margin_bottom = Inches(float(margin.get("bottom", 0.02)))
                ppt_cell.fill.solid()
                ppt_cell.fill.fore_color.rgb = _rgb(style.get("fill", "#FFFFFF"))
                for paragraph in ppt_cell.text_frame.paragraphs:
                    paragraph.alignment = _alignment(style.get("align", "center"))
                    paragraph.font.name = self.font_name
                    paragraph.font.size = Pt(fitted_size)
                    paragraph.font.bold = bool(style.get("bold") or ir_cell.get("is_header"))
                    paragraph.font.color.rgb = _rgb(style.get("color", "#" + self.colors.get("ink", "132D46")))

    def _apply_merges(self, ppt_table, table_ir: dict[str, Any], rows: int, cols: int) -> None:
        for merged in table_ir.get("structure", {}).get("merged_cells", []):
            r0 = int(merged.get("row", 0))
            c0 = int(merged.get("col", 0))
            r1 = r0 + int(merged.get("rowspan", 1)) - 1
            c1 = c0 + int(merged.get("colspan", 1)) - 1
            if 0 <= r0 <= r1 < rows and 0 <= c0 <= c1 < cols and (r0 != r1 or c0 != c1):
                try:
                    ppt_table.cell(r0, c0).merge(ppt_table.cell(r1, c1))
                except ValueError:
                    pass


class ImageTableRenderer:
    """Render source table crop assets while preserving aspect ratio."""

    def __init__(self, project: Path, colors: dict[str, str] | None = None, font_name: str = "Microsoft YaHei"):
        self.project = project
        self.colors = colors or {"ink": "132D46", "muted": "5B6F82", "line": "CCD8E2"}
        self.font_name = font_name

    def render(
        self,
        slide,
        table_ir: dict[str, Any],
        x: float,
        y: float,
        w: float,
        h: float,
        *,
        title: str | None = None,
        note: str | None = None,
    ):
        title_text = TEXT_SANITIZER.sanitize(title if title is not None else table_ir.get("title", ""), component="caption")
        if title_text:
            title_text, title_size, title_overflow = _fit_text(title_text, w, 0.28, 14, 12, "caption")
            box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.28))
            if title_overflow:
                box.name = "TableTitleNeedsRelayout"
            p = box.text_frame.paragraphs[0]
            p.text = title_text
            p.font.name = self.font_name
            p.font.size = Pt(title_size)
            p.font.bold = True
            p.font.color.rgb = _rgb("#" + self.colors.get("ink", "132D46"))
            y += 0.34
            h -= 0.34

        image_path = self._image_path(table_ir)
        if not image_path.exists():
            raise FileNotFoundError(f"Missing table image asset: {image_path}")
        aspect = self._image_aspect(image_path)
        ix, iy, iw, ih = _fit_rect(x, y, w, h - (0.28 if note else 0), aspect)
        picture = slide.shapes.add_picture(str(image_path), Inches(ix), Inches(iy), width=Inches(iw), height=Inches(ih))
        if note:
            note_text, note_size, note_overflow = _fit_text(note, w, 0.22, 10, 8, "table_note")
            box = slide.shapes.add_textbox(Inches(x), Inches(y + h - 0.24), Inches(w), Inches(0.22))
            if note_overflow:
                box.name = "TableNoteNeedsRelayout"
            p = box.text_frame.paragraphs[0]
            p.text = note_text
            p.font.name = self.font_name
            p.font.size = Pt(note_size)
            p.font.color.rgb = _rgb("#" + self.colors.get("muted", "5B6F82"))
        return picture

    def _image_path(self, table_ir: dict[str, Any]) -> Path:
        asset = table_ir.get("assets", {}).get("crop_image") or table_ir.get("assets", {}).get("source_table_image")
        if not asset:
            raise FileNotFoundError(f"Table {table_ir.get('table_id')} has no crop_image asset")
        return (self.project / asset).resolve()

    @staticmethod
    def _image_aspect(path: Path) -> float:
        from PIL import Image

        with Image.open(path) as image:
            return image.width / max(1, image.height)


class HybridTableRenderer:
    """Render source crop plus short, source-grounded conclusions."""

    def __init__(self, project: Path, colors: dict[str, str] | None = None, font_name: str = "Microsoft YaHei"):
        self.project = project
        self.colors = colors or {"ink": "132D46", "muted": "5B6F82", "accent": "0E9F9A", "line": "CCD8E2"}
        self.font_name = font_name
        self.image_renderer = ImageTableRenderer(project, self.colors, font_name)

    def render(
        self,
        slide,
        table_ir: dict[str, Any],
        x: float,
        y: float,
        w: float,
        h: float,
        conclusions: list[str] | None = None,
        layout: str = "right",
    ):
        conclusions = [
            TEXT_SANITIZER.sanitize(str(item), component="hybrid_conclusion")
            for item in (conclusions or [])
            if str(item).strip()
        ][:4]
        if layout == "bottom":
            image_h = h * 0.68
            self.image_renderer.render(slide, table_ir, x, y, w, image_h)
            self._conclusions(slide, x, y + image_h + 0.18, w, h - image_h - 0.18, conclusions)
        else:
            image_w = w * 0.66
            self.image_renderer.render(slide, table_ir, x, y, image_w, h)
            self._conclusions(slide, x + image_w + 0.24, y + 0.18, w - image_w - 0.24, h - 0.36, conclusions)

    def _conclusions(self, slide, x: float, y: float, w: float, h: float, conclusions: list[str]) -> None:
        title = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.28))
        p = title.text_frame.paragraphs[0]
        p.text = "核心结论"
        p.font.name = self.font_name
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = _rgb("#" + self.colors.get("accent", "0E9F9A"))
        row_h = (h - 0.44) / max(1, len(conclusions))
        for index, text in enumerate(conclusions, start=1):
            box = slide.shapes.add_textbox(Inches(x), Inches(y + 0.42 + (index - 1) * row_h), Inches(w), Inches(max(0.32, row_h - 0.04)))
            p = box.text_frame.paragraphs[0]
            fitted_text, fitted_size, needs_relayout = _fit_text(
                f"{index}. {text}",
                w,
                max(0.32, row_h - 0.04),
                14,
                BODY_MIN_PT,
                "hybrid_conclusion",
            )
            if needs_relayout:
                box.name = "HybridConclusionNeedsRelayout"
            p.text = fitted_text
            p.font.name = self.font_name
            p.font.size = Pt(fitted_size)
            p.font.color.rgb = _rgb("#" + self.colors.get("ink", "132D46"))


def render_table(
    slide,
    table_ir: dict[str, Any],
    project: Path,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    conclusions: list[str] | None = None,
    mode: str | None = None,
    colors: dict[str, str] | None = None,
):
    selected = (mode or table_ir.get("render_mode") or "image").lower()
    if selected == "native":
        return NativeTableRenderer(colors).render(slide, table_ir, x, y, w, h)
    if selected == "hybrid":
        return HybridTableRenderer(project, colors).render(slide, table_ir, x, y, w, h, conclusions)
    return ImageTableRenderer(project, colors).render(slide, table_ir, x, y, w, h)
