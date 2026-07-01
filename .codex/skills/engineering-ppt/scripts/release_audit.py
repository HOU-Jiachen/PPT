#!/usr/bin/env python3
"""Audit engineering PPT planning artifacts, SVG output, and the exported PPTX."""

from __future__ import annotations

import argparse
import difflib
import json
import math
import re
import sys
import zipfile
from collections import defaultdict
from pathlib import Path
from typing import Any
from xml.etree import ElementTree as ET

from lxml import etree
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from text_fitting import (
    INTERNAL_FORBIDDEN_PHRASES,
    SlideContentSanitizer,
    normalize_text,
    split_sentences,
)


NUMBER_RE = re.compile(r"(?<![\w.])[-+]?\d+(?:\.\d+)?%?(?![\w.])")
SVG_NS = {"svg": "http://www.w3.org/2000/svg"}
STRUCTURAL_DEFAULT = {"cover", "agenda", "section", "closing"}
ORIGINAL_MODES_DEFAULT = {"ORIGINAL_TEXT", "ORIGINAL_TABLE", "ORIGINAL_FIGURE", "CALCULATION"}
DEFAULT_FORBIDDEN_VISIBLE_PHRASES = [
    "报告原文摘录",
    "工程原文摘录",
    "缺失项沿用提取状态",
    "不做补造",
    "不作补造",
    "审查提示",
    "数据提示",
    "口径A",
    "口径B",
    "来源模式",
    "证据编号",
    "本页用于",
    "证据解读",
    "原表定位",
    "行列规模",
    "密集表按重点行重排",
    "完整数据回看报告原表",
    "source_mode",
    "evidence_ids",
    "visual_proof",
    "layout_pattern",
    "source_note",
    "Source mode",
    "报告对图件的说明",
    "报告对表格的说明",
    "报告计算口径",
    "报告阐述",
    "资料来源：",
    "原始对象",
    "必讲内容",
    "保留理由",
    "按照报告章节顺序进入",
    "关键数值保持源表口径",
]
DEFAULT_FORBIDDEN_VISIBLE_REGEXES = [
    re.compile(r"\bE-\d(?:-[A-Z0-9]+)+\b"),
    re.compile(r"\bT-?\d{3}\b"),
    re.compile(r"\bimage_\d+\.(?:png|jpe?g|wmf|emf)\b", re.IGNORECASE),
    re.compile(r"\b(?:ORIGINAL_TEXT|ORIGINAL_TABLE|ORIGINAL_FIGURE|CALCULATION|INTERPRETATION|CONCLUSION|MANAGEMENT_ACTION)\b"),
    re.compile(r"(?:\.{3,}|…|……)"),
]


def display_path(path: Path, base: Path | None = None) -> str:
    resolved = path.resolve()
    roots = [base.resolve()] if base else []
    roots.append(Path.cwd().resolve())
    for root in roots:
        try:
            return resolved.relative_to(root).as_posix()
        except ValueError:
            continue
    return resolved.name


class Audit:
    def __init__(self) -> None:
        self.items: list[dict] = []

    def add(self, level: str, code: str, message: str, **context: object) -> None:
        self.items.append(
            {"level": level, "code": code, "message": message, "context": context}
        )

    def error(self, code: str, message: str, **context: object) -> None:
        self.add("error", code, message, **context)

    def warning(self, code: str, message: str, **context: object) -> None:
        self.add("warning", code, message, **context)

    def info(self, code: str, message: str, **context: object) -> None:
        self.add("info", code, message, **context)

    def summary(self) -> dict:
        counts = defaultdict(int)
        for item in self.items:
            counts[item["level"]] += 1
        return {
            "errors": counts["error"],
            "warnings": counts["warning"],
            "info": counts["info"],
            "release_ready": counts["error"] == 0,
        }


def load_json(path: Path, audit: Audit, required: bool = True) -> dict:
    if not path.exists():
        if required:
            audit.error("missing-file", f"Required file is missing: {path.name}", path=str(path))
        return {}
    try:
        return json.loads(path.read_text(encoding="utf-8-sig"))
    except Exception as exc:
        audit.error("invalid-json", f"Invalid JSON: {path.name}", error=str(exc))
        return {}


def normalize_number(token: str) -> str:
    suffix = "%" if token.endswith("%") else ""
    raw = token[:-1] if suffix else token
    try:
        value = float(raw)
    except ValueError:
        return token
    if math.isclose(value, round(value), abs_tol=1e-10):
        return f"{int(round(value))}{suffix}"
    return f"{value:.10f}".rstrip("0").rstrip(".") + suffix


def collect_numbers(text: str) -> set[str]:
    return {normalize_number(token) for token in NUMBER_RE.findall(text)}


def visible_svg_text(path: Path) -> tuple[str, list[dict], int]:
    tree = ET.parse(path)
    root = tree.getroot()
    texts: list[str] = []
    text_nodes: list[dict] = []
    for node in root.findall(".//svg:text", SVG_NS):
        value = "".join(node.itertext()).strip()
        if not value:
            continue
        texts.append(value)
        text_nodes.append(
            {
                "text": value,
                "x": node.attrib.get("x"),
                "y": node.attrib.get("y"),
                "font_size": node.attrib.get("font-size"),
                "anchor": node.attrib.get("text-anchor", "start"),
            }
        )
    visual_count = len(root.findall(".//svg:image", SVG_NS))
    visual_count += len(root.findall(".//svg:path", SVG_NS))
    visual_count += len(root.findall(".//svg:polyline", SVG_NS))
    return "\n".join(texts), text_nodes, visual_count


def estimated_text_overflow(node: dict, width: float, height: float) -> bool:
    try:
        x = float(node["x"])
        y = float(node["y"])
        size = float(node["font_size"])
    except (TypeError, ValueError):
        return False
    text = node["text"]
    estimated = sum(1.0 if ord(ch) > 127 else 0.56 for ch in text) * size
    anchor = node["anchor"]
    left = x if anchor == "start" else x - estimated if anchor == "end" else x - estimated / 2
    right = left + estimated
    return left < -2 or right > width + 2 or y - size < -2 or y > height + 2


def text_bbox(node: dict) -> tuple[float, float, float, float] | None:
    try:
        x = float(node["x"])
        y = float(node["y"])
        size = float(node["font_size"])
    except (TypeError, ValueError):
        return None
    text = str(node.get("text", ""))
    if not text:
        return None
    width = sum(1.0 if ord(ch) > 127 else 0.56 for ch in text) * size
    anchor = node.get("anchor", "start")
    left = x if anchor == "start" else x - width if anchor == "end" else x - width / 2
    top = y - size * 0.90
    return (left, top, left + width, y + size * 0.30)


def parse_svg_length(value: object, default: float = 0) -> float:
    if value is None:
        return default
    match = re.search(r"[-+]?\d+(?:\.\d+)?", str(value))
    if not match:
        return default
    return float(match.group(0))


def image_boxes(root: ET.Element, width: float, height: float, policy: dict) -> list[dict]:
    collision = policy.get("layout_collision", {})
    min_area_ratio = float(collision.get("image_min_area_ratio", 0.03))
    min_area = width * height * min_area_ratio
    boxes: list[dict] = []
    for index, node in enumerate(root.findall(".//svg:image", SVG_NS), start=1):
        x = parse_svg_length(node.attrib.get("x"))
        y = parse_svg_length(node.attrib.get("y"))
        w = parse_svg_length(node.attrib.get("width"))
        h = parse_svg_length(node.attrib.get("height"))
        if w <= 0 or h <= 0 or w * h < min_area:
            continue
        boxes.append({"id": f"image-{index}", "bbox": (x, y, x + w, y + h)})
    return boxes


def rect_boxes(root: ET.Element, width: float, height: float, policy: dict) -> list[dict]:
    aesthetic = policy.get("aesthetic_layout", {})
    min_area_ratio = float(aesthetic.get("minimum_box_area_ratio", 0.004))
    min_area = width * height * min_area_ratio
    max_area = width * height * 0.85
    boxes: list[dict] = []
    for index, node in enumerate(root.findall(".//svg:rect", SVG_NS), start=1):
        x = parse_svg_length(node.attrib.get("x"))
        y = parse_svg_length(node.attrib.get("y"))
        w = parse_svg_length(node.attrib.get("width"))
        h = parse_svg_length(node.attrib.get("height"))
        area = w * h
        if w <= 0 or h <= 0 or area < min_area or area > max_area:
            continue
        boxes.append({"id": f"rect-{index}", "bbox": (x, y, x + w, y + h)})
    return boxes


def intersection_area(
    a: tuple[float, float, float, float], b: tuple[float, float, float, float]
) -> float:
    left = max(a[0], b[0])
    top = max(a[1], b[1])
    right = min(a[2], b[2])
    bottom = min(a[3], b[3])
    if right <= left or bottom <= top:
        return 0
    return (right - left) * (bottom - top)


def box_area(box: tuple[float, float, float, float]) -> float:
    return max(0, box[2] - box[0]) * max(0, box[3] - box[1])


def compact_text(value: str, limit: int = 54) -> str:
    value = re.sub(r"\s+", " ", value).strip()
    return value if len(value) <= limit else value[:limit].rstrip()


def sampled_union_area(
    boxes: list[tuple[float, float, float, float]],
    width: float,
    height: float,
    columns: int = 48,
    rows: int = 27,
) -> float:
    if not boxes:
        return 0
    covered = 0
    for row in range(rows):
        y = height * (row + 0.5) / rows
        for column in range(columns):
            x = width * (column + 0.5) / columns
            if any(box[0] <= x <= box[2] and box[1] <= y <= box[3] for box in boxes):
                covered += 1
    return covered / (columns * rows) * width * height


def audit_aesthetic_layout(
    page_number: int,
    root: ET.Element,
    nodes: list[dict],
    image_items: list[dict],
    width: float,
    height: float,
    slide_type: str,
    policy: dict,
    strict: bool,
    audit: Audit,
) -> None:
    aesthetic = policy.get("aesthetic_layout", {})
    if not strict or aesthetic.get("enabled", True) is False:
        return

    for index, node in enumerate(root.findall(".//svg:image", SVG_NS), start=1):
        preserve = str(node.attrib.get("preserveAspectRatio", "")).strip().lower()
        if preserve == "none":
            audit.error(
                "image-aspect-ratio-disabled",
                "SVG image disables aspect-ratio preservation.",
                page=page_number,
                image=f"image-{index}",
            )

    visual_items = list(image_items)
    collision_items = list(image_items)
    visual_items.extend(rect_boxes(root, width, height, policy))
    for index, node in enumerate(nodes, start=1):
        if is_template_or_meta_text(node, width, height, policy):
            continue
        bbox = text_bbox(node)
        if bbox:
            item = {"id": f"text-{index}", "bbox": bbox}
            visual_items.append(item)
            collision_items.append(item)

    boxes = [item["bbox"] for item in visual_items if box_area(item["bbox"]) > 0]
    if not boxes:
        return

    occupied_area = sampled_union_area(boxes, width, height)
    occupied_ratio = occupied_area / max(width * height, 1)
    minimum = float(
        aesthetic.get(
            "minimum_structural_occupied_area_ratio"
            if slide_type in STRUCTURAL_DEFAULT
            else "minimum_occupied_area_ratio",
            0.18,
        )
    )
    maximum = float(aesthetic.get("maximum_occupied_area_ratio", 0.92))
    if occupied_ratio < minimum:
        audit.warning(
            "low-occupied-area",
            "Slide appears visually sparse under deterministic occupied-area sampling.",
            page=page_number,
            occupied_ratio=round(occupied_ratio, 3),
            minimum=minimum,
        )
    if occupied_ratio > maximum:
        audit.warning(
            "overcrowded-occupied-area",
            "Slide appears visually overcrowded under deterministic occupied-area sampling.",
            page=page_number,
            occupied_ratio=round(occupied_ratio, 3),
            maximum=maximum,
        )

    weighted_area = sum(box_area(box) for box in boxes)
    if weighted_area:
        center_x = sum(((box[0] + box[2]) / 2) * box_area(box) for box in boxes) / weighted_area
        center_y = sum(((box[1] + box[3]) / 2) * box_area(box) for box in boxes) / weighted_area
        offset = max(abs(center_x - width / 2) / width, abs(center_y - height / 2) / height)
        warn_offset = float(aesthetic.get("maximum_balance_offset_ratio", 0.22))
        severe_offset = float(aesthetic.get("severe_balance_offset_ratio", 0.32))
        if offset >= severe_offset:
            audit.error(
                "severe-visual-imbalance",
                "Slide visual mass is severely biased toward one side.",
                page=page_number,
                balance_offset=round(offset, 3),
                maximum=severe_offset,
            )
        elif offset >= warn_offset:
            audit.warning(
                "visual-imbalance",
                "Slide visual mass is biased toward one side.",
                page=page_number,
                balance_offset=round(offset, 3),
                maximum=warn_offset,
            )

    overlap_min = float(aesthetic.get("element_overlap_min_ratio", 0.55))
    severe_overlap = float(aesthetic.get("severe_element_overlap_ratio", 0.75))
    for left_index, left_item in enumerate(collision_items):
        for right_item in collision_items[left_index + 1 :]:
            left_box = left_item["bbox"]
            right_box = right_item["bbox"]
            smaller = min(box_area(left_box), box_area(right_box))
            if not smaller:
                continue
            ratio = intersection_area(left_box, right_box) / smaller
            if ratio >= severe_overlap:
                audit.error(
                    "severe-element-overlap",
                    "Two SVG visual elements severely overlap.",
                    page=page_number,
                    element_a=left_item["id"],
                    element_b=right_item["id"],
                    overlap_ratio=round(ratio, 3),
                )
            elif ratio >= overlap_min:
                audit.warning(
                    "element-overlap",
                    "Two SVG visual elements substantially overlap.",
                    page=page_number,
                    element_a=left_item["id"],
                    element_b=right_item["id"],
                    overlap_ratio=round(ratio, 3),
                )


def parse_font_size(value: object) -> float | None:
    if value is None:
        return None
    match = re.search(r"[-+]?\d+(?:\.\d+)?", str(value))
    if not match:
        return None
    return float(match.group(0))


def is_template_or_meta_text(node: dict, width: float, height: float, policy: dict) -> bool:
    exemptions = policy.get("font_exemptions", {})
    text = str(node.get("text", "")).strip()
    try:
        x = float(node.get("x") or 0)
        y = float(node.get("y") or 0)
    except (TypeError, ValueError):
        return False

    footer_y_min = float(exemptions.get("footer_y_min_ratio", 0.90)) * height
    header_y_max = float(exemptions.get("header_y_max_ratio", 0.12)) * height
    if y >= footer_y_min or y <= header_y_max:
        return True

    for phrase in exemptions.get("phrases", ["来源：", "Source:"]):
        if phrase and text.startswith(str(phrase)):
            return True

    page_number_regex = exemptions.get("page_number_regex")
    if page_number_regex and re.fullmatch(str(page_number_regex), text):
        return x <= width * 0.08 or x >= width * 0.92 or y >= footer_y_min

    return False


def audit_layout_collisions(
    page_number: int,
    nodes: list[dict],
    image_items: list[dict],
    width: float,
    height: float,
    policy: dict,
    strict: bool,
    audit: Audit,
) -> None:
    if not strict or not policy.get("layout_collision", {}).get("enabled", True):
        return
    collision = policy.get("layout_collision", {})
    margin = float(collision.get("canvas_margin_px", 2))
    text_overlap_min = float(collision.get("text_overlap_min_ratio", 0.35))
    text_image_overlap_min = float(collision.get("text_image_overlap_min_ratio", 0.15))

    text_items: list[dict] = []
    for index, node in enumerate(nodes, start=1):
        if is_template_or_meta_text(node, width, height, policy):
            continue
        bbox = text_bbox(node)
        if bbox is None:
            continue
        text_items.append({"index": index, "node": node, "bbox": bbox})
        if bbox[0] < margin or bbox[1] < margin or bbox[2] > width - margin or bbox[3] > height - margin:
            audit.error(
                "text-outside-safe-frame",
                "Visible text probably crosses the slide safe frame.",
                page=page_number,
                text=compact_text(str(node.get("text", ""))),
                bbox=[round(value, 1) for value in bbox],
                frame=[margin, margin, width - margin, height - margin],
            )

    for left_index, left_item in enumerate(text_items):
        for right_item in text_items[left_index + 1 :]:
            left_box = left_item["bbox"]
            right_box = right_item["bbox"]
            overlap = intersection_area(left_box, right_box)
            if not overlap:
                continue
            denominator = min(box_area(left_box), box_area(right_box))
            if denominator and overlap / denominator >= text_overlap_min:
                audit.error(
                    "probable-text-overlap",
                    "Two visible text boxes probably overlap or cover each other.",
                    page=page_number,
                    text_a=compact_text(str(left_item["node"].get("text", ""))),
                    text_b=compact_text(str(right_item["node"].get("text", ""))),
                    overlap_ratio=round(overlap / denominator, 3),
                )

    for text_item in text_items:
        text_box = text_item["bbox"]
        text_area = box_area(text_box)
        if not text_area:
            continue
        for image_item in image_items:
            overlap = intersection_area(text_box, image_item["bbox"])
            if overlap / text_area >= text_image_overlap_min:
                audit.error(
                    "text-overlaps-image",
                    "Visible text probably overlaps a source image or figure.",
                    page=page_number,
                    text=compact_text(str(text_item["node"].get("text", ""))),
                    image=image_item["id"],
                    overlap_ratio=round(overlap / text_area, 3),
                )


def audit_plan(
    project: Path,
    plan: dict,
    ledger: dict,
    policy: dict,
    strict: bool,
    audit: Audit,
) -> dict[int, dict]:
    slides = plan.get("slides", [])
    evidence = {item.get("id"): item for item in ledger.get("evidence", []) if item.get("id")}
    structural = set(policy.get("structural_slide_types", STRUCTURAL_DEFAULT))
    original_modes = set(policy.get("original_source_modes", ORIGINAL_MODES_DEFAULT))
    chapter_exemptions = set(
        policy.get("chapter_original_evidence_exemptions", ["项目总览"])
    )
    pages: dict[int, dict] = {}
    chapter_modes: dict[str, set[str]] = defaultdict(set)

    if strict and not slides:
        audit.error("empty-deck-plan", "deck_plan.json has no slides.")
        return pages

    for slide in slides:
        page = slide.get("page")
        if not isinstance(page, int):
            audit.error("invalid-page", "Every slide must have an integer page number.", slide=slide)
            continue
        pages[page] = slide
        slide_type = str(slide.get("type", "")).lower()
        title = str(slide.get("title", "")).strip()
        chapter = str(slide.get("chapter", "")).strip()
        source_mode = str(slide.get("source_mode", "")).upper()
        evidence_ids = slide.get("evidence_ids", [])

        if not title:
            audit.error("missing-title", "Slide title is empty.", page=page)
        if slide_type not in structural:
            if not chapter:
                audit.error("missing-chapter", "Content slide has no chapter.", page=page)
            if not source_mode:
                audit.error("missing-source-mode", "Content slide has no source_mode.", page=page)
            if not evidence_ids:
                audit.error("missing-evidence", "Content slide has no evidence_ids.", page=page)
        for evidence_id in evidence_ids:
            item = evidence.get(evidence_id)
            if item is None:
                audit.error(
                    "unknown-evidence",
                    "Slide references an unknown evidence ID.",
                    page=page,
                    evidence_id=evidence_id,
                )
            elif item.get("verification_status") != "verified":
                level = audit.error if strict else audit.warning
                level(
                    "unverified-evidence",
                    "Evidence must be verified against the original source.",
                    page=page,
                    evidence_id=evidence_id,
                )
            else:
                evidence_text = " ".join(
                    str(item.get(key, ""))
                    for key in ("summary", "exact_text", "values", "notes")
                )
                slide.setdefault("_allowed_numbers", [])
                slide["_allowed_numbers"].extend(sorted(collect_numbers(evidence_text)))
        if chapter and source_mode:
            chapter_modes[chapter].add(source_mode)

    for chapter, modes in chapter_modes.items():
        if chapter not in chapter_exemptions and not modes.intersection(original_modes):
            audit.error(
                "chapter-without-original-evidence",
                "Chapter has no original text, table, figure, or calculation slide.",
                chapter=chapter,
                modes=sorted(modes),
            )

    richness = policy.get("content_richness", {})
    substantive = [
        slide for slide in slides if str(slide.get("type", "")).lower() not in structural
    ]
    if substantive:
        summary_types = {
            str(item).lower()
            for item in richness.get(
                "summary_page_types",
                ["overview", "summary", "chapter_summary", "conclusion", "closing", "management_action"],
            )
        }
        max_summary_pages = int(richness.get("max_summary_or_overall_pages", 0) or 0)
        summary_pages = [
            slide.get("page")
            for slide in slides
            if str(slide.get("type", "")).lower() in summary_types
            or str(slide.get("source_mode", "")).upper() in {"CONCLUSION", "MANAGEMENT_ACTION"}
        ]
        if strict and max_summary_pages and len(summary_pages) > max_summary_pages:
            audit.error(
                "too-many-summary-pages",
                "Summary, overall, conclusion, and management-action pages exceed the configured maximum.",
                pages=summary_pages,
                maximum=max_summary_pages,
            )

        original_count = sum(
            1
            for slide in substantive
            if str(slide.get("source_mode", "")).upper() in original_modes
        )
        ratio = original_count / len(substantive)
        minimum_ratio = float(richness.get("minimum_original_source_ratio", 0))
        if strict and minimum_ratio and ratio < minimum_ratio:
            audit.error(
                "insufficient-original-source-coverage",
                "Substantive slides do not preserve enough original text, tables, figures, or calculations.",
                original_slides=original_count,
                substantive_slides=len(substantive),
                ratio=round(ratio, 3),
                minimum_ratio=minimum_ratio,
            )

        max_consecutive = int(richness.get("max_consecutive_interpretation_slides", 0))
        if strict and max_consecutive:
            run: list[int] = []
            for slide in slides:
                slide_type = str(slide.get("type", "")).lower()
                source_mode = str(slide.get("source_mode", "")).upper()
                if slide_type in structural or source_mode in original_modes:
                    run = []
                    continue
                if source_mode:
                    run.append(slide.get("page"))
                    if len(run) > max_consecutive:
                        audit.error(
                            "too-many-consecutive-summary-slides",
                            "Too many consecutive non-original-evidence slides.",
                            pages=run,
                            max_consecutive=max_consecutive,
                        )
                        run = run[-max_consecutive:]
    return pages


def visible_forbidden_phrases(policy: dict) -> list[str]:
    phrases = list(DEFAULT_FORBIDDEN_VISIBLE_PHRASES)
    phrases.extend(INTERNAL_FORBIDDEN_PHRASES)
    phrases.extend(str(item) for item in policy.get("forbidden_visible_phrases", []))
    return list(dict.fromkeys(item for item in phrases if item))


def scan_forbidden(text: str, phrases: list[str], audit: Audit, artifact: str) -> None:
    sanitizer = SlideContentSanitizer()
    internal_matches = set(sanitizer.visible_text_violations(text))
    for match in internal_matches:
        audit.error(
            "forbidden-internal-slide-text",
            "Backend analysis, prompt constraints, or agent process wording is visible in the deck.",
            phrase=match,
            artifact=artifact,
        )
    for phrase in phrases:
        if phrase in internal_matches:
            continue
        if phrase and phrase in text:
            audit.error(
                "forbidden-visible-phrase",
                "Internal workflow wording is visible in the deck.",
                phrase=phrase,
                artifact=artifact,
            )
    for pattern in DEFAULT_FORBIDDEN_VISIBLE_REGEXES:
        match = pattern.search(text)
        if match:
            audit.error(
                "forbidden-visible-pattern",
                "Internal identifiers or asset filenames are visible in the deck.",
                pattern=pattern.pattern,
                match=match.group(0),
                artifact=artifact,
            )


def audit_final_text_review_pptx(
    presentation: Presentation,
    policy: dict,
    strict: bool,
    audit: Audit,
) -> None:
    settings = policy.get("final_text_review", {})
    if settings.get("enabled", True) is False:
        return
    sanitizer = SlideContentSanitizer(visible_forbidden_phrases(policy))
    max_bullets = int(settings.get("max_body_bullets", 5))
    max_sentence_chars = int(settings.get("max_sentence_chars", 96))
    error_on_internal = settings.get("error_on_internal_terms", True) is not False
    slide_width = int(presentation.slide_width)
    slide_height = int(presentation.slide_height)
    for slide_index, slide in enumerate(presentation.slides, start=1):
        all_text, body_text, visible_chars, _visual_area = slide_text_payload(slide, presentation, policy)
        if not all_text.strip():
            continue
        violations = sanitizer.visible_text_violations(all_text)
        if violations and error_on_internal:
            audit.error(
                "final-text-review-internal-text",
                "Final PPTX body contains internal notes, generation constraints, or process wording.",
                page=slide_index,
                matches=violations[:8],
                text=compact_text(all_text),
            )

        bullet_like_count = 0
        long_sentences: list[str] = []
        for shape in iter_pptx_shapes(slide.shapes):
            if not getattr(shape, "has_text_frame", False):
                continue
            text = re.sub(r"\s+", " ", getattr(shape, "text", "") or "").strip()
            if not text or is_pptx_layout_meta_text(text, shape, slide_width, slide_height, policy):
                continue
            paragraphs = [
                normalize_text(paragraph.text)
                for paragraph in shape.text_frame.paragraphs
                if normalize_text(paragraph.text)
            ]
            bullet_like_count += sum(1 for item in paragraphs if len(item) >= 8)
            for paragraph in paragraphs:
                for sentence in split_sentences(paragraph):
                    value = normalize_text(sentence)
                    if len(value) > max_sentence_chars:
                        long_sentences.append(value)
        if strict and bullet_like_count > max_bullets + 2 and visible_chars > 260:
            audit.warning(
                "final-text-review-too-many-visible-points",
                "A slide has more visible body points than PPT-style reporting should carry.",
                page=slide_index,
                body_points=bullet_like_count,
                preferred_max=max_bullets,
            )
        if strict and long_sentences:
            audit.warning(
                "final-text-review-report-style-sentence",
                "A slide still contains long report-style sentences; compress or move detail to notes.",
                page=slide_index,
                examples=[compact_text(item, 120) for item in long_sentences[:3]],
            )


def audit_ledger(ledger: dict, catalog: dict, strict: bool, audit: Audit) -> None:
    catalog_ids = {item.get("id") for item in catalog.get("entries", []) if item.get("id")}
    seen = set()
    allowed_kinds = {
        "FACT",
        "ORIGINAL_TEXT",
        "ORIGINAL_TABLE",
        "ORIGINAL_FIGURE",
        "CALCULATION",
        "INTERPRETATION",
        "RECOMMENDATION",
    }
    for item in ledger.get("evidence", []):
        evidence_id = item.get("id")
        if not evidence_id:
            audit.error("evidence-without-id", "Evidence record has no ID.")
            continue
        if evidence_id in seen:
            audit.error("duplicate-evidence-id", "Evidence ID is duplicated.", evidence_id=evidence_id)
        seen.add(evidence_id)
        kind = str(item.get("kind", "")).upper()
        if kind not in allowed_kinds:
            audit.error("invalid-evidence-kind", "Evidence kind is invalid.", evidence_id=evidence_id)
        if item.get("verification_status") == "verified":
            required = ("source_file", "source_locator", "catalog_ids")
            for field in required:
                if not item.get(field):
                    audit.error(
                        "incomplete-verified-evidence",
                        "Verified evidence lacks source traceability.",
                        evidence_id=evidence_id,
                        field=field,
                    )
            for catalog_id in item.get("catalog_ids", []):
                if catalog_id not in catalog_ids:
                    audit.error(
                        "unknown-catalog-id",
                        "Verified evidence references an unknown source-catalog ID.",
                        evidence_id=evidence_id,
                        catalog_id=catalog_id,
                    )
            if kind == "CALCULATION":
                if not item.get("formula") or not item.get("inputs"):
                    audit.error(
                        "incomplete-calculation",
                        "Verified calculation must include formula and verified inputs.",
                        evidence_id=evidence_id,
                    )
        elif strict:
            # Slide-level checks identify use of unverified evidence. This catches unused
            # records too, preventing a misleading all-clear ledger.
            audit.warning(
                "ledger-record-not-verified",
                "Evidence ledger contains a record that is not verified.",
                evidence_id=evidence_id,
            )


def audit_content_analysis(project: Path, strict: bool, audit: Audit) -> dict:
    inventory_path = project / "analysis" / "report_content_inventory.json"
    blueprint_path = project / "analysis" / "ppt_content_blueprint.md"
    inventory = load_json(inventory_path, audit, required=strict)

    if strict and not blueprint_path.exists():
        audit.error(
            "missing-file",
            "Required file is missing: ppt_content_blueprint.md",
            path=str(blueprint_path),
        )
    elif blueprint_path.exists():
        text = blueprint_path.read_text(encoding="utf-8", errors="replace")
        if strict and len(text.strip()) < 800:
            audit.error(
                "content-blueprint-too-thin",
                "PPT content blueprint is too thin to support slide planning.",
                characters=len(text.strip()),
            )

    if not inventory:
        return {}

    summary = inventory.get("summary", {})
    if strict and int(summary.get("sections", 0) or 0) == 0:
        audit.error("empty-content-inventory", "Report content inventory has no sections.")
    if strict and int(summary.get("content_units", 0) or 0) == 0:
        audit.error("empty-content-units", "Report content inventory has no PPT content units.")

    for unit in inventory.get("content_units", []):
        if not unit.get("id"):
            audit.error("content-unit-without-id", "PPT content unit has no ID.")
        if not unit.get("catalog_ids"):
            audit.error(
                "content-unit-without-source",
                "PPT content unit has no source catalog IDs.",
                unit=unit.get("id", ""),
            )
        if not unit.get("layout_hint") and strict:
            audit.error(
                "content-unit-without-layout",
                "PPT content unit has no layout recommendation.",
                unit=unit.get("id", ""),
            )
        structure = unit.get("paragraph_structure", {})
        if (
            strict
            and str(unit.get("source_mode", "")).upper() == "ORIGINAL_TEXT"
            and structure.get("requires_numbered_segments")
            and not structure.get("suggested_segments")
        ):
            audit.error(
                "content-unit-missing-paragraph-segments",
                "Multi-part source paragraph is missing numbered segment/title guidance.",
                unit=unit.get("id", ""),
            )
        if strict and not unit.get("emphasis_candidates"):
            audit.warning(
                "content-unit-without-emphasis-candidates",
                "PPT content unit has no key term/value emphasis candidates.",
                unit=unit.get("id", ""),
            )
    return inventory


def audit_svgs(
    project: Path,
    pages: dict[int, dict],
    policy: dict,
    source_numbers: set[str],
    strict: bool,
    audit: Audit,
) -> int:
    root_candidates = [project / "svg_final", project / "svg_output"]
    candidates = list(root_candidates)
    if not any(directory.exists() and any(directory.glob("*.svg")) for directory in root_candidates):
        candidates.extend(project.glob("*/svg_final"))
        candidates.extend(project.glob("*/svg_output"))
    candidate_files = [
        (directory, sorted(directory.glob("*.svg")))
        for directory in candidates
        if directory.exists()
    ]
    svg_dir, files = max(candidate_files, key=lambda item: len(item[1]), default=(project, []))
    if not files:
        audit.warning("no-svg", "No SVG output found; visual-content checks were skipped.")
        return 0
    audit.info("svg-directory", "Auditing SVG directory.", path=display_path(svg_dir), slides=len(files))

    phrases = visible_forbidden_phrases(policy)
    sparse = policy.get("sparse_page", {})
    min_chars = int(sparse.get("minimum_visible_characters_without_visual", 70))
    for path in files:
        match = re.match(r"(\d+)", path.stem)
        page_number = int(match.group(1)) if match else 0
        slide = pages.get(page_number, {})
        slide_type = str(slide.get("type", path.stem.split("_", 1)[-1])).lower()
        text, nodes, visual_count = visible_svg_text(path)
        scan_forbidden(text, phrases, audit, path.name)

        tree = ET.parse(path)
        root = tree.getroot()
        viewbox = [float(value) for value in root.attrib.get("viewBox", "0 0 1280 720").split()]
        width, height = viewbox[2], viewbox[3]
        image_items = image_boxes(root, width, height, policy)
        for node in nodes:
            if estimated_text_overflow(node, width, height):
                audit.error(
                    "probable-text-overflow",
                    "A single SVG text line probably crosses the canvas boundary.",
                    page=page_number,
                    text=node["text"],
                )
            min_font = (
                policy.get("minimum_font_px", {}).get("body_absolute")
                or policy.get("minimum_font_px", {}).get("body_minimum")
                or 14
            )
            size = parse_font_size(node.get("font_size"))
            if (
                strict
                and size is not None
                and size < float(min_font)
                and not is_template_or_meta_text(node, width, height, policy)
            ):
                audit.error(
                    "body-font-too-small",
                    "Visible PPT body text is below the absolute minimum font size.",
                    page=page_number,
                    text=node["text"],
                    font_size=size,
                    minimum=float(min_font),
                )

        audit_layout_collisions(
            page_number,
            nodes,
            image_items,
            width,
            height,
            policy,
            strict,
            audit,
        )
        audit_aesthetic_layout(
            page_number,
            root,
            nodes,
            image_items,
            width,
            height,
            slide_type,
            policy,
            strict,
            audit,
        )

        if slide_type not in STRUCTURAL_DEFAULT and not slide.get("density_exempt_reason"):
            visible_chars = len(re.sub(r"\s+", "", text))
            if visible_chars < min_chars and visual_count == 0:
                audit.warning(
                    "sparse-page",
                    "Page has little visible text and no figure/path-based visual proof.",
                    page=page_number,
                    visible_characters=visible_chars,
                )

        allowed_numbers = set(source_numbers)
        allowed_numbers.update(slide.get("_allowed_numbers", []))
        ignored = {str(page_number)}
        ignored.update(collect_numbers(str(slide.get("source_note", ""))))
        chapter_match = re.search(r"第\s*(\d+)\s*章", str(slide.get("chapter", "")))
        if chapter_match:
            ignored.add(chapter_match.group(1))
        unsupported = sorted(collect_numbers(text) - allowed_numbers - ignored)
        if unsupported:
            level = audit.error if strict else audit.warning
            level(
                "unsupported-number",
                "Visible number was not found in the normalized source corpus.",
                page=page_number,
                numbers=unsupported,
            )
    return len(files)


def pptx_text(archive: zipfile.ZipFile) -> str:
    values = []
    for name in archive.namelist():
        if not name.startswith("ppt/slides/slide") or not name.endswith(".xml"):
            continue
        root = etree.fromstring(archive.read(name))
        values.extend(root.xpath("//a:t/text()", namespaces={"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}))
    return "\n".join(values)


def iter_pptx_shapes(shapes):
    for shape in shapes:
        yield shape
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP and hasattr(shape, "shapes"):
            yield from iter_pptx_shapes(shape.shapes)


def pptx_paragraph_font_size(paragraph) -> float | None:
    if paragraph.font.size:
        return round(paragraph.font.size.pt, 2)
    sizes = [run.font.size.pt for run in paragraph.runs if run.font.size]
    return round(min(sizes), 2) if sizes else None


def is_pptx_meta_text(text: str, shape, slide_width: int, slide_height: int, policy: dict) -> bool:
    text = re.sub(r"\s+", " ", text or "").strip()
    if not text:
        return True
    exemptions = policy.get("font_exemptions", {})
    y = int(getattr(shape, "top", 0) or 0)
    x = int(getattr(shape, "left", 0) or 0)
    footer_y_min = float(exemptions.get("footer_y_min_ratio", 0.90)) * slide_height
    header_y_max = float(exemptions.get("header_y_max_ratio", 0.12)) * slide_height
    if y >= footer_y_min or y <= header_y_max:
        return True
    page_number_regex = exemptions.get("page_number_regex")
    if page_number_regex and re.fullmatch(str(page_number_regex), text):
        return x <= slide_width * 0.10 or x >= slide_width * 0.85
    if text in {"CONTENTS", "目录"}:
        return True
    return False


def is_pptx_layout_meta_text(text: str, shape, slide_width: int, slide_height: int, policy: dict) -> bool:
    text = re.sub(r"\s+", " ", text or "").strip()
    if not text:
        return True
    exemptions = policy.get("font_exemptions", {})
    y = int(getattr(shape, "top", 0) or 0)
    x = int(getattr(shape, "left", 0) or 0)
    footer_y_min = float(exemptions.get("footer_y_min_ratio", 0.90)) * slide_height
    if y >= footer_y_min:
        return True
    page_number_regex = exemptions.get("page_number_regex")
    if page_number_regex and re.fullmatch(str(page_number_regex), text):
        return x <= slide_width * 0.10 or x >= slide_width * 0.85 or y >= footer_y_min
    return False


def pptx_shape_bbox(shape) -> tuple[int, int, int, int] | None:
    try:
        left = int(shape.left)
        top = int(shape.top)
        width = int(shape.width)
        height = int(shape.height)
    except Exception:
        return None
    if width <= 0 or height <= 0:
        return None
    return (left, top, left + width, top + height)


def emu_to_pt(value: int | float) -> float:
    return float(value) / 914400 * 72


def visible_text_units(text: str) -> float:
    units = 0.0
    for char in text:
        if char.isspace():
            units += 0.28
        elif ord(char) > 127:
            units += 1.0
        else:
            units += 0.56
    return units


def estimate_text_lines(text: str, width_pt: float, font_size: float) -> int:
    text = re.sub(r"[ \t]+", " ", text or "").strip()
    if not text:
        return 0
    max_units = max(4.0, width_pt / max(font_size * 0.88, 1))
    lines = 0
    for paragraph in re.split(r"\n+", text):
        paragraph = paragraph.strip()
        if not paragraph:
            continue
        lines += max(1, math.ceil(visible_text_units(paragraph) / max_units))
    return lines


def paragraph_size_or_default(paragraph, default: float) -> float:
    size = pptx_paragraph_font_size(paragraph)
    return float(size or default)


def paragraph_has_marker_or_title(text: str, policy: dict) -> bool:
    settings = policy.get("paragraph_structure", {})
    marker = settings.get("marker_regex", r"^(?:\d{1,2}[、.．]|\(\d+\)|（\d+）)")
    value = re.sub(r"\s+", " ", text or "").strip()
    if not value:
        return True
    if re.search(str(marker), value):
        return True
    delimiters = str(settings.get("title_delimiters", "：:"))
    first_part = re.split(r"[。；;\n]", value, maxsplit=1)[0]
    return any(delimiter in first_part[:18] for delimiter in delimiters)


def run_has_emphasis(run, default_rgb: str | None = None) -> bool:
    if getattr(run.font, "bold", False):
        return True
    try:
        rgb_value = run.font.color.rgb
    except Exception:
        rgb_value = None
    if rgb_value is not None:
        value = str(rgb_value).upper()
        if default_rgb is None or value != default_rgb.upper():
            return True
    try:
        if run.font.highlight_color is not None:
            return True
    except Exception:
        pass
    return False


def shape_emphasis_run_count(shape, policy: dict) -> int:
    default_rgb = str(policy.get("default_body_rgb", "01203C"))
    count = 0
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    count += sum(1 for run in paragraph.runs if run.text.strip() and run_has_emphasis(run, default_rgb))
        return count
    if not getattr(shape, "has_text_frame", False):
        return 0
    for paragraph in shape.text_frame.paragraphs:
        runs = [run for run in paragraph.runs if run.text.strip()]
        count += sum(1 for run in runs if run_has_emphasis(run, default_rgb))
        if not runs and paragraph.text.strip():
            if getattr(paragraph.font, "bold", False):
                count += 1
            else:
                try:
                    rgb_value = paragraph.font.color.rgb
                except Exception:
                    rgb_value = None
                if rgb_value is not None and str(rgb_value).upper() != default_rgb.upper():
                    count += 1
    return count


def text_needs_emphasis(text: str, policy: dict) -> bool:
    settings = policy.get("emphasis", {})
    if settings.get("check_only_when_numbers_or_terms", True) is False:
        return True
    if collect_numbers(text):
        return True
    terms = [str(item) for item in settings.get("terms", [])]
    return any(term and term in text for term in terms)


def audit_pptx_text_structure_and_emphasis(
    presentation: Presentation,
    policy: dict,
    strict: bool,
    audit: Audit,
) -> None:
    if not strict:
        return
    structure = policy.get("paragraph_structure", {})
    structure_enabled = structure.get("enabled", True) is not False
    min_unlabeled = int(structure.get("min_unlabeled_paragraphs", 3))
    emphasis = policy.get("emphasis", {})
    emphasis_required = emphasis.get("required", True) is not False
    min_emphasis = int(emphasis.get("min_emphasized_runs_per_content_slide", 1))

    slide_width = int(presentation.slide_width)
    slide_height = int(presentation.slide_height)
    for slide_index, slide in enumerate(presentation.slides, start=1):
        all_text, body_text, visible_chars, _visual_area = slide_text_payload(slide, presentation, policy)
        if is_structural_pptx_slide(all_text):
            continue

        slide_emphasis_count = 0
        for shape in iter_pptx_shapes(slide.shapes):
            if getattr(shape, "has_table", False):
                slide_emphasis_count += shape_emphasis_run_count(shape, policy)
                continue
            if not getattr(shape, "has_text_frame", False):
                continue
            text = re.sub(r"\s+", " ", getattr(shape, "text", "") or "").strip()
            if is_pptx_layout_meta_text(text, shape, slide_width, slide_height, policy):
                continue
            slide_emphasis_count += shape_emphasis_run_count(shape, policy)
            paragraphs = [
                re.sub(r"\s+", " ", paragraph.text or "").strip()
                for paragraph in shape.text_frame.paragraphs
                if re.sub(r"\s+", " ", paragraph.text or "").strip()
            ]
            body_like = [item for item in paragraphs if len(item) >= 12]
            if (
                structure_enabled
                and len(body_like) >= min_unlabeled
                and sum(1 for item in body_like if paragraph_has_marker_or_title(item, policy)) < len(body_like) * 0.6
            ):
                audit.error(
                    "pptx-unlabeled-paragraph-group",
                    "A visible multi-paragraph text block lacks item numbers or short paragraph titles.",
                    page=slide_index,
                    paragraph_count=len(body_like),
                    text=compact_text(" ".join(body_like)),
                )

        if (
            emphasis_required
            and visible_chars >= 80
            and text_needs_emphasis(body_text, policy)
            and slide_emphasis_count < min_emphasis
        ):
            audit.error(
                "pptx-missing-emphasis",
                "Content slide with key terms or values has no visible bold/color/highlight emphasis runs.",
                page=slide_index,
                emphasized_runs=slide_emphasis_count,
                minimum=min_emphasis,
            )


def audit_pptx_text_capacity(
    slide_index: int,
    text: str,
    shape,
    font_size: float,
    policy: dict,
    audit: Audit,
    code: str = "pptx-text-overflows-box",
) -> None:
    if not text.strip():
        return
    settings = policy.get("text_fit", {})
    if settings.get("enabled", True) is False:
        return
    width_pt = emu_to_pt(int(getattr(shape, "width", 0) or 0))
    height_pt = emu_to_pt(int(getattr(shape, "height", 0) or 0))
    if width_pt <= 0 or height_pt <= 0:
        return
    line_height = font_size * float(settings.get("line_height_ratio", 1.24))
    available = max(1.0, height_pt / max(line_height, 1))
    required = estimate_text_lines(text, width_pt, font_size)
    tolerance = float(settings.get("overflow_tolerance_lines", 0.35))
    if required > available + tolerance:
        audit.error(
            code,
            "Visible PPTX text needs more lines than the textbox/table cell can hold.",
            page=slide_index,
            text=compact_text(text),
            required_lines=required,
            available_lines=round(available, 2),
            font_size=font_size,
        )


def slide_text_payload(slide, presentation: Presentation, policy: dict) -> tuple[str, str, int, float]:
    slide_width = int(presentation.slide_width)
    slide_height = int(presentation.slide_height)
    all_texts: list[str] = []
    body_texts: list[str] = []
    visual_area = 0.0
    slide_area = max(1, slide_width * slide_height)
    for shape in iter_pptx_shapes(slide.shapes):
        bbox = pptx_shape_bbox(shape)
        if getattr(shape, "has_table", False):
            if bbox:
                visual_area += box_area(bbox) / slide_area
            for row in shape.table.rows:
                for cell in row.cells:
                    text = re.sub(r"\s+", " ", cell.text or "").strip()
                    if text:
                        all_texts.append(text)
                        body_texts.append(text)
            continue
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
            if bbox:
                visual_area += box_area(bbox) / slide_area
            continue
        if getattr(shape, "has_text_frame", False):
            text = re.sub(r"\s+", " ", getattr(shape, "text", "") or "").strip()
            if text:
                all_texts.append(text)
                if not is_pptx_layout_meta_text(text, shape, slide_width, slide_height, policy):
                    body_texts.append(text)
    return "\n".join(all_texts), "\n".join(body_texts), len(re.sub(r"\s+", "", "\n".join(body_texts))), visual_area


def normalize_cell_text(value: str) -> str:
    value = re.sub(r"\s+", "", value or "")
    value = re.sub(r"[，,。；;：:（）()\[\]【】]", "", value)
    return value


def table_rows_from_shape(shape) -> list[list[str]]:
    rows: list[list[str]] = []
    if not getattr(shape, "has_table", False):
        return rows
    for row in shape.table.rows:
        rows.append([re.sub(r"\s+", " ", cell.text or "").strip() for cell in row.cells])
    return rows


def pptx_table_has_merge(shape) -> bool:
    if not getattr(shape, "has_table", False):
        return False
    for row in shape.table.rows:
        for cell in row.cells:
            if bool(getattr(cell, "is_merge_origin", False)) or bool(getattr(cell, "is_spanned", False)):
                return True
    return False


def source_table_candidates(catalog: dict, docx_models: dict) -> list[dict[str, Any]]:
    candidates: list[dict[str, Any]] = []
    model_numbers_by_locator: dict[str, set[str]] = {}
    for model in docx_models.get("tables", []) or []:
        locator = str(model.get("locator", ""))
        rows = model.get("rows") or []
        if not locator or not rows:
            continue
        model_numbers_by_locator[locator] = collect_numbers(
            "\n".join(" | ".join(map(str, row)) for row in rows)
        )
    for entry in catalog.get("entries", []):
        if entry.get("kind") != "table":
            continue
        rows = entry.get("rows") or []
        if not rows:
            continue
        locator = str(entry.get("locator", ""))
        numbers = set(collect_numbers(entry.get("text", "")))
        numbers.update(model_numbers_by_locator.get(locator, set()))
        candidates.append(
            {
                "id": entry.get("id", ""),
                "locator": locator,
                "source": entry.get("source", ""),
                "rows": rows,
                "row_count": int(entry.get("row_count") or len(rows)),
                "column_count": int(entry.get("column_count") or max((len(row) for row in rows), default=0)),
                "has_merged_cells": False,
                "numbers": numbers,
                "model": None,
            }
        )
    for model in docx_models.get("tables", []) or []:
        rows = model.get("rows") or []
        candidates.append(
            {
                "id": model.get("locator", ""),
                "locator": model.get("locator", ""),
                "source": model.get("source", docx_models.get("source_file", "")),
                "rows": rows,
                "row_count": int(model.get("row_count") or len(rows)),
                "column_count": int(model.get("column_count") or max((len(row) for row in rows), default=0)),
                "has_merged_cells": bool(model.get("has_merged_cells")),
                "numbers": collect_numbers("\n".join(" | ".join(map(str, row)) for row in rows)),
                "model": model,
            }
        )
    return candidates


def add_table_ir_candidates(candidates: list[dict[str, Any]], table_ir: dict) -> None:
    for table in table_ir.get("tables", []) if isinstance(table_ir, dict) else []:
        rows = table.get("rows") or []
        structure = table.get("structure", {})
        candidates.append(
            {
                "id": table.get("table_id", ""),
                "locator": table.get("locator", table.get("table_id", "")),
                "source": table.get("source", {}).get("file", ""),
                "rows": rows,
                "row_count": int(structure.get("rows") or len(rows)),
                "column_count": int(structure.get("cols") or max((len(row) for row in rows), default=0)),
                "has_merged_cells": bool(structure.get("merged_cells")),
                "numbers": collect_numbers("\n".join(" | ".join(map(str, row)) for row in rows)),
                "model": table,
                "render_mode": table.get("render_mode", "auto"),
            }
        )


def audit_table_ir(project: Path, table_ir: dict, catalog: dict, policy: dict, strict: bool, audit: Audit) -> None:
    settings = policy.get("table_fidelity", {})
    if not strict or settings.get("enabled", True) is False:
        return
    docx_or_pdf_sources = [
        item
        for item in catalog.get("sources", [])
        if str(item.get("name", "")).lower().endswith((".docx", ".pdf"))
    ]
    if settings.get("require_table_ir", True) and docx_or_pdf_sources and not table_ir:
        audit.error(
            "missing-table-ir",
            "DOCX/PDF engineering decks must generate analysis/table_ir.json so LLM does not reconstruct tables.",
            sources=[item.get("name") for item in docx_or_pdf_sources],
        )
        return

    allowed = set(settings.get("allowed_render_modes", ["native", "image", "hybrid", "auto"]))
    tables = table_ir.get("tables", []) if isinstance(table_ir, dict) else []
    catalog_table_ids = {
        str(item.get("table_id"))
        for item in catalog.get("entries", [])
        if item.get("kind") == "table" and item.get("table_id")
    }
    ir_table_ids = {str(item.get("table_id")) for item in tables if item.get("table_id")}
    missing_from_ir = sorted(catalog_table_ids - ir_table_ids)
    if missing_from_ir:
        audit.error(
            "catalog-table-missing-ir",
            "Source catalog table entries reference table_id values absent from analysis/table_ir.json.",
            table_ids=missing_from_ir[:12],
        )

    for table in tables:
        table_id = str(table.get("table_id", ""))
        mode = str(table.get("render_mode", "")).lower()
        if mode not in allowed:
            audit.error(
                "invalid-table-render-mode",
                "Table IR render_mode must be one of native, image, hybrid, or auto.",
                table_id=table_id,
                render_mode=mode,
            )
        structure = table.get("structure", {})
        complex_by_rule = (
            bool(structure.get("is_cross_page"))
            or bool(structure.get("has_diagonal_header"))
            or bool(structure.get("has_nested_table"))
            or int(structure.get("rows") or 0) > int(settings.get("small_medium_max_rows", 10))
            or int(structure.get("cols") or 0) > int(settings.get("small_medium_max_cols", 7))
            or int(structure.get("merged_cell_count") or len(structure.get("merged_cells") or [])) >= 3
            or float(structure.get("border_complexity") or 0.0) > 0.6
        )
        if settings.get("complex_tables_default_to_image", True) and complex_by_rule and mode == "native":
            audit.error(
                "complex-table-rendered-native",
                "Complex source tables must default to image/hybrid instead of native reconstruction.",
                table_id=table_id,
                rows=structure.get("rows"),
                cols=structure.get("cols"),
                merged_cell_count=structure.get("merged_cell_count"),
            )
        if settings.get("image_modes_require_crop_asset", True) and mode in {"image", "hybrid"}:
            asset = table.get("assets", {}).get("crop_image") or table.get("assets", {}).get("source_table_image")
            if not asset:
                audit.error(
                    "table-image-mode-missing-asset",
                    "Image/hybrid table render modes require a source-derived table crop asset.",
                    table_id=table_id,
                )
            elif not (project / asset).exists():
                audit.error(
                    "table-image-asset-missing",
                    "Table crop asset referenced by Table IR does not exist.",
                    table_id=table_id,
                    asset=asset,
                )


def flatten_cells(rows: list[list[str]]) -> list[str]:
    values: list[str] = []
    for row in rows:
        for cell in row:
            value = normalize_cell_text(str(cell))
            if value:
                values.append(value)
    return values


def table_match_score(ppt_rows: list[list[str]], source_rows: list[list[str]]) -> float:
    ppt_cells = set(flatten_cells(ppt_rows))
    source_cells = set(flatten_cells(source_rows))
    if not ppt_cells or not source_cells:
        return 0.0
    exact = len(ppt_cells.intersection(source_cells))
    partial = 0
    for cell in ppt_cells - source_cells:
        if len(cell) >= 4 and any(cell in source or source in cell for source in source_cells if len(source) >= 4):
            partial += 1
    return (exact + partial * 0.5) / max(len(ppt_cells), 1)


def best_source_table_match(ppt_rows: list[list[str]], candidates: list[dict[str, Any]]) -> tuple[dict[str, Any] | None, float]:
    best: dict[str, Any] | None = None
    best_score = 0.0
    for candidate in candidates:
        score = table_match_score(ppt_rows, candidate["rows"])
        if score > best_score:
            best = candidate
            best_score = score
    return best, best_score


def row_text_set(row: list[str]) -> set[str]:
    return {normalize_cell_text(str(cell)) for cell in row if normalize_cell_text(str(cell))}


def slide_has_split_or_excerpt_note(all_text: str) -> bool:
    return bool(re.search(r"节选|摘录|重点行|重点列|拆分|续表|分表|局部|裁剪", all_text or ""))


def audit_pptx_table_fidelity(
    project: Path,
    presentation: Presentation,
    catalog: dict,
    table_ir: dict,
    pages: dict[int, dict],
    policy: dict,
    strict: bool,
    audit: Audit,
) -> None:
    settings = policy.get("table_fidelity", {})
    if not strict or settings.get("enabled", True) is False:
        return

    docx_sources = [
        item
        for item in catalog.get("sources", [])
        if str(item.get("name", "")).lower().endswith(".docx")
    ]
    docx_models_path = project / "analysis" / "docx_table_models.json"
    docx_models = load_json(docx_models_path, audit, required=False) if docx_models_path.exists() else {}
    planned_table_pages = [
        page
        for page in pages.values()
        if str(page.get("source_mode", "")).upper() == "ORIGINAL_TABLE"
    ]
    if (
        settings.get("require_docx_table_models_for_docx_sources", True)
        and docx_sources
        and planned_table_pages
        and not docx_models
    ):
        audit.error(
            "missing-docx-table-models",
            "DOCX table-heavy deck needs analysis/docx_table_models.json before strict table rendering/checking.",
            docx_sources=[item.get("name") for item in docx_sources],
            planned_table_pages=[page.get("page") for page in planned_table_pages],
        )

    candidates = source_table_candidates(catalog, docx_models)
    add_table_ir_candidates(candidates, table_ir)
    if not candidates:
        return

    matched_tables = 0
    native_tables = 0
    small_rows = int(settings.get("small_medium_max_rows", 10))
    small_cols = int(settings.get("small_medium_max_cols", 7))
    header_missing_error_ratio = float(settings.get("header_missing_error_ratio", 0.4))

    for slide_index, slide in enumerate(presentation.slides, start=1):
        all_text, _body_text, _visible_chars, _visual_area = slide_text_payload(slide, presentation, policy)
        for shape_index, shape in enumerate(iter_pptx_shapes(slide.shapes), start=1):
            if not getattr(shape, "has_table", False):
                continue
            native_tables += 1
            ppt_rows = table_rows_from_shape(shape)
            if not ppt_rows:
                continue
            candidate, score = best_source_table_match(ppt_rows, candidates)
            if candidate is None or score < 0.16:
                audit.warning(
                    "pptx-table-source-unmatched",
                    "Native PPTX table could not be confidently matched to a source table.",
                    page=slide_index,
                    table=shape_index,
                    match_score=round(score, 3),
                )
                continue
            matched_tables += 1

            ppt_row_count = len(ppt_rows)
            ppt_col_count = max((len(row) for row in ppt_rows), default=0)
            source_row_count = int(candidate["row_count"])
            source_col_count = int(candidate["column_count"])
            source_is_small = source_row_count <= small_rows and source_col_count <= small_cols
            has_excerpt_note = slide_has_split_or_excerpt_note(all_text)

            if source_is_small and not has_excerpt_note and (ppt_row_count != source_row_count or ppt_col_count != source_col_count):
                audit.error(
                    "pptx-table-shape-mismatch",
                    "Small/medium source table should preserve row and column scale unless the slide explicitly records a split/excerpt.",
                    page=slide_index,
                    table=shape_index,
                    source=candidate["locator"],
                    ppt_rows=ppt_row_count,
                    ppt_cols=ppt_col_count,
                    source_rows=source_row_count,
                    source_cols=source_col_count,
                    match_score=round(score, 3),
                )

            source_header = row_text_set(candidate["rows"][0] if candidate["rows"] else [])
            ppt_header = row_text_set(ppt_rows[0] if ppt_rows else [])
            if source_header:
                missing = [cell for cell in source_header if not any(cell in item or item in cell for item in ppt_header)]
                missing_ratio = len(missing) / max(len(source_header), 1)
                if missing_ratio >= header_missing_error_ratio:
                    audit.error(
                        "pptx-table-header-mismatch",
                        "PPTX table header does not preserve enough source header cells.",
                        page=slide_index,
                        table=shape_index,
                        source=candidate["locator"],
                        missing_header_cells=missing[:6],
                        missing_ratio=round(missing_ratio, 3),
                    )

            ppt_numbers = collect_numbers("\n".join(" | ".join(row) for row in ppt_rows))
            source_numbers = set(candidate["numbers"])
            unsupported = sorted(ppt_numbers - source_numbers)
            if unsupported and settings.get("unsupported_table_number_is_error", True):
                audit.error(
                    "pptx-table-number-not-in-source-table",
                    "PPTX table contains numbers that are not present in its matched source table.",
                    page=slide_index,
                    table=shape_index,
                    source=candidate["locator"],
                    numbers=unsupported,
                )

            if (
                settings.get("native_merge_required_for_merged_small_tables", True)
                and source_is_small
                and candidate.get("has_merged_cells")
                and not pptx_table_has_merge(shape)
            ):
                audit.error(
                    "pptx-merged-table-flattened",
                    "Merged-cell source table was rendered without native PPTX merged cells.",
                    page=slide_index,
                    table=shape_index,
                    source=candidate["locator"],
                )

    if planned_table_pages and matched_tables == 0:
        if native_tables:
            audit.error(
                "planned-table-without-matched-native-table",
                "Deck plan includes original table slides, but no native PPTX table matched the source catalog/models.",
                planned_table_pages=[page.get("page") for page in planned_table_pages],
            )
        else:
            audit.warning(
                "planned-table-rendered-without-native-table",
                "Deck plan includes original table slides but no native PPTX table was available for structural comparison; verify source crops visually.",
                planned_table_pages=[page.get("page") for page in planned_table_pages],
            )


def table_ir_lookup(table_ir: dict) -> dict[str, dict]:
    lookup: dict[str, dict] = {}
    for table in table_ir.get("tables", []) or []:
        for key in (table.get("table_id"), table.get("locator")):
            if key:
                lookup[str(key)] = table
                lookup[str(key).replace("-", "")] = table
    return lookup


def min_table_source_font(table: dict) -> float:
    sizes: list[float] = []
    for cell in table.get("cells", []) or []:
        if not str(cell.get("text", "")).strip():
            continue
        try:
            sizes.append(float(cell.get("style", {}).get("font_size", 10.5)))
        except (TypeError, ValueError):
            continue
    return min(sizes) if sizes else 10.5


def audit_pptx_image_table_readability(
    project: Path,
    presentation: Presentation,
    table_ir: dict,
    policy: dict,
    strict: bool,
    audit: Audit,
) -> None:
    if not strict:
        return
    settings = policy.get("table_fidelity", {})
    minimum = float(settings.get("image_table_min_effective_pt", 10))
    source_dpi = float(settings.get("table_image_source_dpi", 220))
    lookup = table_ir_lookup(table_ir)
    if not lookup:
        return

    for slide_index, slide in enumerate(presentation.slides, start=1):
        for shape in iter_pptx_shapes(slide.shapes):
            name = str(getattr(shape, "name", "") or "")
            if "TableImage" not in name:
                continue
            match = re.search(r"TableImage(?:NeedsSplit)?[:：]([A-Za-z]-?\d{3})", name)
            table = lookup.get(match.group(1)) if match else None
            if table is None:
                continue
            asset = table.get("assets", {}).get("crop_image") or table.get("assets", {}).get("source_table_image")
            image_path = project / str(asset or "")
            if not image_path.exists():
                continue
            try:
                image_width, image_height = shape.image.size
            except Exception:
                with Image.open(image_path) as image:
                    image_width, image_height = image.size
            original_w_pt = image_width / source_dpi * 72.0
            original_h_pt = image_height / source_dpi * 72.0
            display_w_pt = emu_to_pt(int(getattr(shape, "width", 0) or 0))
            display_h_pt = emu_to_pt(int(getattr(shape, "height", 0) or 0))
            scale = min(
                display_w_pt / max(original_w_pt, 1.0),
                display_h_pt / max(original_h_pt, 1.0),
            )
            effective_pt = min_table_source_font(table) * scale
            if effective_pt < minimum:
                audit.error(
                    "pptx-image-table-font-too-small",
                    "Source table crop is displayed below the readable table-text floor; enlarge it or split it onto a new slide.",
                    page=slide_index,
                    table=table.get("locator") or table.get("table_id"),
                    effective_font_size=round(effective_pt, 2),
                    minimum=minimum,
                    displayed_size_pt=[round(display_w_pt, 1), round(display_h_pt, 1)],
                )


def is_structural_pptx_slide(all_text: str) -> bool:
    text = re.sub(r"\s+", " ", all_text or "").strip()
    if not text:
        return True
    if "CONTENTS" in text or text == "目录" or text.startswith("目录 "):
        return True
    if re.search(r"第\s*\d+\s*章\s*/\s*共\s*\d+\s*章", text):
        return True
    if "章节内容保留报告中的源表" in text and "后续页面以该章节证据为核心" in text:
        return True
    if "工程评审汇报" in text and "专项勘察" in text and "依据报告章节" in text:
        return True
    return False


def normalized_body_for_duplicate(text: str) -> str:
    text = re.sub(r"\b\d{1,3}\b", " ", text or "")
    text = re.sub(r"第\s*\d+\s*章\s*/\s*共\s*\d+\s*章", " ", text)
    text = re.sub(r"\s+", "", text)
    return text


def pptx_large_object_boxes(presentation: Presentation, slide) -> list[dict]:
    slide_area = int(presentation.slide_width) * int(presentation.slide_height)
    items: list[dict] = []
    for index, shape in enumerate(iter_pptx_shapes(slide.shapes), start=1):
        bbox = pptx_shape_bbox(shape)
        if bbox is None:
            continue
        area = box_area(bbox)
        if getattr(shape, "has_table", False):
            items.append({"id": f"table-{index}", "bbox": bbox})
        elif getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE and area >= slide_area * 0.03:
            items.append({"id": f"picture-{index}", "bbox": bbox})
    return items


def audit_pptx_layout_geometry(
    presentation: Presentation,
    policy: dict,
    strict: bool,
    audit: Audit,
) -> None:
    if not strict or not policy.get("layout_collision", {}).get("enabled", True):
        return
    collision = policy.get("layout_collision", {})
    slide_width = int(presentation.slide_width)
    slide_height = int(presentation.slide_height)
    margin_px = float(collision.get("canvas_margin_px", 2))
    margin = int(round(slide_width / 1280 * margin_px))
    text_overlap_min = float(collision.get("text_overlap_min_ratio", 0.35))
    text_object_overlap_min = float(collision.get("text_image_overlap_min_ratio", 0.15))

    for slide_index, slide in enumerate(presentation.slides, start=1):
        object_items = pptx_large_object_boxes(presentation, slide)
        text_items: list[dict] = []
        slide_full_text = " ".join(
            re.sub(r"\s+", " ", getattr(shape, "text", "") or "").strip()
            for shape in iter_pptx_shapes(slide.shapes)
            if getattr(shape, "has_text_frame", False)
        )
        slide_has_marker_2 = (chr(0xFF08) + "2" + chr(0xFF09)) in slide_full_text or "(2)" in slide_full_text
        for shape_index, shape in enumerate(iter_pptx_shapes(slide.shapes), start=1):
            if getattr(shape, "has_table", False) or not getattr(shape, "has_text_frame", False):
                continue
            text = re.sub(r"\s+", " ", getattr(shape, "text", "") or "").strip()
            if is_pptx_layout_meta_text(text, shape, slide_width, slide_height, policy):
                continue
            bbox = pptx_shape_bbox(shape)
            if bbox is None:
                continue
            text_items.append({"index": shape_index, "text": text, "bbox": bbox})
            sizes = [
                paragraph_size_or_default(paragraph, float(policy.get("minimum_font_px", {}).get("body_absolute", 14)))
                for paragraph in shape.text_frame.paragraphs
                if paragraph.text.strip()
            ]
            if sizes:
                audit_pptx_text_capacity(
                    slide_index,
                    getattr(shape, "text", "") or "",
                    shape,
                    min(sizes),
                    policy,
                    audit,
                )
            fullwidth_1 = chr(0xFF08) + "1" + chr(0xFF09)
            fullwidth_2 = chr(0xFF08) + "2" + chr(0xFF09)
            marker_positions = [pos for pos in [text.find(fullwidth_1), text.find("(1)")] if pos >= 0]
            marker_position = min(marker_positions) if marker_positions else -1
            if (
                marker_position >= 0
                and fullwidth_2 not in text
                and "(2)" not in text
                and not slide_has_marker_2
                and (
                    marker_position > len(text) * 0.55
                    or text.rstrip().endswith(fullwidth_1)
                    or text.rstrip().endswith("(1)")
                    or len(text[marker_position:]) <= 90
                )
            ):
                audit.error(
                    "pptx-orphan-list-marker",
                    "Visible text contains an isolated Chinese list marker without the next item.",
                    page=slide_index,
                    text=compact_text(text),
                )
            if re.search(r"(?:；(?:份|幅|个|单位|数量)；\d|单位；数量|取得成果；单位；数量)", text):
                audit.error(
                    "pptx-table-row-leaked-into-body",
                    "Table-like row text appears in a narrative text box.",
                    page=slide_index,
                    text=compact_text(text),
                )
            if len(text) >= 45 and (text[-1] in "，,、；;：" or re.search(r"(?:与|和|及|至|为|在|位于)$", text)):
                audit.error(
                    "pptx-probable-incomplete-sentence",
                    "Visible paragraph appears to end before the sentence is complete.",
                    page=slide_index,
                    text=compact_text(text),
                )
            if (
                bbox[0] < margin
                or bbox[1] < margin
                or bbox[2] > slide_width - margin
                or bbox[3] > slide_height - margin
            ):
                audit.error(
                    "pptx-text-outside-safe-frame",
                    "PPTX text box crosses the slide safe frame after final font/layout adjustments.",
                    page=slide_index,
                    text=compact_text(text),
                    bbox=list(bbox),
                )

        for left_index, left_item in enumerate(text_items):
            for right_item in text_items[left_index + 1 :]:
                overlap = intersection_area(left_item["bbox"], right_item["bbox"])
                if not overlap:
                    continue
                denominator = min(box_area(left_item["bbox"]), box_area(right_item["bbox"]))
                if denominator and overlap / denominator >= text_overlap_min:
                    audit.error(
                        "pptx-probable-text-overlap",
                        "Two PPTX text boxes overlap after final font/layout adjustments.",
                        page=slide_index,
                        text_a=compact_text(left_item["text"]),
                        text_b=compact_text(right_item["text"]),
                        overlap_ratio=round(overlap / denominator, 3),
                    )

        for text_item in text_items:
            text_area = box_area(text_item["bbox"])
            if not text_area:
                continue
            for object_item in object_items:
                overlap = intersection_area(text_item["bbox"], object_item["bbox"])
                if overlap / text_area >= text_object_overlap_min:
                    audit.error(
                        "pptx-text-overlaps-object",
                        "PPTX text box overlaps a source figure/table after final font/layout adjustments.",
                        page=slide_index,
                        text=compact_text(text_item["text"]),
                        object=object_item["id"],
                        overlap_ratio=round(overlap / text_area, 3),
                    )


def audit_pptx_font_sizes(presentation: Presentation, policy: dict, audit: Audit) -> None:
    minimums = policy.get("minimum_font_px", {})
    body_min = float(minimums.get("body_absolute", 14))
    table_min = float(minimums.get("table_absolute", 12))
    if table_min <= 0:
        table_min = 12

    slide_width = int(presentation.slide_width)
    slide_height = int(presentation.slide_height)
    for slide_index, slide in enumerate(presentation.slides, start=1):
        for shape in iter_pptx_shapes(slide.shapes):
            if getattr(shape, "has_table", False):
                row_count = max(1, len(shape.table.rows))
                col_count = max(1, len(shape.table.columns))
                for row_index, row in enumerate(shape.table.rows, start=1):
                    for col_index, cell in enumerate(row.cells, start=1):
                        pseudo_cell = type(
                            "CellBox",
                            (),
                            {
                                "width": int(getattr(shape, "width", 0) or 0) / col_count,
                                "height": int(getattr(row, "height", 0) or 0)
                                or int(getattr(shape, "height", 0) or 0) / row_count,
                            },
                        )()
                        cell_sizes: list[float] = []
                        for paragraph in cell.text_frame.paragraphs:
                            text = paragraph.text.strip()
                            size = pptx_paragraph_font_size(paragraph)
                            if text:
                                cell_sizes.append(float(size or table_min))
                            if text and size is not None and size < table_min:
                                audit.error(
                                    "pptx-table-font-too-small",
                                    "PPTX table text is below the minimum font size.",
                                    page=slide_index,
                                    row=row_index,
                                    column=col_index,
                                    text=compact_text(text),
                                    font_size=size,
                                    minimum=table_min,
                                )
                        if cell.text.strip() and cell_sizes:
                            audit_pptx_text_capacity(
                                slide_index,
                                cell.text,
                                pseudo_cell,
                                min(cell_sizes),
                                policy,
                                audit,
                                code="pptx-table-cell-overflow",
                            )
                continue

            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                size = pptx_paragraph_font_size(paragraph)
                if (
                    text
                    and size is not None
                    and size < body_min
                    and not is_pptx_meta_text(text, shape, slide_width, slide_height, policy)
                ):
                    audit.error(
                        "pptx-body-font-too-small",
                        "PPTX body text is below the minimum font size.",
                        page=slide_index,
                        text=compact_text(text),
                        font_size=size,
                        minimum=body_min,
                    )


def audit_pptx_density_and_duplicates(
    presentation: Presentation,
    policy: dict,
    strict: bool,
    audit: Audit,
) -> None:
    if not strict:
        return
    sparse = policy.get("sparse_page", {})
    min_chars_without_visual = int(
        sparse.get(
            "minimum_visible_characters_without_visual",
            sparse.get("minimum_visible_characters_native_no_visual", 160),
        )
    )
    min_chars_with_visual = int(sparse.get("minimum_visible_characters_with_visual", 70))
    min_visual_area = float(sparse.get("minimum_visual_area_ratio", 0.18))
    duplicate_policy = policy.get("duplicate_content", {})
    duplicate_enabled = duplicate_policy.get("enabled", True) is not False
    duplicate_ratio = float(duplicate_policy.get("max_consecutive_body_similarity", 0.90))
    min_duplicate_chars = int(duplicate_policy.get("minimum_body_characters", 80))

    previous_body = ""
    previous_index = 0
    for slide_index, slide in enumerate(presentation.slides, start=1):
        all_text, body_text, visible_chars, visual_area = slide_text_payload(slide, presentation, policy)
        structural = is_structural_pptx_slide(all_text)
        if not structural:
            minimum = min_chars_with_visual if visual_area >= min_visual_area else min_chars_without_visual
            if visible_chars < minimum:
                audit.error(
                    "pptx-sparse-page",
                    "Content slide has too little visible material for its visual area.",
                    page=slide_index,
                    visible_characters=visible_chars,
                    visual_area_ratio=round(visual_area, 3),
                    minimum=minimum,
                )

        normalized = normalized_body_for_duplicate(body_text)
        if (
            duplicate_enabled
            and not structural
            and previous_body
            and len(normalized) >= min_duplicate_chars
            and len(previous_body) >= min_duplicate_chars
        ):
            ratio = difflib.SequenceMatcher(None, previous_body, normalized).ratio()
            if ratio >= duplicate_ratio:
                audit.error(
                    "pptx-duplicate-consecutive-content",
                    "Two consecutive content slides have near-identical body text.",
                    page=slide_index,
                    previous_page=previous_index,
                    similarity=round(ratio, 3),
                )
        if structural:
            previous_body = ""
            previous_index = 0
        else:
            previous_body = normalized
            previous_index = slide_index


def planned_section_divider_count(plan: dict, policy: dict) -> int:
    settings = policy.get("section_dividers", {})
    if settings.get("required", True) is False:
        return 0
    structural = set(policy.get("structural_slide_types", STRUCTURAL_DEFAULT))
    structural.update({"cover", "agenda", "section", "closing"})
    content_chapters: list[str] = []

    def add_once(chapter: str) -> None:
        chapter = re.sub(r"\s+", " ", chapter or "").strip()
        if chapter and chapter not in content_chapters:
            content_chapters.append(chapter)

    for slide in plan.get("slides", []):
        if str(slide.get("type", "")).lower() not in structural:
            add_once(str(slide.get("chapter", "")))
    ordered: list[str] = []
    for chapter in plan.get("deck", {}).get("chapter_order", []) or []:
        chapter = re.sub(r"\s+", " ", str(chapter)).strip()
        if chapter in content_chapters and chapter not in ordered:
            ordered.append(chapter)
    for chapter in content_chapters:
        if chapter not in ordered:
            ordered.append(chapter)
    return len(ordered)


def expected_pptx_slide_count(plan: dict, pages: dict[int, dict], svg_count: int, policy: dict) -> int:
    base = len(pages) or svg_count
    if not base:
        return 0
    return base + planned_section_divider_count(plan, policy)


def audit_pptx(
    project: Path,
    path: Path,
    expected_slides: int,
    policy: dict,
    catalog: dict,
    table_ir: dict,
    pages: dict[int, dict],
    audit: Audit,
    strict: bool,
) -> None:
    if not path.exists():
        audit.error("missing-pptx", "Requested PPTX does not exist.", path=display_path(path))
        return
    try:
        with zipfile.ZipFile(path) as archive:
            bad = archive.testzip()
            if bad:
                audit.error("corrupt-pptx", "PPTX contains a corrupt ZIP entry.", entry=bad)
            xml_parts = [
                name for name in archive.namelist() if name.endswith(".xml") or name.endswith(".rels")
            ]
            for name in xml_parts:
                etree.fromstring(archive.read(name))
            for name in archive.namelist():
                if name.startswith("ppt/media/") and not name.endswith("/") and not archive.read(name):
                    audit.error("empty-media", "PPTX contains an empty media file.", entry=name)
            scan_forbidden(
                pptx_text(archive),
                visible_forbidden_phrases(policy),
                audit,
                path.name,
            )
        presentation = Presentation(path)
        actual = len(presentation.slides)
        audit_pptx_font_sizes(presentation, policy, audit)
        audit_final_text_review_pptx(presentation, policy, strict, audit)
        audit_pptx_text_structure_and_emphasis(presentation, policy, strict, audit)
        audit_pptx_layout_geometry(presentation, policy, strict, audit)
        audit_pptx_table_fidelity(project, presentation, catalog, table_ir, pages, policy, strict, audit)
        audit_pptx_image_table_readability(project, presentation, table_ir, policy, strict, audit)
        audit_pptx_density_and_duplicates(presentation, policy, strict, audit)
        if expected_slides and actual != expected_slides:
            audit.error(
                "slide-count-mismatch",
                "PPTX slide count does not match deck plan.",
                expected=expected_slides,
                actual=actual,
            )
        ratio = presentation.slide_width / presentation.slide_height
        if not math.isclose(ratio, 16 / 9, rel_tol=0.01):
            audit.warning("unexpected-aspect-ratio", "PPTX is not 16:9.", ratio=ratio)
        audit.info("pptx-valid", "PPTX package, XML, media, and parser checks passed.", slides=actual)
    except Exception as exc:
        audit.error("pptx-parse-failed", "PPTX validation failed.", error=str(exc))


def write_reports(project: Path, audit: Audit, metadata: dict) -> None:
    qa = project / "qa"
    qa.mkdir(parents=True, exist_ok=True)
    payload = {"summary": audit.summary(), "metadata": metadata, "findings": audit.items}
    (qa / "release_audit.json").write_text(
        json.dumps(payload, ensure_ascii=False, indent=2) + "\n",
        encoding="utf-8",
    )
    lines = [
        "# Release Audit",
        "",
        f"- Errors: {payload['summary']['errors']}",
        f"- Warnings: {payload['summary']['warnings']}",
        f"- Release ready: {payload['summary']['release_ready']}",
        "",
    ]
    for item in audit.items:
        context = json.dumps(item["context"], ensure_ascii=False)
        lines.append(f"- **{item['level'].upper()} {item['code']}**: {item['message']} `{context}`")
    lines.append("")
    (qa / "release_audit.md").write_text("\n".join(lines), encoding="utf-8")


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("project", type=Path)
    parser.add_argument("--pptx", type=Path)
    parser.add_argument("--strict", action="store_true")
    args = parser.parse_args()
    project = args.project.resolve()
    audit = Audit()
    policy = load_json(project / "qa" / "release_policy.json", audit, required=args.strict)
    plan = load_json(project / "deck_plan.json", audit, required=args.strict)
    ledger = load_json(project / "evidence_ledger.json", audit, required=args.strict)
    catalog = load_json(project / "analysis" / "source_catalog.json", audit, required=args.strict)
    table_ir = load_json(project / "analysis" / "table_ir.json", audit, required=False)
    content_inventory = audit_content_analysis(project, args.strict, audit)
    audit_ledger(ledger, catalog, args.strict, audit)
    audit_table_ir(project, table_ir, catalog, policy, args.strict, audit)
    source_text = "\n".join(item.get("text", "") for item in catalog.get("entries", []))
    source_numbers = collect_numbers(source_text)
    pages = audit_plan(project, plan, ledger, policy, args.strict, audit)
    svg_count = audit_svgs(project, pages, policy, source_numbers, args.strict, audit)
    expected_pptx_slides = expected_pptx_slide_count(plan, pages, svg_count, policy)
    if args.pptx:
        audit_pptx(
            project,
            args.pptx.resolve(),
            expected_pptx_slides,
            policy,
            catalog,
            table_ir,
            pages,
            audit,
            args.strict,
        )
    metadata = {
        "project": display_path(project),
        "strict": args.strict,
        "planned_slides": len(pages),
        "expected_pptx_slides": expected_pptx_slides,
        "svg_slides": svg_count,
        "content_units": content_inventory.get("summary", {}).get("content_units", 0) if content_inventory else 0,
        "pptx": display_path(args.pptx) if args.pptx else None,
    }
    write_reports(project, audit, metadata)
    print(json.dumps({"summary": audit.summary(), **metadata}, ensure_ascii=False))
    sys.exit(1 if audit.summary()["errors"] else 0)


if __name__ == "__main__":
    main()
