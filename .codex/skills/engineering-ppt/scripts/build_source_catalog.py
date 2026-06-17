#!/usr/bin/env python3
"""Build a source inventory and addressable evidence catalog from report files."""

from __future__ import annotations

import argparse
import hashlib
import json
import re
import zipfile
from pathlib import Path
from typing import Iterable

from lxml import etree
from openpyxl import load_workbook
from pptx import Presentation

try:
    import fitz
except ImportError:  # pragma: no cover
    fitz = None


SUPPORTED = {".docx", ".pdf", ".txt", ".md", ".markdown", ".xlsx", ".xlsm", ".pptx"}
GENERATED_SOURCE_NAMES = {
    "chapter_coverage.md",
    "deck_plan.json",
    "design_spec.md",
    "evidence_ledger.json",
    "project_config.json",
    "spec_lock.md",
}
ROOT_SOURCE_SUFFIXES = {".docx", ".pdf", ".xlsx", ".xlsm", ".pptx"}
NUMBER_RE = re.compile(r"(?<![\w.])[-+]?\d+(?:\.\d+)?%?(?![\w.])")
WORD_NS = {
    "w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "pr": "http://schemas.openxmlformats.org/package/2006/relationships",
}


def display_path(path: Path, base: Path | None = None) -> str:
    resolved = path.resolve()
    roots = [base.resolve()] if base else []
    roots.append(Path.cwd().resolve())
    for root in roots:
        try:
            return resolved.relative_to(root).as_posix()
        except ValueError:
            continue
    return resolved.name


def clean_text(value: str) -> str:
    return re.sub(r"\s+", " ", value or "").strip()


def numbers(text: str) -> list[str]:
    return NUMBER_RE.findall(text)


def short_hash(path: Path) -> str:
    digest = hashlib.sha1(str(path.resolve()).encode("utf-8")).hexdigest()
    return digest[:8]


def entry(source: Path, locator: str, kind: str, text: str, **extra: object) -> dict:
    cleaned = clean_text(text)
    item = {
        "id": f"{short_hash(source)}:{locator}",
        "source": source.name,
        "source_path": display_path(source),
        "locator": locator,
        "kind": kind,
        "text": cleaned,
        "numbers": numbers(cleaned),
    }
    item.update(extra)
    return item


def extract_docx(path: Path) -> list[dict]:
    items: list[dict] = []
    with zipfile.ZipFile(path) as archive:
        document = etree.fromstring(archive.read("word/document.xml"))
        style_names: dict[str, str] = {}
        if "word/styles.xml" in archive.namelist():
            styles = etree.fromstring(archive.read("word/styles.xml"))
            for style in styles.xpath("//w:style", namespaces=WORD_NS):
                style_id = style.get(f"{{{WORD_NS['w']}}}styleId", "")
                names = style.xpath("./w:name/@w:val", namespaces=WORD_NS)
                style_names[style_id] = names[0] if names else style_id

        paragraph_index = 0
        table_index = 0
        body = document.find("w:body", WORD_NS)
        for block in list(body) if body is not None else []:
            local = etree.QName(block).localname
            if local == "p":
                paragraph_index += 1
                text = clean_text("".join(block.xpath(".//w:t/text()", namespaces=WORD_NS)))
                if not text:
                    continue
                style_ids = block.xpath("./w:pPr/w:pStyle/@w:val", namespaces=WORD_NS)
                style_id = style_ids[0] if style_ids else ""
                style = style_names.get(style_id, style_id)
                kind = (
                    "heading"
                    if style.lower().startswith("heading") or "标题" in style
                    else "paragraph"
                )
                items.append(entry(path, f"P{paragraph_index:04d}", kind, text, style=style))
            elif local == "tbl":
                table_index += 1
                rows = []
                for row in block.xpath("./w:tr", namespaces=WORD_NS):
                    cells = []
                    for cell in row.xpath("./w:tc", namespaces=WORD_NS):
                        cells.append(
                            clean_text(" ".join(cell.xpath(".//w:t/text()", namespaces=WORD_NS)))
                        )
                    rows.append(cells)
                text = "\n".join(" | ".join(row) for row in rows)
                items.append(
                    entry(
                        path,
                        f"T{table_index:03d}",
                        "table",
                        text,
                        rows=rows,
                        row_count=len(rows),
                        column_count=max((len(row) for row in rows), default=0),
                    )
                )

        for index, name in enumerate(
            sorted(item for item in archive.namelist() if item.startswith("word/media/")),
            start=1,
        ):
            info = archive.getinfo(name)
            items.append(
                entry(
                    path,
                    f"I{index:04d}",
                    "figure",
                    name,
                    asset=name,
                    bytes=info.file_size,
                )
            )
    return items


def extract_pdf(path: Path) -> list[dict]:
    if fitz is None:
        raise RuntimeError("PyMuPDF is required to catalog PDF sources.")
    items: list[dict] = []
    document = fitz.open(path)
    for page_number, page in enumerate(document, start=1):
        text = clean_text(page.get_text("text"))
        if text:
            items.append(entry(path, f"PAGE{page_number:04d}", "page_text", text))
        for image_index, image in enumerate(page.get_images(full=True), start=1):
            items.append(
                entry(
                    path,
                    f"PAGE{page_number:04d}-I{image_index:03d}",
                    "figure",
                    f"PDF embedded image xref={image[0]}",
                    page=page_number,
                    xref=image[0],
                )
            )
    return items


def extract_text(path: Path) -> list[dict]:
    text = path.read_text(encoding="utf-8", errors="replace")
    items: list[dict] = []
    for index, line in enumerate(text.splitlines(), start=1):
        cleaned = clean_text(line)
        if cleaned:
            kind = "heading" if cleaned.startswith("#") else "paragraph"
            items.append(entry(path, f"L{index:05d}", kind, cleaned))
    return items


def extract_workbook(path: Path) -> list[dict]:
    workbook = load_workbook(path, read_only=True, data_only=False)
    items: list[dict] = []
    for sheet in workbook.worksheets:
        rows = []
        for row in sheet.iter_rows(values_only=True):
            values = [clean_text("" if value is None else str(value)) for value in row]
            if any(values):
                rows.append(values)
        text = "\n".join(" | ".join(row) for row in rows)
        if text:
            items.append(
                entry(
                    path,
                    f"SHEET:{sheet.title}",
                    "table",
                    text,
                    sheet=sheet.title,
                    rows=rows,
                )
            )
    return items


def extract_pptx(path: Path) -> list[dict]:
    presentation = Presentation(path)
    items: list[dict] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and clean_text(shape.text):
                texts.append(clean_text(shape.text))
        if texts:
            items.append(entry(path, f"SLIDE{slide_number:04d}", "slide_text", "\n".join(texts)))
    return items


def discover_sources(project: Path, explicit: list[Path]) -> list[Path]:
    if explicit:
        candidates = explicit
    else:
        candidates = []
        sources_root = project / "sources"
        if sources_root.exists():
            candidates.extend(sources_root.glob("*"))
        if project.exists():
            candidates.extend(
                path
                for path in project.glob("*")
                if path.suffix.lower() in ROOT_SOURCE_SUFFIXES
            )
    unique: dict[Path, Path] = {}
    for candidate in candidates:
        path = candidate.resolve()
        if (
            path.is_file()
            and path.suffix.lower() in SUPPORTED
            and not path.name.startswith("~$")
            and path.name not in GENERATED_SOURCE_NAMES
            and "analysis" not in path.parts
            and "qa" not in path.parts
            and "exports" not in path.parts
        ):
            unique[path] = path
    return sorted(unique.values(), key=lambda value: value.name.lower())


def extract(path: Path) -> list[dict]:
    suffix = path.suffix.lower()
    if suffix == ".docx":
        return extract_docx(path)
    if suffix == ".pdf":
        return extract_pdf(path)
    if suffix in {".txt", ".md", ".markdown"}:
        return extract_text(path)
    if suffix in {".xlsx", ".xlsm"}:
        return extract_workbook(path)
    if suffix == ".pptx":
        return extract_pptx(path)
    return []


def write_inventory(project: Path, sources: Iterable[Path], catalog: list[dict]) -> None:
    analysis = project / "analysis"
    analysis.mkdir(parents=True, exist_ok=True)
    inventory = []
    for source in sources:
        inventory.append(
            {
                "name": source.name,
                "path": display_path(source, project),
                "size": source.stat().st_size,
                "modified": source.stat().st_mtime,
            }
        )
    payload = {
        "schema_version": "1.0",
        "sources": inventory,
        "entries": catalog,
    }
    (analysis / "source_catalog.json").write_text(
        json.dumps(payload, ensure_ascii=False, indent=2) + "\n",
        encoding="utf-8",
    )
    corpus = "\n".join(item["text"] for item in catalog if item["kind"] != "figure")
    (analysis / "source_corpus.txt").write_text(corpus + "\n", encoding="utf-8")

    lines = ["# Source Inventory", "", "| File | Size | Catalog entries |", "|---|---:|---:|"]
    for source in sources:
        count = sum(1 for item in catalog if item["source_path"] == display_path(source))
        lines.append(f"| {source.name} | {source.stat().st_size} | {count} |")
    lines.extend(
        [
            "",
            "Use `analysis/source_catalog.json` IDs in `evidence_ledger.json`.",
            "Do not treat extraction output as verified when tables or OCR are visibly damaged.",
            "",
        ]
    )
    (analysis / "source_inventory.md").write_text("\n".join(lines), encoding="utf-8")


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("project", type=Path)
    parser.add_argument("--source", action="append", type=Path, default=[])
    args = parser.parse_args()
    project = args.project.resolve()
    sources = discover_sources(project, args.source)
    if not sources:
        raise SystemExit("No supported source files found.")
    catalog: list[dict] = []
    failures = []
    for source in sources:
        try:
            catalog.extend(extract(source))
        except Exception as exc:  # continue inventorying other authoritative files
            failures.append({"source": str(source), "error": str(exc)})
    write_inventory(project, sources, catalog)
    print(
        json.dumps(
            {
                "sources": len(sources),
                "catalog_entries": len(catalog),
                "failures": failures,
                "output": str(project / "analysis" / "source_catalog.json"),
            },
            ensure_ascii=False,
        )
    )
    if failures:
        raise SystemExit(2)


if __name__ == "__main__":
    main()
