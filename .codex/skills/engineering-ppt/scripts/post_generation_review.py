#!/usr/bin/env python3
"""Post-generation review and auto-repair pipeline for engineering PPTX files."""

from __future__ import annotations

import argparse
import difflib
import json
import math
import re
import subprocess
import sys
from copy import deepcopy
from dataclasses import asdict, dataclass
from datetime import datetime
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pptx.util import Pt

from text_fitting import INTERNAL_FORBIDDEN_PHRASES, SlideContentSanitizer, compress_for_ppt, normalize_text


MAX_REVIEW_REPAIR_ROUNDS = 3
ISSUE_TYPES = {
    "text_overflow",
    "text_overlap",
    "content_too_long",
    "internal_note_leak",
    "factual_inconsistency",
    "missing_required_content",
    "table_format_error",
    "image_distortion",
    "layout_unbalanced",
    "font_below_minimum",
    "caption_too_long",
    "duplicate_content",
    "duplicate_slide",
    "duplicate_table",
    "duplicate_figure",
    "empty_placeholder",
    "low_readability",
    "source_untraceable",
}
PLACEHOLDER_RE = re.compile(r"\{\{?[^{}]{1,40}\}?\}|待补充|此处插入图片|此处填写|未命名章节|\b(?:XXX|None|null|N/A)\b", re.I)
NUMBER_RE = re.compile(r"(?<![\w.])[-+]?\d+(?:\.\d+)?%?(?![\w.])")
META_TEXT_RE = re.compile(r"^\d{1,3}$|^第\d{1,2}章|^CONTENTS$|^目录$")
BODY_MIN_PT = 14.0
TABLE_MIN_PT = 12.0
EMU_PER_INCH = 914400
EMPHASIS_RGB = RGBColor(13, 114, 181)
EMPHASIS_TERMS = [
    "范围",
    "目标",
    "责任",
    "措施",
    "风险",
    "结论",
    "投资",
    "验收",
    "监测",
    "水土流失",
    "防治责任范围",
]


@dataclass
class ReviewIssue:
    issue_id: str
    slide_index: int
    issue_type: str
    severity: str
    location: str
    description: str
    auto_fixable: bool
    suggested_fix: str
    status: str = "pending"
    review_stage: str = "PostGenerationReview"
    repair_action: str = ""


def compact(text: str, limit: int = 90) -> str:
    value = normalize_text(text)
    return value if len(value) <= limit else value[:limit].rstrip(" ，,；;。") + "..."


def issue(
    issues: list[ReviewIssue],
    round_idx: int,
    slide_index: int,
    issue_type: str,
    severity: str,
    location: str,
    description: str,
    auto_fixable: bool,
    suggested_fix: str,
    stage: str,
) -> None:
    if issue_type not in ISSUE_TYPES:
        issue_type = "layout_unbalanced"
    issues.append(
        ReviewIssue(
            issue_id=f"S{slide_index:03d}-{issue_type.upper().replace('_', '-')}-{len(issues)+1:03d}",
            slide_index=slide_index,
            issue_type=issue_type,
            severity=severity,
            location=location,
            description=description,
            auto_fixable=auto_fixable,
            suggested_fix=suggested_fix,
            review_stage=stage,
        )
    )


def iter_shapes(slide):
    for shape in slide.shapes:
        yield shape


def shape_location(shape, fallback: str) -> str:
    return str(getattr(shape, "name", "") or fallback)


def shape_bbox(shape) -> tuple[int, int, int, int]:
    left = int(getattr(shape, "left", 0) or 0)
    top = int(getattr(shape, "top", 0) or 0)
    width = int(getattr(shape, "width", 0) or 0)
    height = int(getattr(shape, "height", 0) or 0)
    return left, top, left + width, top + height


def box_area(box: tuple[int, int, int, int]) -> int:
    return max(0, box[2] - box[0]) * max(0, box[3] - box[1])


def intersection_area(a: tuple[int, int, int, int], b: tuple[int, int, int, int]) -> int:
    left = max(a[0], b[0])
    top = max(a[1], b[1])
    right = min(a[2], b[2])
    bottom = min(a[3], b[3])
    return max(0, right - left) * max(0, bottom - top)


def paragraph_font_size(paragraph, fallback: float = BODY_MIN_PT) -> float:
    sizes = [run.font.size.pt for run in paragraph.runs if run.text.strip() and run.font.size is not None]
    if sizes:
        return min(float(size) for size in sizes)
    if paragraph.font.size is not None:
        return float(paragraph.font.size.pt)
    return fallback


def is_meta_text(text: str) -> bool:
    value = normalize_text(text)
    return not value or bool(META_TEXT_RE.fullmatch(value))


def source_numbers(project: Path) -> set[str]:
    catalog_path = project / "analysis" / "source_catalog.json"
    if not catalog_path.exists():
        return set()
    catalog = json.loads(catalog_path.read_text(encoding="utf-8-sig"))
    text = "\n".join(str(item.get("text", "")) for item in catalog.get("entries", []))
    return {normalize_number(value) for value in NUMBER_RE.findall(text)}


def load_policy(project: Path) -> dict[str, Any]:
    policy_path = project / "qa" / "release_policy.json"
    if not policy_path.exists():
        return {}
    try:
        return json.loads(policy_path.read_text(encoding="utf-8-sig"))
    except Exception:
        return {}


def normalize_number(value: str) -> str:
    suffix = "%" if value.endswith("%") else ""
    raw = value[:-1] if suffix else value
    try:
        number = float(raw)
    except ValueError:
        return value
    if math.isclose(number, round(number), abs_tol=1e-10):
        return f"{int(round(number))}{suffix}"
    return f"{number:.10f}".rstrip("0").rstrip(".") + suffix


def collect_body_text(slide) -> str:
    parts: list[str] = []
    for shape in iter_shapes(slide):
        if getattr(shape, "has_text_frame", False):
            text = normalize_text(getattr(shape, "text", "") or "")
            if text and not is_meta_text(text):
                parts.append(text)
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                for cell in row.cells:
                    text = normalize_text(cell.text)
                    if text:
                        parts.append(text)
    return "\n".join(parts)


def canonical_duplicate_text(text: str) -> str:
    value = normalize_text(text)
    value = re.sub(r"\b\d{1,3}\b", " ", value)
    value = re.sub(r"[^\w\u4e00-\u9fff]+", "", value)
    return value.lower()


def text_similarity(left: str, right: str) -> float:
    if not left or not right:
        return 0.0
    return difflib.SequenceMatcher(None, left, right).ratio()


def shape_area_ratio(prs: Presentation, shape) -> float:
    slide_area = max(1, int(prs.slide_width) * int(prs.slide_height))
    return box_area(shape_bbox(shape)) / slide_area


def picture_sha(shape) -> str:
    try:
        return str(shape.image.sha1)
    except Exception:
        try:
            return str(hash(shape.image.blob))
        except Exception:
            return ""


def table_fingerprint(shape) -> str:
    name = shape_location(shape, "")
    if name.startswith("TableImage:"):
        return f"image-table:{name.split(':', 1)[1].strip()}"
    if getattr(shape, "has_table", False):
        cells: list[str] = []
        for row in shape.table.rows:
            for cell in row.cells:
                value = canonical_duplicate_text(cell.text)
                if value:
                    cells.append(value)
        return "native-table:" + "|".join(cells)
    if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE and "table" in name.lower():
        sha = picture_sha(shape)
        return f"image-table:{sha}" if sha else ""
    return ""


def figure_fingerprint(shape) -> str:
    if getattr(shape, "shape_type", None) != MSO_SHAPE_TYPE.PICTURE:
        return ""
    name = shape_location(shape, "").lower()
    if "table" in name:
        return ""
    sha = picture_sha(shape)
    return f"figure:{sha}" if sha else ""


def is_structural_slide_text(text: str) -> bool:
    value = normalize_text(text)
    if not value:
        return True
    if value.startswith("目录") or "CONTENTS" in value:
        return True
    if "水土保持方案报告表" in value and "技术汇报" in value:
        return True
    if "章节内容保留报告中的源表" in value and "后续页面以该章节证据为核心" in value:
        return True
    if re.search(r"第\s*\d+\s*章\s*/\s*共\s*\d+\s*章", value):
        return True
    return False


def review_content(prs: Presentation, round_idx: int) -> list[ReviewIssue]:
    issues: list[ReviewIssue] = []
    sanitizer = SlideContentSanitizer()
    for slide_index, slide in enumerate(prs.slides, start=1):
        for shape_index, shape in enumerate(iter_shapes(slide), start=1):
            if not getattr(shape, "has_text_frame", False):
                continue
            location = shape_location(shape, f"text_box_{shape_index}")
            text = normalize_text(getattr(shape, "text", "") or "")
            if not text or is_meta_text(text):
                continue
            violations = sanitizer.visible_text_violations(text)
            if violations:
                issue(issues, round_idx, slide_index, "internal_note_leak", "critical", location, f"正文包含后台/生成过程文字：{', '.join(violations[:4])}", True, "删除、改写或移入 internal_notes。", "ContentReview")
            if PLACEHOLDER_RE.search(text):
                issue(issues, round_idx, slide_index, "empty_placeholder", "critical", location, f"正文包含未替换占位符：{compact(text)}", True, "删除占位符或用源报告内容替换。", "ContentReview")
            paragraphs = [normalize_text(p.text) for p in shape.text_frame.paragraphs if normalize_text(p.text)]
            body_paragraphs = [p for p in paragraphs if not is_meta_text(p)]
            if len(body_paragraphs) > 5:
                issue(issues, round_idx, slide_index, "content_too_long", "high", location, "单页正文要点超过 5 条。", True, "压缩到 5 条以内，次要内容移入备注。", "ContentReview")
            for para_index, paragraph in enumerate(body_paragraphs, start=1):
                if len(paragraph) > 58:
                    issue(issues, round_idx, slide_index, "content_too_long", "high", f"{location}/p{para_index}", f"要点过长：{compact(paragraph)}", True, "压缩为 35-45 个中文字符左右。", "ContentReview")
                if para_index > 1 and paragraph == body_paragraphs[para_index - 2]:
                    issue(issues, round_idx, slide_index, "duplicate_content", "high", f"{location}/p{para_index}", f"重复表达：{compact(paragraph)}", True, "删除重复要点。", "ContentReview")
    return issues


def review_duplicates(prs: Presentation, round_idx: int, policy: dict[str, Any]) -> list[ReviewIssue]:
    issues: list[ReviewIssue] = []
    duplicate_policy = policy.get("duplicate_content", {})
    if duplicate_policy.get("enabled", True) is False:
        return issues
    check_slides = duplicate_policy.get("check_duplicate_slides", True) is not False
    check_tables = duplicate_policy.get("check_duplicate_tables", True) is not False
    check_figures = duplicate_policy.get("check_duplicate_figures", True) is not False
    min_text_chars = int(duplicate_policy.get("minimum_body_characters", 80))
    paragraph_threshold = float(duplicate_policy.get("paragraph_similarity", 0.92))
    slide_threshold = float(duplicate_policy.get("max_slide_similarity", 0.94))
    slide_visual_threshold = float(duplicate_policy.get("max_slide_similarity_with_shared_visual", 0.88))
    previous_slides: list[dict[str, Any]] = []
    seen_tables: dict[str, tuple[int, str]] = {}
    seen_figures: dict[str, tuple[int, str]] = {}

    for slide_index, slide in enumerate(prs.slides, start=1):
        body_text = collect_body_text(slide)
        normalized = canonical_duplicate_text(body_text)
        structural = is_structural_slide_text(body_text)
        table_keys: list[tuple[str, str]] = []
        figure_keys: list[tuple[str, str]] = []
        paragraphs: list[tuple[str, str]] = []

        for shape_index, shape in enumerate(iter_shapes(slide), start=1):
            location = shape_location(shape, f"shape_{shape_index}")
            if getattr(shape, "has_text_frame", False):
                for para_index, paragraph in enumerate(shape.text_frame.paragraphs, start=1):
                    text = normalize_text(paragraph.text)
                    key = canonical_duplicate_text(text)
                    if len(key) >= 18 and not is_meta_text(text):
                        paragraphs.append((f"{location}/p{para_index}", key))
            table_key = table_fingerprint(shape)
            if table_key and len(table_key) > 18:
                table_keys.append((location, table_key))
            figure_key = figure_fingerprint(shape)
            if figure_key and shape_area_ratio(prs, shape) >= 0.025:
                figure_keys.append((location, figure_key))

        for idx, (location, key) in enumerate(paragraphs):
            for previous_location, previous_key in paragraphs[:idx]:
                if text_similarity(previous_key, key) >= paragraph_threshold:
                    issue(
                        issues,
                        round_idx,
                        slide_index,
                        "duplicate_content",
                        "high",
                        location,
                        "同页存在高度重复的正文表述。",
                        True,
                        "删除后一条重复表述或合并为一个要点。",
                        "ContentReview",
                    )
                    break

        if check_slides and not structural and len(normalized) >= min_text_chars:
            for previous in previous_slides:
                ratio = text_similarity(previous["text"], normalized)
                shared_visuals = bool(set(previous["tables"]).intersection(key for _loc, key in table_keys) or set(previous["figures"]).intersection(key for _loc, key in figure_keys))
                if ratio >= slide_threshold or (ratio >= slide_visual_threshold and shared_visuals):
                    issue(
                        issues,
                        round_idx,
                        slide_index,
                        "duplicate_slide",
                        "high",
                        f"slide|previous:{previous['slide_index']}",
                        f"页面与第 {previous['slide_index']} 页内容高度重复，相似度 {ratio:.2f}。",
                        True,
                        "删除后一页或合并两页内容。",
                        "ContentReview",
                    )
                    break
            previous_slides.append(
                {
                    "slide_index": slide_index,
                    "text": normalized,
                    "tables": [key for _loc, key in table_keys],
                    "figures": [key for _loc, key in figure_keys],
                }
            )

        for location, key in table_keys if check_tables else []:
            previous = seen_tables.get(key)
            if previous:
                severity = "high" if previous[0] == slide_index - 1 else "medium"
                issue(
                    issues,
                    round_idx,
                    slide_index,
                    "duplicate_table",
                    severity,
                    location,
                    f"表格与第 {previous[0]} 页的 {previous[1]} 重复。",
                    severity == "high",
                    "删除后一处重复表格；若确需复用，应在备注中说明用途差异。",
                    "VisualReview",
                )
            else:
                seen_tables[key] = (slide_index, location)

        for location, key in figure_keys if check_figures else []:
            previous = seen_figures.get(key)
            if previous:
                severity = "high" if previous[0] == slide_index - 1 else "medium"
                issue(
                    issues,
                    round_idx,
                    slide_index,
                    "duplicate_figure",
                    severity,
                    location,
                    f"图件与第 {previous[0]} 页的 {previous[1]} 重复。",
                    severity == "high",
                    "删除后一处重复图件；若确需复用，应改为引用说明或合并页面。",
                    "VisualReview",
                )
            else:
                seen_figures[key] = (slide_index, location)
    return issues


def review_format(prs: Presentation, round_idx: int) -> list[ReviewIssue]:
    issues: list[ReviewIssue] = []
    slide_width = int(prs.slide_width)
    slide_height = int(prs.slide_height)
    for slide_index, slide in enumerate(prs.slides, start=1):
        text_boxes: list[tuple[Any, tuple[int, int, int, int], str]] = []
        object_boxes: list[tuple[Any, tuple[int, int, int, int], str]] = []
        for shape_index, shape in enumerate(iter_shapes(slide), start=1):
            location = shape_location(shape, f"shape_{shape_index}")
            bbox = shape_bbox(shape)
            if bbox[0] < 0 or bbox[1] < 0 or bbox[2] > slide_width or bbox[3] > slide_height:
                issue(issues, round_idx, slide_index, "text_overflow" if getattr(shape, "has_text_frame", False) else "layout_unbalanced", "high", location, "对象超出页面边界。", True, "调整尺寸或移入页面安全范围。", "FormatReview")
            if getattr(shape, "has_text_frame", False):
                text = normalize_text(getattr(shape, "text", "") or "")
                if text and not is_meta_text(text):
                    text_boxes.append((shape, bbox, location))
                    for para_index, paragraph in enumerate(shape.text_frame.paragraphs, start=1):
                        para_text = normalize_text(paragraph.text)
                        if not para_text:
                            continue
                        size = paragraph_font_size(paragraph)
                        if size < BODY_MIN_PT and not is_meta_text(para_text):
                            issue(issues, round_idx, slide_index, "font_below_minimum", "high", f"{location}/p{para_index}", f"正文字号 {size:.1f}pt 低于最小值 {BODY_MIN_PT:.0f}pt。", True, "恢复到最小字号并压缩或拆页。", "FormatReview")
                        if estimated_lines(para_text, max(1, bbox[2] - bbox[0]), max(size, BODY_MIN_PT)) > available_lines(max(1, bbox[3] - bbox[1]), max(size, BODY_MIN_PT)):
                            issue(issues, round_idx, slide_index, "text_overflow", "high", f"{location}/p{para_index}", f"文字可能超出文本框：{compact(para_text)}", True, "压缩文字、扩大文本框或拆页。", "FormatReview")
            if getattr(shape, "has_table", False):
                object_boxes.append((shape, bbox, location))
                for row_index, row in enumerate(shape.table.rows, start=1):
                    for col_index, cell in enumerate(row.cells, start=1):
                        for paragraph in cell.text_frame.paragraphs:
                            if paragraph.text.strip() and paragraph_font_size(paragraph, TABLE_MIN_PT) < TABLE_MIN_PT:
                                issue(issues, round_idx, slide_index, "font_below_minimum", "high", f"{location}/r{row_index}c{col_index}", "表格单元格字号低于最小值。", True, "调整布局或删除低可读表格。", "FormatReview")
            if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
                object_boxes.append((shape, bbox, location))
                if "TableImageNeedsSplit" in location:
                    issue(issues, round_idx, slide_index, "low_readability", "high", location, "表格截图低于可读字号或需要跨页展示。", True, "直接删除该表格，不放入 PPT。", "FormatReview")
                try:
                    image_w, image_h = shape.image.size
                    image_ratio = image_w / max(1, image_h)
                    shape_ratio = max(1, bbox[2] - bbox[0]) / max(1, bbox[3] - bbox[1])
                    if abs(math.log(image_ratio / shape_ratio)) > 0.12:
                        issue(issues, round_idx, slide_index, "image_distortion", "high", location, "图片或表格截图宽高比被改变。", True, "按原始宽高比重新缩放。", "FormatReview")
                except Exception:
                    pass
        for left_index, (left_shape, left_box, left_location) in enumerate(text_boxes):
            for right_shape, right_box, right_location in text_boxes[left_index + 1 :]:
                smaller = min(box_area(left_box), box_area(right_box))
                if smaller and intersection_area(left_box, right_box) / smaller > 0.35:
                    issue(issues, round_idx, slide_index, "text_overlap", "high", f"{left_location}|{right_location}", "两个正文文本框明显重叠。", True, "压缩内容或移动文本框。", "FormatReview")
        for text_shape, text_box, text_location in text_boxes:
            text_area = max(1, box_area(text_box))
            for object_shape, object_box, object_location in object_boxes:
                if text_shape is object_shape:
                    continue
                if intersection_area(text_box, object_box) / text_area > 0.15:
                    issue(issues, round_idx, slide_index, "text_overlap", "high", f"{text_location}|{object_location}", "正文压住图片或表格区域。", True, "删除次要文字、移动文本框或删除低可读表格。", "FormatReview")
    return issues


def estimated_lines(text: str, width_emu: int, font_size: float) -> int:
    width_pt = width_emu / 12700
    units = sum(0.55 if ord(ch) < 128 else 1.0 for ch in text)
    chars_per_line = max(4.0, width_pt / max(font_size * 0.92, 1.0))
    return max(1, math.ceil(units / chars_per_line))


def available_lines(height_emu: int, font_size: float) -> float:
    height_pt = height_emu / 12700
    return max(1.0, height_pt / max(font_size * 1.24, 1.0))


def review_fact_consistency(project: Path, prs: Presentation, round_idx: int) -> list[ReviewIssue]:
    issues: list[ReviewIssue] = []
    numbers = source_numbers(project)
    if not numbers:
        return issues
    for slide_index, slide in enumerate(prs.slides, start=1):
        for shape_index, shape in enumerate(iter_shapes(slide), start=1):
            if not getattr(shape, "has_text_frame", False):
                continue
            location = shape_location(shape, f"text_box_{shape_index}")
            for para_index, paragraph in enumerate(shape.text_frame.paragraphs, start=1):
                text = normalize_text(paragraph.text)
                if not text or is_meta_text(text):
                    continue
                unsupported = [
                    value
                    for value in {normalize_number(item) for item in NUMBER_RE.findall(text)}
                    if value not in numbers and value not in {str(slide_index), "1", "2", "3", "4", "5"}
                ]
                if unsupported:
                    issue(issues, round_idx, slide_index, "factual_inconsistency", "medium", f"{location}/p{para_index}", f"数值未能在源报告中直接追溯：{', '.join(unsupported[:5])}", True, "删除或改写该段，无法确认的数值移入备注。", "FactConsistencyReview")
    return issues


def review_visual(prs: Presentation, round_idx: int) -> list[ReviewIssue]:
    issues: list[ReviewIssue] = []
    slide_width = int(prs.slide_width)
    slide_height = int(prs.slide_height)
    slide_area = max(1, slide_width * slide_height)
    for slide_index, slide in enumerate(prs.slides, start=1):
        boxes = [shape_bbox(shape) for shape in iter_shapes(slide)]
        occupied = sum(box_area(box) for box in boxes) / slide_area
        if occupied > 1.45:
            issue(issues, round_idx, slide_index, "layout_unbalanced", "medium", "slide", "页面元素密度偏高，可能显得拥挤。", False, "人工复核或重新排版。", "VisualReview")
        if occupied < 0.08 and len(boxes) > 1:
            issue(issues, round_idx, slide_index, "layout_unbalanced", "low", "slide", "页面可见内容偏少，可能留白过大。", False, "可补充源报告内容或保留为过渡页。", "VisualReview")
    return issues


def run_release_audit(project: Path, pptx: Path) -> dict[str, Any]:
    script = Path(__file__).resolve().with_name("release_audit.py")
    subprocess.run(
        [sys.executable, str(script), str(project), "--strict", "--pptx", str(pptx)],
        check=False,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        text=True,
    )
    audit_path = project / "qa" / "release_audit.json"
    if audit_path.exists():
        return json.loads(audit_path.read_text(encoding="utf-8-sig"))
    return {"summary": {}, "findings": []}


def issues_from_release_audit(audit: dict[str, Any], round_idx: int) -> list[ReviewIssue]:
    mapping = {
        "forbidden": ("internal_note_leak", "critical"),
        "internal": ("internal_note_leak", "critical"),
        "placeholder": ("empty_placeholder", "critical"),
        "overflow": ("text_overflow", "high"),
        "overlap": ("text_overlap", "high"),
        "font-too-small": ("font_below_minimum", "high"),
        "image-table-font-too-small": ("low_readability", "high"),
        "missing-emphasis": ("low_readability", "high"),
        "duplicate": ("duplicate_content", "high"),
        "sparse": ("layout_unbalanced", "medium"),
        "number-not-in-source": ("factual_inconsistency", "high"),
        "source-unmatched": ("source_untraceable", "medium"),
    }
    issues: list[ReviewIssue] = []
    for item in audit.get("findings", []):
        if item.get("level") not in {"error", "warning"}:
            continue
        code = str(item.get("code", "audit"))
        issue_type = "layout_unbalanced"
        severity = "medium" if item.get("level") == "warning" else "high"
        for key, mapped in mapping.items():
            if key in code:
                issue_type, severity = mapped
                break
        context = item.get("context", {}) or {}
        slide_index = int(context.get("page") or context.get("slide") or 0)
        issue(issues, round_idx, slide_index, issue_type, severity, str(context.get("text") or context.get("object") or code), str(item.get("message", code)), severity in {"critical", "high"}, "按 issue 类型自动修复或删除问题内容。", "ReleaseAudit")
    return issues


def review_ppt(project: Path, pptx: Path, round_idx: int) -> list[ReviewIssue]:
    prs = Presentation(pptx)
    policy = load_policy(project)
    issues = []
    issues.extend(review_content(prs, round_idx))
    issues.extend(review_duplicates(prs, round_idx, policy))
    issues.extend(review_format(prs, round_idx))
    issues.extend(review_fact_consistency(project, prs, round_idx))
    issues.extend(review_visual(prs, round_idx))
    issues.extend(issues_from_release_audit(run_release_audit(project, pptx), round_idx))
    return dedupe_issues(issues)


def dedupe_issues(issues: list[ReviewIssue]) -> list[ReviewIssue]:
    seen: set[tuple[Any, ...]] = set()
    result: list[ReviewIssue] = []
    for item in issues:
        key = (item.slide_index, item.issue_type, item.location, item.description[:80])
        if key in seen:
            continue
        seen.add(key)
        item.issue_id = f"S{item.slide_index:03d}-{item.issue_type.upper().replace('_', '-')}-{len(result)+1:03d}"
        result.append(item)
    return result


def sanitize_text_value(text: str) -> str:
    sanitizer = SlideContentSanitizer()
    value = sanitizer.sanitize(text, component="body_text")
    value = PLACEHOLDER_RE.sub("", value)
    return normalize_text(value)


def repair_text_frame(shape, max_paragraphs: int = 5) -> bool:
    if not getattr(shape, "has_text_frame", False):
        return False
    paragraphs = [sanitize_text_value(paragraph.text) for paragraph in shape.text_frame.paragraphs]
    paragraphs = [item for item in paragraphs if item and not PLACEHOLDER_RE.search(item)]
    if not paragraphs:
        shape.text_frame.clear()
        return True
    clean: list[str] = []
    seen_keys: list[str] = []
    for item in paragraphs[:max_paragraphs]:
        key = canonical_duplicate_text(item)
        if key and any(text_similarity(key, previous) >= 0.92 for previous in seen_keys):
            continue
        seen_keys.append(key)
        clean.append(compress_for_ppt(item, max_chars=45))
    if not clean:
        shape.text_frame.clear()
        return True
    original = normalize_text(getattr(shape, "text", "") or "")
    if normalize_text("\n".join(clean)) == original:
        return False
    shape.text_frame.clear()
    paragraph = shape.text_frame.paragraphs[0]
    paragraph.text = clean[0]
    paragraph.font.size = Pt(max(paragraph_font_size(paragraph), BODY_MIN_PT))
    for item in clean[1:]:
        new_paragraph = shape.text_frame.add_paragraph()
        new_paragraph.text = item
        new_paragraph.font.size = Pt(BODY_MIN_PT)
    return True


def delete_shape(shape) -> bool:
    try:
        element = shape._element
        element.getparent().remove(element)
        return True
    except Exception:
        return False


def delete_slide(prs: Presentation, slide_index: int) -> bool:
    if slide_index < 1 or slide_index > len(prs.slides):
        return False
    try:
        slides = prs.slides
        slide_id = slides._sldIdLst[slide_index - 1]
        rel_id = slide_id.rId
        slides._sldIdLst.remove(slide_id)
        prs.part.drop_rel(rel_id)
        return True
    except Exception:
        return False


def restore_picture_aspect(shape) -> bool:
    if getattr(shape, "shape_type", None) != MSO_SHAPE_TYPE.PICTURE:
        return False
    try:
        image_w, image_h = shape.image.size
    except Exception:
        return False
    aspect = image_w / max(1, image_h)
    box_w = int(shape.width)
    box_h = int(shape.height)
    if box_w <= 0 or box_h <= 0:
        return False
    box_aspect = box_w / box_h
    center_x = int(shape.left) + box_w // 2
    center_y = int(shape.top) + box_h // 2
    if box_aspect > aspect:
        new_w = int(box_h * aspect)
        new_h = box_h
    else:
        new_w = box_w
        new_h = int(box_w / aspect)
    shape.width = new_w
    shape.height = new_h
    shape.left = center_x - new_w // 2
    shape.top = center_y - new_h // 2
    return True


def is_body_candidate_for_emphasis(shape) -> bool:
    if not getattr(shape, "has_text_frame", False):
        return False
    name = shape_location(shape, "").lower()
    if any(marker in name for marker in ("title", "page_meta", "caption", "table_note")):
        return False
    text = normalize_text(getattr(shape, "text", "") or "")
    if len(text) < 20 or META_TEXT_RE.search(text):
        return False
    if int(getattr(shape, "top", 0) or 0) < int(1.05 * EMU_PER_INCH):
        return False
    return bool(NUMBER_RE.search(text) or any(term in text for term in EMPHASIS_TERMS))


def paragraph_text_has_body_emphasis(paragraph) -> bool:
    for run in paragraph.runs:
        if not (run.text or "").strip():
            continue
        if getattr(run.font, "bold", False):
            return True
        try:
            if run.font.color.rgb is not None and str(run.font.color.rgb).upper() != "222222":
                return True
        except Exception:
            pass
    return False


def add_body_emphasis(shape) -> bool:
    if not is_body_candidate_for_emphasis(shape):
        return False
    if any(paragraph_text_has_body_emphasis(paragraph) for paragraph in shape.text_frame.paragraphs):
        return False
    pattern = re.compile("|".join(re.escape(term) for term in EMPHASIS_TERMS) + r"|[-+]?\d+(?:\.\d+)?%?")
    for paragraph in shape.text_frame.paragraphs:
        text = paragraph.text or ""
        match = pattern.search(text)
        if not match:
            continue
        before, focus, after = text[: match.start()], match.group(0), text[match.end() :]
        try:
            paragraph.clear()
            if before:
                paragraph.add_run().text = before
            run = paragraph.add_run()
            run.text = focus
            run.font.bold = True
            run.font.color.rgb = EMPHASIS_RGB
            if after:
                paragraph.add_run().text = after
            return True
        except Exception:
            for run in paragraph.runs:
                if focus in (run.text or ""):
                    run.font.bold = True
                    run.font.color.rgb = EMPHASIS_RGB
                    return True
    return False


def repair_ppt(project: Path, pptx: Path, issues: list[ReviewIssue], round_idx: int) -> tuple[Path, list[ReviewIssue]]:
    prs = Presentation(pptx)
    repaired = deepcopy(issues)
    changed = False

    duplicate_slide_issues = [
        item for item in repaired if item.auto_fixable and item.issue_type == "duplicate_slide"
    ]
    if duplicate_slide_issues:
        for item in sorted(duplicate_slide_issues, key=lambda issue_item: issue_item.slide_index, reverse=True):
            if delete_slide(prs, item.slide_index):
                changed = True
                item.status = "fixed"
                item.repair_action = "deleted duplicate slide"
        if changed:
            repaired_path = pptx.with_name(f"{pptx.stem}_review_round{round_idx}{pptx.suffix}")
            prs.save(repaired_path)
            return repaired_path, repaired

    by_slide: dict[int, list[ReviewIssue]] = {}
    for item in repaired:
        by_slide.setdefault(item.slide_index, []).append(item)

    for slide_index, slide in enumerate(prs.slides, start=1):
        slide_issues = by_slide.get(slide_index, [])
        if not slide_issues:
            continue
        issue_types = {item.issue_type for item in slide_issues if item.auto_fixable}
        needs_body_emphasis = any(
            item.auto_fixable and item.location == "pptx-missing-emphasis"
            for item in slide_issues
        )
        for shape in list(iter_shapes(slide)):
            name = shape_location(shape, "")
            duplicate_object_issue = next(
                (
                    item
                    for item in slide_issues
                    if item.auto_fixable
                    and item.issue_type in {"duplicate_table", "duplicate_figure"}
                    and item.location == name
                ),
                None,
            )
            if duplicate_object_issue and delete_shape(shape):
                changed = True
                duplicate_object_issue.status = "fixed"
                duplicate_object_issue.repair_action = f"deleted repeated {duplicate_object_issue.issue_type.replace('duplicate_', '')}"
                continue
            if "low_readability" in issue_types and getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE and ("TableImage" in name or "Table" in name):
                if delete_shape(shape):
                    changed = True
                    for item in slide_issues:
                        if item.issue_type == "low_readability":
                            item.status = "fixed"
                            item.repair_action = "deleted low-readability table image"
                continue
            if "image_distortion" in issue_types and restore_picture_aspect(shape):
                changed = True
            if getattr(shape, "has_text_frame", False) and issue_types.intersection({"internal_note_leak", "empty_placeholder", "content_too_long", "text_overflow", "font_below_minimum", "factual_inconsistency", "duplicate_content", "text_overlap"}):
                if repair_text_frame(shape):
                    changed = True
            if needs_body_emphasis and add_body_emphasis(shape):
                changed = True
                for item in slide_issues:
                    if item.location == "pptx-missing-emphasis":
                        item.status = "fixed"
                        item.repair_action = "added emphasis to body content only"

        for item in slide_issues:
            if item.status == "pending" and item.auto_fixable and changed:
                item.status = "fixed"
                item.repair_action = item.repair_action or "sanitized/compressed affected slide content"

    if not changed:
        return pptx, repaired
    repaired_path = pptx.with_name(f"{pptx.stem}_review_round{round_idx}{pptx.suffix}")
    prs.save(repaired_path)
    return repaired_path, repaired


def write_issues(project: Path, round_idx: int, issues: list[ReviewIssue]) -> None:
    qa = project / "qa"
    qa.mkdir(parents=True, exist_ok=True)
    (qa / f"issue_list_round_{round_idx}.json").write_text(
        json.dumps([asdict(item) for item in issues], ensure_ascii=False, indent=2) + "\n",
        encoding="utf-8",
    )


def write_report(project: Path, final_pptx: Path, rounds: list[dict[str, Any]], final_issues: list[ReviewIssue]) -> Path:
    qa = project / "qa"
    remaining = [item for item in final_issues if item.status != "fixed"]
    blocking = [item for item in remaining if item.severity in {"critical", "high"}]
    repaired_issue_records = [
        issue_item
        for round_item in rounds
        for issue_item in round_item.get("repaired_issues", [])
        if issue_item.get("status") == "fixed"
    ]
    report = {
        "pipeline": "PostGenerationReviewPipeline",
        "created_at": datetime.now().isoformat(timespec="seconds"),
        "final_pptx": str(final_pptx),
        "total_slides": len(Presentation(final_pptx).slides),
        "review_rounds": len(rounds),
        "issues_found": sum(len(round_item["issues"]) for round_item in rounds),
        "issues_fixed": len(repaired_issue_records),
        "issues_remaining": len(remaining),
        "remaining_blocking": len(blocking),
        "remaining_issues": [asdict(item) for item in remaining],
        "fixed_issues": repaired_issue_records,
        "rounds": rounds,
        "final_status": "passed" if not blocking else "blocked",
    }
    path = qa / "review_report.json"
    path.write_text(json.dumps(report, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    return path


def run_pipeline(project: Path, pptx: Path) -> dict[str, Any]:
    current = pptx.resolve()
    qa = project / "qa"
    qa.mkdir(parents=True, exist_ok=True)
    for stale_issue_list in qa.glob("issue_list_round_*.json"):
        try:
            stale_issue_list.unlink()
        except OSError:
            pass
    rounds: list[dict[str, Any]] = []
    final_issues: list[ReviewIssue] = []
    last_reviewed_pptx: Path | None = None
    for round_idx in range(1, MAX_REVIEW_REPAIR_ROUNDS + 1):
        issues = review_ppt(project, current, round_idx)
        last_reviewed_pptx = current
        write_issues(project, round_idx, issues)
        blocking = [item for item in issues if item.severity in {"critical", "high"}]
        round_record = {"round": round_idx, "pptx": str(current), "issues": [asdict(item) for item in issues]}
        rounds.append(round_record)
        final_issues = issues
        if not blocking:
            break
        repairable = [item for item in blocking if item.auto_fixable]
        if not repairable:
            break
        repaired_path, repaired_issues = repair_ppt(project, current, repairable, round_idx)
        round_record["repair_attempted"] = True
        round_record["repaired_pptx"] = str(repaired_path)
        round_record["repaired_issues"] = [asdict(item) for item in repaired_issues]
        if repaired_path == current:
            break
        current = repaired_path
    if last_reviewed_pptx != current:
        final_issues = review_ppt(project, current, len(rounds) + 1)
        write_issues(project, len(rounds) + 1, final_issues)
    report_path = write_report(project, current, rounds, final_issues)
    remaining_blocking = [item for item in final_issues if item.severity in {"critical", "high"}]
    remaining_issues = [item for item in final_issues if item.status != "fixed"]
    return {
        "passed": not remaining_blocking,
        "final_pptx": str(current),
        "review_report": str(report_path),
        "review_rounds": len(rounds),
        "remaining_blocking": len(remaining_blocking),
        "issues_remaining": len(remaining_issues),
    }


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("project", type=Path)
    parser.add_argument("--pptx", type=Path, required=True)
    parser.add_argument("--no-fail", action="store_true")
    args = parser.parse_args()
    result = run_pipeline(args.project.resolve(), args.pptx.resolve())
    print(json.dumps(result, ensure_ascii=False))
    if not result["passed"] and not args.no_fail:
        sys.exit(1)


if __name__ == "__main__":
    main()
