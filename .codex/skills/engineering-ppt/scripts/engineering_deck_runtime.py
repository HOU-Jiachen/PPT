#!/usr/bin/env python3
"""Reusable helpers for source-faithful engineering review PPTX generation.

Project scripts should orchestrate project-specific page choices. Shared behavior
such as template profiling, visible-text hygiene, source evidence extraction, and
basic PPTX drawing belongs here so each project does not grow a one-off runtime.
"""

from __future__ import annotations

import json
import re
import zipfile
from collections import Counter
from datetime import datetime
from pathlib import Path
from typing import Any

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

from figure_layout_optimizer import choose_figure_layout


DEFAULT_COLORS = {
    "bg": "F7F9FC",
    "paper": "FFFFFF",
    "ink": "01203C",
    "muted": "566B7F",
    "primary": "01203C",
    "secondary": "EAF2F8",
    "accent": "4489C8",
    "line": "CCD8E2",
    "soft": "F1F6FA",
    "good": "2E7D32",
    "warn": "A33A2B",
}

BODY_MIN_PT = 14.0
TABLE_MIN_PT = 12.0

IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tif", ".tiff", ".webp"}
EMPHASIS_TERM_RE = re.compile(
    r"[-+]?\d+(?:\.\d+)?\s*(?:%|万?m3|万?m³|km2|km²|hm2|hm²|m|km|万元|亿元|t|d|年|月|日)?"
    r"|防治责任范围|水土流失|表土剥离|临时防护|排水沟|沉沙池|投资估算|监测时段|验收管理|风险|结论|措施|责任|目标"
)

INTERNAL_VISIBLE_PATTERNS = [
    r"来源模式[:：]?\s*\S+",
    r"证据编号[:：]?\s*\S+",
    r"本页用于[:：]?.*",
    r"(?:^|[。；;]\s*)[-•]?\s*(?:原始对象|必讲内容|保留理由|证据)[:：][^。；;]*",
    r"原表定位[:：]?.*",
    r"行列规模[:：]?.*",
    r"密集表按重点行重排.*",
    r"完整数据回看报告原表.*",
    r"\bE-\d(?:-[A-Z0-9]+)+\b",
    r"\bimage_\d+\.(?:png|jpe?g|wmf|emf)\b",
]

GENERIC_PANEL_HEADINGS = [
    "报告对图件的说明",
    "报告对表格的说明",
    "报告计算口径",
    "报告阐述",
]


def rgb(hex_color: str) -> RGBColor:
    h = hex_color.strip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def pt(value: float):
    return Pt(value)


def load_json(path: Path) -> Any:
    return json.loads(path.read_text(encoding="utf-8"))


def sanitize_visible_text(text: str) -> str:
    text = re.sub(r"\s+", " ", text or "").strip()
    text = re.sub(r"\s*\|\s*", "；", text)
    text = text.replace("……", "").replace("...", "").replace("…", "")
    for pattern in INTERNAL_VISIBLE_PATTERNS:
        text = re.sub(pattern, "", text, flags=re.I)
    text = text.replace("证据解读", "报告说明")
    return re.sub(r"\s+", " ", text).strip(" ，,;；")


def ensure_sentence_end(text: str) -> str:
    text = (text or "").strip(" ，,;；")
    if text.endswith(("：", ":")):
        return text[:-1].rstrip(" ，,;；") + "。"
    if text and text[-1] not in "。；;！!？?：:）】》”’0123456789%":
        text += "。"
    return text


def remove_orphan_list_fragment(text: str) -> str:
    """Drop a partial numbered-list tail when only item (1) is visible."""

    value = text or ""
    fullwidth_1 = chr(0xFF08) + "1" + chr(0xFF09)
    fullwidth_2 = chr(0xFF08) + "2" + chr(0xFF09)
    marker_positions = [pos for pos in [value.find(fullwidth_1), value.find("(1)")] if pos >= 0]
    if not marker_positions or fullwidth_2 in value or "(2)" in value:
        return text
    marker_position = min(marker_positions)
    if marker_position < 20:
        return text
    return text[:marker_position].strip(" ，,;；。")


def text_preview(text: str, max_len: int = 480, complete_sentence: bool = False) -> str:
    text = sanitize_visible_text(text)
    if complete_sentence:
        text = remove_orphan_list_fragment(text)
    if len(text) <= max_len:
        return ensure_sentence_end(text) if complete_sentence else text
    clipped = text[:max_len]
    sentence_boundary = max(clipped.rfind(mark) for mark in ("。", "；", ";", "！", "？"))
    phrase_boundary = max(clipped.rfind(mark) for mark in ("，", ",", "、", " "))
    boundary = sentence_boundary if sentence_boundary >= max_len * 0.45 else phrase_boundary
    if boundary >= max_len * 0.58:
        clipped = clipped[: boundary + (1 if sentence_boundary == boundary else 0)]
    clipped = clipped.strip(" ，,;；")
    if complete_sentence:
        clipped = remove_orphan_list_fragment(clipped)
    return ensure_sentence_end(clipped) if complete_sentence else clipped


def complete_semantic_point(text: str, max_len: int = 130) -> str:
    """Return a compact point that does not end on a dangling clause."""

    point = text_preview(text, max_len, complete_sentence=True)
    weak_tail = (
        "并",
        "及",
        "与",
        "和",
        "或",
        "对",
        "将",
        "为",
        "按",
        "由",
        "在",
        "其中",
        "包括",
        "主要",
        "分别",
        "以及",
        "同时",
        "通过",
        "依据",
    )
    stripped = point.rstrip("。；;，,、 ")
    if any(stripped.endswith(tail) for tail in weak_tail):
        for mark in ("；", ";", "，", ",", "、"):
            boundary = stripped.rfind(mark)
            if boundary >= max(12, len(stripped) * 0.45):
                return ensure_sentence_end(stripped[:boundary].strip())
    return point


def catalog_locator_number(entry: dict[str, Any]) -> int | None:
    locator = str(entry.get("locator") or "")
    match = re.search(r"P(\d{4})", locator)
    if match:
        return int(match.group(1))
    entry_id = str(entry.get("id") or "")
    match = re.search(r":P(\d{4})", entry_id)
    return int(match.group(1)) if match else None


def heading_level(entry: dict[str, Any]) -> int:
    style = str(entry.get("style") or "")
    match = re.search(r"heading\s*(\d+)", style, flags=re.I)
    return int(match.group(1)) if match else 6


def catalog_title_context(entries: list[dict[str, Any]]) -> dict[str, dict[str, str]]:
    """Map catalog IDs to report-native section and figure/table captions."""

    context: dict[str, dict[str, str]] = {}
    heading_stack: dict[int, str] = {}
    last_caption = ""
    last_caption_kind = ""

    for entry in entries:
        kind = str(entry.get("kind") or "")
        text = sanitize_visible_text(entry.get("text", ""))
        if not text:
            continue

        if kind == "heading":
            level = heading_level(entry)
            heading_stack = {k: v for k, v in heading_stack.items() if k < level}
            heading_stack[level] = text
            context[entry["id"]] = {"heading": text, "section": text}
            last_caption = ""
            last_caption_kind = ""
            continue

        style = str(entry.get("style") or "")
        caption_kind = ""
        if style == "表名" or text.startswith("表"):
            caption_kind = "table"
        elif style == "图名" or text.startswith("图"):
            caption_kind = "figure"

        deepest = heading_stack[max(heading_stack)] if heading_stack else ""
        if caption_kind:
            last_caption = strip_figure_table_label(text)
            last_caption_kind = caption_kind
            context[entry["id"]] = {"heading": deepest, "caption": last_caption, "caption_kind": caption_kind}
            continue

        title = {"heading": deepest}
        if kind == "table" and last_caption_kind == "table":
            title["caption"] = last_caption
            title["caption_kind"] = "table"
        elif kind in {"image", "figure"} and last_caption_kind == "figure":
            title["caption"] = last_caption
            title["caption_kind"] = "figure"
        context[entry["id"]] = title
        if kind not in {"table", "image", "figure"}:
            last_caption = ""
            last_caption_kind = ""

    return context


def strip_figure_table_label(text: str) -> str:
    text = sanitize_visible_text(text)
    text = re.sub(r"^[图表]\s*\d+(?:\.\d+)?-\d+\s*", "", text)
    return text.strip(" ：:，,。") or sanitize_visible_text(text)


def locator_range_entries(
    page: dict[str, Any],
    entries: list[dict[str, Any]],
    include_trailing_heading_block: bool = True,
) -> list[dict[str, Any]]:
    source_note = str(page.get("source_note") or "")
    ranges: list[tuple[int, int]] = []
    for start, end in re.findall(r"P(\d{4})\s*[-—~至]\s*P?(\d{4})", source_note):
        ranges.append((int(start), int(end)))
    for single in re.findall(r"(?<![-—~至])P(\d{4})(?!\s*[-—~至]\s*P?\d{4})", source_note):
        value = int(single)
        ranges.append((value, value))
    if not ranges:
        return []

    selected: list[dict[str, Any]] = []
    locators = [catalog_locator_number(entry) for entry in entries]
    for start, end in ranges:
        for index, entry in enumerate(entries):
            number = locators[index]
            if number is not None and start <= number <= end:
                selected.append(entry)
        if not include_trailing_heading_block:
            continue
        end_indices = [i for i, number in enumerate(locators) if number == end]
        if not end_indices:
            continue
        last_index = end_indices[-1]
        if str(entries[last_index].get("kind")) != "heading":
            continue
        for entry in entries[last_index + 1 :]:
            if str(entry.get("kind")) == "heading":
                break
            selected.append(entry)

    deduped: list[dict[str, Any]] = []
    seen: set[str] = set()
    for entry in selected:
        entry_id = str(entry.get("id"))
        if entry_id not in seen:
            seen.add(entry_id)
            deduped.append(entry)
    return deduped


def source_topic_heading(
    page: dict[str, Any],
    evidence: list[dict[str, Any]],
    title_context: dict[str, dict[str, str]],
    catalog_entries: list[dict[str, Any]] | None = None,
    max_len: int = 32,
) -> str:
    """Prefer report section/caption titles over agent-generated summaries."""

    mode = str(page.get("source_mode", "")).upper()
    candidates: list[str] = []

    if mode == "ORIGINAL_FIGURE" and catalog_entries:
        figure_ids = extract_figure_ids(" ".join([str(page.get("source_note", "")), str(page.get("visual_proof", ""))]))
        for figure_id in figure_ids:
            for entry in catalog_entries:
                text = sanitize_visible_text(entry.get("text", ""))
                if figure_id in text and (str(entry.get("style")) == "图名" or text.startswith("图")):
                    candidates.append(strip_figure_table_label(text))
                    break

    if mode == "ORIGINAL_TABLE" and catalog_entries:
        table_refs = extract_table_refs(
            " ".join([str(page.get("source_note", "")), str(page.get("visual_proof", "")), str(page.get("title", ""))])
        )
        for table_ref in table_refs:
            for entry in catalog_entries:
                text = sanitize_visible_text(entry.get("text", ""))
                if table_ref in text and (str(entry.get("style")) == "表名" or text.startswith("表")):
                    candidates.append(strip_figure_table_label(text))
                    break

    if catalog_entries and mode in {"ORIGINAL_TEXT", "CALCULATION"}:
        range_entries = locator_range_entries(page, catalog_entries, include_trailing_heading_block=False)
        range_headings = []
        if range_entries:
            leading_heading = title_context.get(str(range_entries[0].get("id")), {}).get("heading", "")
            if leading_heading:
                range_headings.append(leading_heading)
        range_headings.extend(
            sanitize_visible_text(entry.get("text", ""))
            for entry in range_entries
            if str(entry.get("kind")) == "heading"
        )
        if range_headings:
            candidates.append(join_short_titles(range_headings[-3:]))

    for ev in evidence:
        contexts = [title_context.get(catalog_id, {}) for catalog_id in ev.get("catalog_ids", [])]
        if mode == "ORIGINAL_TABLE":
            candidates.extend(ctx.get("caption", "") for ctx in contexts if ctx.get("caption_kind") == "table")
        elif mode == "ORIGINAL_FIGURE":
            candidates.extend(ctx.get("caption", "") for ctx in contexts if ctx.get("caption_kind") == "figure")
        else:
            candidates.extend(ctx.get("heading", "") for ctx in contexts if ctx.get("heading"))

    if catalog_entries and mode == "ORIGINAL_TABLE" and not candidates:
        for entry in locator_range_entries(page, catalog_entries, include_trailing_heading_block=False):
            text = sanitize_visible_text(entry.get("text", ""))
            if str(entry.get("style")) == "表名" or text.startswith("表"):
                candidates.append(strip_figure_table_label(text))

    candidates.append(str(page.get("chapter", "")))
    candidates.append(str(page.get("title", "")))

    for candidate in candidates:
        value = sanitize_visible_text(candidate)
        if value and value not in GENERIC_PANEL_HEADINGS:
            return text_preview(value, max_len)
    return text_preview(report_heading(mode), max_len)


def join_short_titles(titles: list[str], limit: int = 32) -> str:
    clean = [strip_figure_table_label(title) for title in titles if sanitize_visible_text(title)]
    if not clean:
        return ""
    if len(clean) == 1:
        return clean[0]
    joined = "、".join(clean[:-1]) + "与" + clean[-1]
    return text_preview(joined, limit)


def extract_docx_media(source_docx: Path, output_dir: Path) -> list[Path]:
    """Extract embedded DOCX media for source-figure slides."""

    if not source_docx.exists():
        raise FileNotFoundError(f"Source DOCX not found: {source_docx}")
    output_dir.mkdir(parents=True, exist_ok=True)
    extracted: list[Path] = []
    with zipfile.ZipFile(source_docx) as archive:
        for name in archive.namelist():
            path = Path(name)
            if path.parent.as_posix() != "word/media" or path.suffix.lower() not in IMAGE_EXTENSIONS:
                continue
            target = output_dir / path.name
            target.write_bytes(archive.read(name))
            extracted.append(target)
    if not extracted:
        raise ValueError(f"No embedded media found in DOCX: {source_docx}")
    return sorted(extracted, key=lambda item: item.name.lower())


def require_media_paths(media_dir: Path, names: list[str], context: str = "source figure") -> list[Path]:
    """Return existing media paths or fail the build before a placeholder can leak."""

    paths = [media_dir / name for name in names]
    missing = [str(path) for path in paths if not path.exists()]
    if missing:
        joined = "; ".join(missing)
        raise FileNotFoundError(f"Missing required {context} media: {joined}")
    return paths


def split_report_points(text: str, max_points: int = 5, max_len: int = 130) -> list[str]:
    text = sanitize_visible_text(text)
    text = remove_orphan_list_fragment(text)
    if not text:
        return []
    sentences = [item.strip(" ；;。") for item in re.split(r"[。；;]\s*", text) if item.strip()]
    if len(sentences) <= 1:
        sentences = [item.strip() for item in re.split(r"，|,\s*", text) if item.strip()]

    points: list[str] = []
    current = ""
    for sentence in sentences:
        candidate = sentence if not current else f"{current}；{sentence}"
        if len(candidate) <= max_len:
            current = candidate
            continue
        if current:
            points.append(complete_semantic_point(current, max_len))
        current = sentence
        if len(points) >= max_points:
            break
    if current and len(points) < max_points:
        points.append(complete_semantic_point(current, max_len))
    return points[:max_points] or [complete_semantic_point(text, max_len)]


def emphasis_terms_from_text(text: str, max_terms: int = 8) -> list[str]:
    """Pick source-visible numbers and engineering terms for bold/color emphasis."""

    terms: list[str] = []
    for match in EMPHASIS_TERM_RE.findall(sanitize_visible_text(text)):
        term = match.strip()
        if term and term not in terms:
            terms.append(term)
        if len(terms) >= max_terms:
            break
    return terms


def point_title(point: str, fallback: str, max_len: int = 14) -> str:
    value = sanitize_visible_text(point)
    title = re.split(r"[：:，,；;。]", value, maxsplit=1)[0].strip()
    if not title or len(title) < 3:
        title = fallback
    return text_preview(title, max_len)


def structured_points(points: list[str | dict[str, Any]], max_items: int = 5) -> list[dict[str, str]]:
    """Normalize split paragraphs into numbered items with a short visible title."""

    normalized: list[dict[str, str]] = []
    for index, item in enumerate(points[:max_items], start=1):
        if isinstance(item, dict):
            body = sanitize_visible_text(str(item.get("text") or item.get("text_preview") or ""))
            title = sanitize_visible_text(str(item.get("title") or point_title(body, f"要点{index}")))
            number = sanitize_visible_text(str(item.get("number") or f"{index:02d}"))
        else:
            body = sanitize_visible_text(str(item))
            title = point_title(body, f"要点{index}")
            number = f"{index:02d}"
        if body:
            normalized.append(
                {
                    "number": number,
                    "title": text_preview(title, 14),
                    "text": text_preview(body, 120, complete_sentence=True),
                }
            )
    return normalized


def report_heading(source_mode: str) -> str:
    mode = source_mode.upper()
    if mode == "ORIGINAL_FIGURE":
        return "图件要点"
    if mode == "ORIGINAL_TABLE":
        return "表格要点"
    if mode == "CALCULATION":
        return "计算要点"
    if mode in {"CONCLUSION", "RECOMMENDATION", "MANAGEMENT_ACTION"}:
        return "结论与建议"
    return "内容要点"


def concise_topic_heading(page: dict[str, Any], max_len: int = 32) -> str:
    """Fallback report-topic heading without workflow prefixes."""

    mode = str(page.get("source_mode", "")).upper()
    title = sanitize_visible_text(page.get("title", ""))
    title = re.sub(r"^(报告|本页|该页|图件|表格|计算|原文|内容|重点)[：:，,\s]+", "", title)
    title = title.strip(" ：:，,。")
    if not title:
        title = sanitize_visible_text(page.get("chapter", "")) or report_heading(mode)
    return text_preview(title, max_len)


def template_pptx(project_dir: Path) -> Path | None:
    candidates = sorted((project_dir / "template").glob("*.pptx"))
    return candidates[0] if candidates else None


def iter_shapes(shapes):
    for shape in shapes:
        yield shape
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP and hasattr(shape, "shapes"):
            yield from iter_shapes(shape.shapes)


def collect_template_profile(project_dir: Path) -> dict[str, Any]:
    path = template_pptx(project_dir)
    if not path:
        return {"available": False}
    try:
        prs = Presentation(path)
    except Exception as exc:
        return {"available": False, "error": str(exc)}

    fonts: Counter[str] = Counter()
    sizes: Counter[float] = Counter()
    text_colors: Counter[str] = Counter()
    fill_colors: Counter[str] = Counter()
    picture_counts: Counter[int] = Counter()
    table_counts: Counter[int] = Counter()
    layout_keys: Counter[tuple[int, int, int]] = Counter()
    samples: list[dict[str, Any]] = []

    for index, slide in enumerate(prs.slides, start=1):
        pictures = 0
        tables = 0
        texts: list[str] = []
        text_count = 0
        for shape in iter_shapes(slide.shapes):
            if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
                pictures += 1
            if getattr(shape, "has_table", False):
                tables += 1
            try:
                fill = shape.fill
                if fill and fill.type and fill.fore_color.type == 1:
                    fill_colors[str(fill.fore_color.rgb)] += 1
            except Exception:
                pass
            if getattr(shape, "has_text_frame", False):
                value = shape.text.strip()
                if value:
                    text_count += 1
                    texts.append(value.replace("\n", " | ")[:120])
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.name:
                            fonts[run.font.name] += 1
                        if run.font.size:
                            sizes[round(run.font.size.pt, 1)] += 1
                        try:
                            if run.font.color.rgb:
                                text_colors[str(run.font.color.rgb)] += 1
                        except Exception:
                            pass
        picture_counts[pictures] += 1
        table_counts[tables] += 1
        layout_keys[(min(pictures, 4), min(tables, 2), min(text_count, 8))] += 1
        if index <= 12 or index == len(prs.slides):
            samples.append(
                {
                    "slide": index,
                    "pictures": pictures,
                    "tables": tables,
                    "sample_text": texts[:5],
                }
            )

    return {
        "available": True,
        "path": str(path.relative_to(project_dir)),
        "slide_count": len(prs.slides),
        "canvas_inches": [
            round(prs.slide_width / 914400, 3),
            round(prs.slide_height / 914400, 3),
        ],
        "dominant_fonts": fonts.most_common(12),
        "dominant_font_sizes": sizes.most_common(12),
        "dominant_text_colors": text_colors.most_common(12),
        "dominant_fill_colors": fill_colors.most_common(12),
        "picture_count_distribution": picture_counts.most_common(),
        "table_count_distribution": table_counts.most_common(),
        "layout_keys": [
            {"pictures": k[0], "tables": k[1], "text_boxes": k[2], "count": v}
            for k, v in layout_keys.most_common(16)
        ],
        "sample_slides": samples,
        "adopted_rules": [
            "16:9 canvas, deep-blue technical review identity",
            "large section titles and restrained body text",
            "source figures shown without asset filenames or extraction IDs",
            "dense tables are split or simplified into readable key rows",
            "engineering report pages pair original evidence with report-facing explanation",
        ],
    }


def extract_figure_ids(text: str) -> list[str]:
    normalized = (text or "").replace("\\", "")
    figure_ids = re.findall(r"图\s*\d+(?:\.\d+)?-\d+", normalized)
    for prefix, start, end in re.findall(
        r"图\s*(\d+(?:\.\d+)?)-(\d+)\s*(?:至|到|~|—|-)\s*(?:图\s*)?(?:\d+(?:\.\d+)?-)?(\d+)",
        normalized,
    ):
        for number in range(int(start), int(end) + 1):
            figure_ids.append(f"图{prefix}-{number}")
    return list(dict.fromkeys(item.replace(" ", "") for item in figure_ids))


def extract_table_refs(text: str) -> list[str]:
    normalized = (text or "").replace("\\", "")
    refs = re.findall(r"表\s*\d+(?:\.\d+)?-\d+", normalized)
    for prefix, start, end in re.findall(
        r"表\s*(\d+(?:\.\d+)?)-(\d+)\s*(?:至|到|~|—|-)\s*(?:表\s*)?(?:\d+(?:\.\d+)?-)?(\d+)",
        normalized,
    ):
        for number in range(int(start), int(end) + 1):
            refs.append(f"表{prefix}-{number}")
    refs.extend(re.findall(r"\bT\d{3}\b", normalized, flags=re.I))
    return list(dict.fromkeys(item.replace(" ", "") for item in refs))


def figure_image_index(markdown_path: Path) -> dict[str, list[str]]:
    if not markdown_path.exists():
        return {}
    lines = markdown_path.read_text(encoding="utf-8").splitlines()
    index: dict[str, list[str]] = {}
    pending_images: list[str] = []
    image_re = re.compile(r"image_\d+\.(?:png|jpe?g|wmf)", re.I)
    for line in lines:
        names = image_re.findall(line)
        if names:
            pending_images = names
        figure_ids = extract_figure_ids(line)
        if figure_ids and pending_images:
            for figure_id in figure_ids:
                index.setdefault(figure_id, [])
                for name in pending_images:
                    if name not in index[figure_id]:
                        index[figure_id].append(name)
            pending_images = []
    return index


def image_paths_from_text(
    media_dir: Path,
    text: str,
    figure_index: dict[str, list[str]] | None = None,
) -> list[Path]:
    names = re.findall(r"image_\d+\.(?:png|jpe?g|wmf)", text, flags=re.I)
    if figure_index:
        for figure_id in extract_figure_ids(text):
            names.extend(figure_index.get(figure_id, []))
    paths: list[Path] = []
    seen: set[Path] = set()
    for name in names:
        path = media_dir / name
        if path.exists() and path.suffix.lower() != ".wmf" and path not in seen:
            seen.add(path)
            paths.append(path)
    return paths


def source_text_for_refs(
    catalog_entries: list[dict[str, Any]],
    refs: list[str],
    max_len: int = 520,
) -> str:
    parts: list[str] = []
    seen: set[str] = set()
    if not refs:
        return ""
    for index, entry in enumerate(catalog_entries):
        text = sanitize_visible_text(entry.get("text", ""))
        if not text or not any(ref in text for ref in refs):
            continue
        if str(entry.get("kind")) == "table":
            continue
        is_caption = str(entry.get("style")) in {"图名", "表名"} or text.startswith(("图", "表"))
        if not is_caption:
            caption_indices = [
                candidate_index
                for candidate_index, candidate in enumerate(catalog_entries)
                if candidate.get("source") == entry.get("source")
                and str(candidate.get("kind")) != "table"
                and any(ref in sanitize_visible_text(candidate.get("text", "")) for ref in refs)
                and (
                    str(candidate.get("style")) in {"图名", "表名"}
                    or sanitize_visible_text(candidate.get("text", "")).startswith(("图", "表"))
                )
            ]
            if caption_indices and index != caption_indices[0]:
                continue
        for candidate in [strip_figure_table_label(text)]:
            if candidate and candidate not in seen:
                seen.add(candidate)
                parts.append(candidate)

        def usable_neighbor(candidate_entry: dict[str, Any]) -> str:
            value = sanitize_visible_text(candidate_entry.get("text", ""))
            if not value:
                return ""
            if str(candidate_entry.get("kind")) in {"heading", "figure", "image", "table"}:
                return ""
            if value.startswith(("图名", "表名", "图", "表")):
                return ""
            if re.search(r"\bimage_\d+\.(?:png|jpe?g|wmf|emf)\b", value, re.I):
                return ""
            return value

        if len("".join(parts)) < max_len * 0.35:
            for neighbor in reversed(catalog_entries[max(0, index - 5) : index]):
                if neighbor.get("source") != entry.get("source"):
                    continue
                neighbor_text = usable_neighbor(neighbor)
                if not any(ref in neighbor_text for ref in refs):
                    continue
                if neighbor_text and neighbor_text not in seen:
                    seen.add(neighbor_text)
                    parts.append(neighbor_text)
                    if len("".join(parts)) >= max_len * 0.55:
                        break
        for follower in catalog_entries[index + 1 : index + 8]:
            if follower.get("source") != entry.get("source"):
                break
            if str(follower.get("kind")) == "heading":
                break
            follower_text = usable_neighbor(follower)
            if not follower_text:
                continue
            if follower_text not in seen:
                seen.add(follower_text)
                parts.append(follower_text)
            if len("".join(parts)) >= max_len:
                break
        if parts:
            break
    return text_preview(" ".join(parts), max_len, complete_sentence=True)


def source_text_for_note_heading(
    page: dict[str, Any],
    catalog_entries: list[dict[str, Any]],
    max_len: int = 560,
) -> str:
    note = sanitize_visible_text(page.get("source_note", ""))
    title = sanitize_visible_text(page.get("title", ""))
    probes = []
    for value in [note, title, page.get("chapter", "")]:
        clean = sanitize_visible_text(value)
        clean = re.sub(r"(章节|小节|图|表|P\d+|T\d+)", "", clean)
        for token in re.split(r"[、，,；;\s]+", clean):
            token = token.strip()
            if len(token) >= 4 and token not in probes:
                probes.append(token)
    if not probes:
        return ""
    for index, entry in enumerate(catalog_entries):
        if str(entry.get("kind")) != "heading":
            continue
        heading = sanitize_visible_text(entry.get("text", ""))
        if not any(probe in heading or heading in probe for probe in probes):
            continue
        parts = []
        for follower in catalog_entries[index + 1 :]:
            if follower.get("source") != entry.get("source"):
                break
            if str(follower.get("kind")) == "heading":
                if parts:
                    break
                continue
            text = sanitize_visible_text(follower.get("text", ""))
            if text and not text.startswith(("图名", "表名", "图", "表")):
                parts.append(text)
            if len("".join(parts)) >= max_len:
                break
        if parts:
            return text_preview(" ".join(parts), max_len, complete_sentence=True)
    return ""


def evidence_lookup(ledger: dict[str, Any]) -> dict[str, dict[str, Any]]:
    return {item["id"]: item for item in ledger["evidence"]}


def slide_evidence(page: dict[str, Any], ev_by_id: dict[str, dict[str, Any]]) -> list[dict[str, Any]]:
    return [ev_by_id[eid] for eid in page.get("evidence_ids", []) if eid in ev_by_id]


def table_from_catalog(
    catalog_entries: dict[str, dict[str, Any]],
    evidence: list[dict[str, Any]],
) -> dict[str, Any] | None:
    for ev in evidence:
        for catalog_id in ev.get("catalog_ids", []):
            entry = catalog_entries.get(catalog_id)
            if entry and entry.get("kind") == "table" and entry.get("rows"):
                return entry
    return None


def table_from_page(
    page: dict[str, Any],
    catalog_entries_by_id: dict[str, dict[str, Any]],
    catalog_entries: list[dict[str, Any]],
    evidence: list[dict[str, Any]],
) -> dict[str, Any] | None:
    refs = extract_table_refs(
        " ".join(
            [
                str(page.get("source_note", "")),
                str(page.get("visual_proof", "")),
                str(page.get("title", "")),
            ]
        )
    )
    for ref in refs:
        if re.fullmatch(r"T\d{3}", ref, flags=re.I):
            for ev in evidence:
                for catalog_id in ev.get("catalog_ids", []):
                    entry = catalog_entries_by_id.get(catalog_id)
                    if entry and str(entry.get("locator", "")).upper() == ref.upper() and entry.get("kind") == "table":
                        return entry
    for ref in refs:
        for index, entry in enumerate(catalog_entries):
            text = sanitize_visible_text(entry.get("text", ""))
            if ref not in text:
                continue
            if not (text.startswith("表名") or text.startswith("表")):
                continue
            for follower in catalog_entries[index + 1 : index + 5]:
                if follower.get("source") != entry.get("source"):
                    break
                if follower.get("kind") == "table" and follower.get("rows"):
                    return follower
    return table_from_catalog(catalog_entries_by_id, evidence)


def report_text_from_evidence(
    evidence: list[dict[str, Any]],
    catalog_entries: dict[str, dict[str, Any]],
    max_len: int = 620,
) -> str:
    parts: list[str] = []
    seen: set[str] = set()

    def add(value: str | None) -> None:
        value = sanitize_visible_text(value or "")
        if not value or value in seen:
            return
        seen.add(value)
        parts.append(value)

    for ev in evidence:
        add(ev.get("exact_text"))
    for ev in evidence:
        for catalog_id in ev.get("catalog_ids", []):
            entry = catalog_entries.get(catalog_id)
            if entry and entry.get("kind") == "paragraph":
                add(entry.get("text"))
    for ev in evidence:
        add(ev.get("summary"))

    return text_preview(" ".join(parts), max_len, complete_sentence=True)


def report_text_from_page(
    page: dict[str, Any],
    evidence: list[dict[str, Any]],
    catalog_entries_by_id: dict[str, dict[str, Any]],
    catalog_entries: list[dict[str, Any]],
    max_len: int = 760,
) -> str:
    """Build visible report text, supplementing sparse evidence from source_note ranges."""

    parts: list[str] = []
    seen: set[str] = set()

    def add(value: str | None) -> None:
        value = sanitize_visible_text(value or "")
        if not value or value in seen:
            return
        seen.add(value)
        parts.append(value)

    for entry in locator_range_entries(page, catalog_entries):
        kind = str(entry.get("kind") or "")
        if kind in {"heading", "table", "figure", "image"}:
            continue
        text = sanitize_visible_text(entry.get("text", ""))
        if not text or text.startswith(("图", "表")):
            continue
        add(text)
        if len("".join(parts)) >= max_len:
            break

    base = report_text_from_evidence(evidence, catalog_entries_by_id, max_len=max_len)
    add(base)

    if len("".join(parts)) < max_len * 0.45:
        add(source_text_for_note_heading(page, catalog_entries, max_len=max_len))

    return text_preview(" ".join(parts), max_len, complete_sentence=True)


def report_points_from_evidence(
    evidence: list[dict[str, Any]],
    catalog_entries: dict[str, dict[str, Any]],
    max_points: int = 5,
) -> list[str]:
    points: list[str] = []
    seen: set[str] = set()

    for ev in evidence:
        raw = ev.get("exact_text") or ev.get("summary") or ""
        if ev.get("summary") and (
            str(ev.get("kind", "")).upper() in {"ORIGINAL_TABLE", "CALCULATION"}
            or "..." in str(raw)
        ):
            raw = ev.get("summary")
        value = sanitize_visible_text(raw)
        if value and value not in seen:
            seen.add(value)
            if len(value) > 150:
                for point in split_report_points(value, max_points=max_points - len(points), max_len=108):
                    if point and point not in seen:
                        seen.add(point)
                        points.append(point)
                    if len(points) >= max_points:
                        return points
            else:
                points.append(text_preview(value, 112, complete_sentence=True))
        if len(points) >= max_points:
            return points

    for ev in evidence:
        for item in ev.get("values", []) or []:
            basis = sanitize_visible_text(str(item.get("time_basis", "")))
            value = sanitize_visible_text(str(item.get("value", "")))
            unit = sanitize_visible_text(str(item.get("unit", "")))
            point = f"{basis}：{value}{unit}" if basis else f"{value}{unit}"
            if point and point not in seen:
                seen.add(point)
                points.append(text_preview(point, 90, complete_sentence=True))
            if len(points) >= max_points:
                return points

    if not points:
        merged = report_text_from_evidence(evidence, catalog_entries, max_len=680)
        points = split_report_points(merged, max_points=max_points)
    return points[:max_points]


def report_points_from_page(
    page: dict[str, Any],
    evidence: list[dict[str, Any]],
    catalog_entries_by_id: dict[str, dict[str, Any]],
    catalog_entries: list[dict[str, Any]],
    max_points: int = 5,
) -> list[str]:
    points: list[str] = []
    range_entries = locator_range_entries(page, catalog_entries)
    if range_entries:
        range_text = " ".join(
            sanitize_visible_text(entry.get("text", ""))
            for entry in range_entries
            if str(entry.get("kind")) != "heading"
            and not sanitize_visible_text(entry.get("text", "")).startswith(("图", "表"))
        )
        points.extend(split_report_points(range_text, max_points=max_points, max_len=138))
    for point in report_points_from_evidence(evidence, catalog_entries_by_id, max_points=max_points):
        if point and not point_repeats_existing(point, points):
            points.append(point)
        if len(points) >= max_points:
            return points[:max_points]
    if len(points) >= max_points:
        return points
    range_text = report_text_from_page(
        page,
        evidence,
        catalog_entries_by_id,
        catalog_entries,
        max_len=900,
    )
    for point in split_report_points(range_text, max_points=max_points, max_len=138):
        if point and not point_repeats_existing(point, points):
            points.append(point)
        if len(points) >= max_points:
            break
    return points[:max_points]


def point_repeats_existing(point: str, existing: list[str]) -> bool:
    normalized = re.sub(r"\W+", "", point)
    if not normalized:
        return True
    for item in existing:
        other = re.sub(r"\W+", "", item)
        if normalized == other:
            return True
        probe = normalized[:32]
        if len(probe) >= 18 and probe in other:
            return True
        reverse_probe = other[:32]
        if len(reverse_probe) >= 18 and reverse_probe in normalized:
            return True
    return False


def add_textbox(
    slide,
    x,
    y,
    w,
    h,
    text,
    size=18,
    color="ink",
    bold=False,
    align=None,
    font="Microsoft YaHei",
    colors: dict[str, str] | None = None,
):
    palette = colors or DEFAULT_COLORS
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    p = tf.paragraphs[0]
    p.text = sanitize_visible_text(text)
    p.font.name = font
    p.font.size = pt(size)
    p.font.bold = bold
    p.font.color.rgb = rgb(palette[color])
    if align:
        p.alignment = align
    return box


def add_emphasis_textbox(
    slide,
    x,
    y,
    w,
    h,
    text,
    size=18,
    color="ink",
    emphasis_terms: list[str] | None = None,
    emphasis_color="accent",
    bold=False,
    align=None,
    font="Microsoft YaHei",
    colors: dict[str, str] | None = None,
):
    """Add editable text with key terms/numbers highlighted as separate runs."""

    palette = colors or DEFAULT_COLORS
    clean = sanitize_visible_text(text)
    terms = [term for term in (emphasis_terms or emphasis_terms_from_text(clean)) if term]
    terms = sorted(dict.fromkeys(terms), key=len, reverse=True)
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    p = tf.paragraphs[0]
    if align:
        p.alignment = align
    if not clean:
        return box
    if not terms:
        p.text = clean
        p.font.name = font
        p.font.size = pt(size)
        p.font.bold = bold
        p.font.color.rgb = rgb(palette[color])
        return box

    pattern = re.compile("|".join(re.escape(term) for term in terms))
    pos = 0
    for match in pattern.finditer(clean):
        if match.start() > pos:
            run = p.add_run()
            run.text = clean[pos : match.start()]
            run.font.name = font
            run.font.size = pt(size)
            run.font.bold = bold
            run.font.color.rgb = rgb(palette[color])
        run = p.add_run()
        run.text = match.group(0)
        run.font.name = font
        run.font.size = pt(size)
        run.font.bold = True
        run.font.color.rgb = rgb(palette[emphasis_color])
        pos = match.end()
    if pos < len(clean):
        run = p.add_run()
        run.text = clean[pos:]
        run.font.name = font
        run.font.size = pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(palette[color])
    return box


def add_fill(slide, x, y, w, h, color="paper", line="line", colors: dict[str, str] | None = None):
    palette = colors or DEFAULT_COLORS
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(palette[color])
    shape.line.color.rgb = rgb(palette[line])
    shape.line.width = Pt(0.6)
    return shape


def add_title(slide, page: dict[str, Any], colors: dict[str, str] | None = None):
    palette = colors or DEFAULT_COLORS
    display_page = int(page.get("display_page", page["page"]))
    add_textbox(slide, 0.42, 0.20, 11.85, 0.50, page["title"], 24, "ink", True, colors=palette)
    add_textbox(slide, 0.43, 0.68, 7.0, 0.24, page["chapter"], 9.5, "accent", colors=palette)
    add_textbox(
        slide,
        11.65,
        0.66,
        0.7,
        0.20,
        f"{display_page:02d}",
        9,
        "muted",
        align=PP_ALIGN.RIGHT,
        colors=palette,
    )
    line = slide.shapes.add_shape(1, Inches(0.42), Inches(0.88), Inches(11.84), Inches(0.015))
    line.fill.solid()
    line.fill.fore_color.rgb = rgb(palette["accent"])
    line.line.fill.background()


def add_footer(slide, source_note: str, colors: dict[str, str] | None = None):
    # Source notes remain in backend contracts and speaker notes; visible slides omit the bottom-left footer.
    return None


def add_table(
    slide,
    entry: dict[str, Any],
    x,
    y,
    w,
    h,
    max_rows=8,
    max_cols=5,
    colors: dict[str, str] | None = None,
):
    palette = colors or DEFAULT_COLORS
    rows = entry.get("rows") or []
    if not rows:
        add_textbox(slide, x, y, w, h, text_preview(entry.get("text", ""), 720), BODY_MIN_PT, "ink", colors=palette)
        return None
    cell_limit = max(24, min(62, int((w / max(1, max_cols)) * (h / max(1, max_rows)) * 22)))
    clean = [[text_preview(str(cell), cell_limit) for cell in row[:max_cols]] for row in rows[:max_rows]]
    col_count = max(len(row) for row in clean) if clean else 1
    for row in clean:
        while len(row) < col_count:
            row.append("")
    shape = slide.shapes.add_table(len(clean), col_count, Inches(x), Inches(y), Inches(w), Inches(h))
    table = shape.table
    row_height = Inches(h / max(1, len(clean)))
    for row in table.rows:
        row.height = row_height
    col_width = Inches(w / max(1, col_count))
    for column in table.columns:
        column.width = col_width
    for i, row in enumerate(clean):
        for j, cell_text in enumerate(row):
            cell = table.cell(i, j)
            cell.text = cell_text
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.04)
            cell.margin_right = Inches(0.04)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                paragraph.font.name = "Microsoft YaHei"
                paragraph.font.size = pt(TABLE_MIN_PT)
                paragraph.font.color.rgb = rgb(palette["ink"])
                if i == 0:
                    paragraph.font.bold = True
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb(palette["secondary"] if i == 0 else "FFFFFF")
    return shape


def add_image(slide, path: Path, x, y, w, h, colors: dict[str, str] | None = None):
    try:
        with Image.open(path) as im:
            iw, ih = im.size
    except Exception:
        add_textbox(slide, x, y, w, h, "报告图件暂无法嵌入", 14, "warn", colors=colors)
        return
    ratio = iw / ih
    box_ratio = w / h
    if ratio > box_ratio:
        final_w = w
        final_h = w / ratio
        final_x = x
        final_y = y + (h - final_h) / 2
    else:
        final_h = h
        final_w = h * ratio
        final_x = x + (w - final_w) / 2
        final_y = y
    slide.shapes.add_picture(
        str(path),
        Inches(final_x),
        Inches(final_y),
        width=Inches(final_w),
        height=Inches(final_h),
    )


def add_image_panel(
    slide,
    paths: list[Path],
    x,
    y,
    w,
    h,
    colors: dict[str, str] | None = None,
    allow_placeholder: bool = False,
):
    palette = colors or DEFAULT_COLORS
    if not paths:
        if not allow_placeholder:
            raise ValueError("Required source-figure panel has no image paths")
        add_fill(slide, x, y, w, h, "soft", colors=palette)
        add_textbox(
            slide,
            x + 0.25,
            y + 0.25,
            w - 0.5,
            h - 0.5,
            "报告相关图件",
            18,
            "muted",
            True,
            align=PP_ALIGN.CENTER,
            colors=palette,
        )
        return
    n = min(len(paths), 4)
    paths = paths[:n]
    if n == 1:
        add_image(slide, paths[0], x, y, w, h, colors=palette)
        return

    layout = choose_figure_layout(paths, w, h)
    cols = int(layout["cols"])
    rows = int(layout["rows"])
    gap = float(layout["gap"])
    cell_w = (w - gap * (cols - 1)) / cols
    cell_h = (h - gap * (rows - 1)) / rows
    for idx, path in enumerate(paths):
        cx = x + (idx % cols) * (cell_w + gap)
        cy = y + (idx // cols) * (cell_h + gap)
        add_image(slide, path, cx + 0.03, cy + 0.03, cell_w - 0.06, cell_h - 0.06, colors=palette)


def add_bullets(
    slide,
    x,
    y,
    w,
    h,
    items: list[str],
    size=BODY_MIN_PT,
    colors: dict[str, str] | None = None,
):
    palette = colors or DEFAULT_COLORS
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    for idx, item in enumerate(items):
        paragraph = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        paragraph.text = sanitize_visible_text(item)
        paragraph.level = 0
        paragraph.font.name = "Microsoft YaHei"
        paragraph.font.size = pt(max(float(size), BODY_MIN_PT))
        paragraph.font.color.rgb = rgb(palette["ink"])
        paragraph.space_after = pt(4)
    return box


def add_numbered_points(
    slide,
    x,
    y,
    w,
    h,
    items: list[str | dict[str, Any]],
    size=14,
    colors: dict[str, str] | None = None,
    emphasis_terms: list[str] | None = None,
):
    """Render multi-paragraph source text as numbered, named, editable blocks."""

    palette = colors or DEFAULT_COLORS
    points = structured_points(items, max_items=5)
    if not points:
        return []
    gap = 0.10
    row_h = (h - gap * (len(points) - 1)) / len(points)
    shapes = []
    for idx, point in enumerate(points):
        row_y = y + idx * (row_h + gap)
        shapes.append(add_fill(slide, x, row_y, w, row_h, "paper", "line", colors=palette))
        shapes.append(add_fill(slide, x + 0.12, row_y + 0.12, 0.46, 0.34, "accent", "accent", colors=palette))
        shapes.append(
            add_textbox(
                slide,
                x + 0.16,
                row_y + 0.17,
                0.38,
                0.18,
                point["number"],
                10.5,
                "paper",
                True,
                align=PP_ALIGN.CENTER,
                colors=palette,
            )
        )
        shapes.append(
            add_textbox(
                slide,
                x + 0.72,
                row_y + 0.08,
                w - 0.86,
                0.22,
                point["title"],
                max(size, 14),
                "ink",
                True,
                colors=palette,
            )
        )
        terms = emphasis_terms or emphasis_terms_from_text(point["text"])
        shapes.append(
            add_emphasis_textbox(
                slide,
                x + 0.72,
                row_y + 0.36,
                w - 0.88,
                max(0.22, row_h - 0.44),
                point["text"],
                max(size - 1, BODY_MIN_PT),
                "muted",
                emphasis_terms=terms,
                colors=palette,
            )
        )
    return shapes


def add_summary_dashboard_slide(
    slide,
    title: str,
    subtitle: str,
    sections: list[dict[str, str]],
    facts: list[str],
    closing: str = "",
    colors: dict[str, str] | None = None,
) -> None:
    """Create one of the allowed high-level summary/overall pages."""

    palette = colors or DEFAULT_COLORS
    add_fill(slide, 0, 0, 13.333, 7.5, "paper", "paper", colors=palette)
    add_textbox(slide, 0.62, 0.48, 7.6, 0.52, title, 28, "primary", True, colors=palette)
    add_emphasis_textbox(slide, 0.65, 1.12, 8.8, 0.38, subtitle, 14, "muted", colors=palette)
    add_fill(slide, 0.64, 1.68, 12.0, 0.02, "accent", "accent", colors=palette)

    for index, section in enumerate(sections[:4], start=1):
        col = (index - 1) % 2
        row = (index - 1) // 2
        card_x = 0.72 + col * 6.1
        card_y = 1.98 + row * 1.55
        add_fill(slide, card_x, card_y, 5.72, 1.24, "soft", "line", colors=palette)
        add_textbox(slide, card_x + 0.18, card_y + 0.14, 0.44, 0.26, f"{index:02d}", 12, "accent", True, colors=palette)
        add_textbox(slide, card_x + 0.70, card_y + 0.11, 4.76, 0.30, section.get("title", ""), 15, "ink", True, colors=palette)
        add_emphasis_textbox(
            slide,
            card_x + 0.20,
            card_y + 0.52,
            5.25,
            0.42,
            section.get("text", ""),
            11.5,
            "muted",
            colors=palette,
        )

    add_textbox(slide, 0.76, 5.28, 1.6, 0.26, "关键数据", 14, "primary", True, colors=palette)
    for index, fact in enumerate(facts[:4]):
        x0 = 2.18 + index * 2.58
        add_fill(slide, x0, 5.16, 2.25, 0.82, "paper", "line", colors=palette)
        add_emphasis_textbox(slide, x0 + 0.14, 5.34, 1.96, 0.26, fact, 13, "ink", colors=palette)
    if closing:
        add_fill(slide, 0.76, 6.35, 11.78, 0.48, "secondary", "secondary", colors=palette)
        add_emphasis_textbox(slide, 0.98, 6.49, 11.2, 0.22, closing, 12.5, "primary", colors=palette)


def add_agenda_slide(
    slide,
    plan: dict[str, Any],
    page: dict[str, Any],
    colors: dict[str, str] | None = None,
) -> None:
    palette = colors or DEFAULT_COLORS
    display_page = int(page.get("display_page", page["page"]))
    add_fill(slide, 0, 0, 13.333, 7.5, "paper", "paper", colors=palette)
    add_textbox(slide, 0.75, 0.62, 5.8, 0.62, "目录", 30, "primary", True, colors=palette)
    add_textbox(slide, 0.78, 1.30, 4.0, 0.25, "CONTENTS", 10, "accent", colors=palette)
    chapters = list(plan["coverage"]["chapter_allocation"].keys())
    for idx, chapter in enumerate(chapters, start=1):
        col = 0 if idx <= 5 else 1
        row = idx - 1 if idx <= 5 else idx - 6
        x = 0.85 + col * 6.0
        y = 1.75 + row * 0.82
        add_textbox(slide, x, y, 0.42, 0.28, f"{idx:02d}", 14, "accent", True, colors=palette)
        add_textbox(slide, x + 0.55, y - 0.03, 4.95, 0.36, chapter, 14, "ink", True, colors=palette)
    add_textbox(
        slide,
        0.8,
        6.75,
        9.8,
        0.30,
        "汇报按报告章节展开，依次呈现项目背景、工程概况、评价预测、措施体系、投资效益与管理验收。",
        10.5,
        "muted",
        colors=palette,
    )
    add_textbox(
        slide,
        11.65,
        6.82,
        0.7,
        0.2,
        f"{display_page:02d}",
        9,
        "muted",
        align=PP_ALIGN.RIGHT,
        colors=palette,
    )


def chapter_order_from_plan(plan: dict[str, Any]) -> list[str]:
    """Return substantive report chapters in the intended deck order."""

    structural = {"cover", "agenda", "section", "closing"}
    slide_chapters: list[str] = []

    def add_once(chapter: str) -> None:
        chapter = sanitize_visible_text(chapter)
        if chapter and chapter not in slide_chapters:
            slide_chapters.append(chapter)

    for chapter in plan.get("deck", {}).get("chapter_order", []) or []:
        add_once(str(chapter))
    for page in plan.get("slides", []):
        slide_type = str(page.get("type", "")).lower()
        if slide_type not in structural:
            add_once(str(page.get("chapter", "")))
    return slide_chapters


def section_divider_count(plan: dict[str, Any]) -> int:
    return len(chapter_order_from_plan(plan))


def engineering_page_count(plan: dict[str, Any]) -> int:
    return len(plan.get("slides", [])) + section_divider_count(plan)


def add_section_divider_slide(
    slide,
    chapter_index: int,
    chapter: str,
    total_chapters: int,
    subtitle: str | None = None,
    colors: dict[str, str] | None = None,
) -> None:
    palette = colors or DEFAULT_COLORS
    add_fill(slide, 0, 0, 13.333, 7.5, "primary", "primary", colors=palette)
    add_fill(slide, 0.72, 1.05, 0.08, 5.18, "accent", "accent", colors=palette)
    add_textbox(
        slide,
        0.82,
        0.70,
        4.6,
        0.35,
        f"第{chapter_index:02d}章 / 共{total_chapters:02d}章",
        16,
        "secondary",
        True,
        colors=palette,
    )
    add_textbox(slide, 0.82, 2.08, 1.42, 0.88, f"{chapter_index:02d}", 54, "paper", True, colors=palette)
    add_textbox(slide, 2.38, 2.20, 9.8, 0.82, chapter, 30, "paper", True, colors=palette)
    add_textbox(
        slide,
        2.42,
        3.24,
        9.0,
        0.36,
        subtitle or f"本章围绕{chapter}展开，提炼关键源证据、计算口径和评审判断。",
        15,
        "secondary",
        colors=palette,
    )
    add_fill(slide, 2.42, 4.10, 7.0, 0.02, "accent", "accent", colors=palette)


def create_engineering_design_spec(
    plan: dict[str, Any],
    media_dir: Path,
    template_profile: dict[str, Any],
    project_title: str,
    audience: str,
    use_case: str,
    design_style: str,
    image_names: list[str] | None = None,
) -> str:
    allocation = plan["coverage"]["chapter_allocation"]
    lines = [
        f"# {plan.get('project', 'Engineering Project')} - Design Spec",
        "",
        "## I. Project Information",
        "",
        "| Item | Value |",
        "| ---- | ----- |",
        f"| **Project Name** | {project_title} |",
        "| **Canvas Format** | PPT 16:9 (1280x720) |",
        f"| **Page Count** | {engineering_page_count(plan)} |",
        f"| **Design Style** | {design_style} |",
        f"| **Target Audience** | {audience} |",
        f"| **Use Case** | {use_case} |",
        f"| **Created Date** | {datetime.now().strftime('%Y-%m-%d')} |",
        "",
        "## II. Canvas Specification",
        "",
        "| Property | Value |",
        "| -------- | ----- |",
        "| **Format** | PPT 16:9 |",
        "| **Dimensions** | 1280x720 |",
        "| **viewBox** | `0 0 1280 720` |",
        "| **Margins** | left/right 40-60px, top 30px, footer 24px |",
        "| **Content Area** | 1200x600px |",
        "",
        "## III. Visual Theme",
        "",
        "- **Style**: source-faithful Chinese engineering review presentation",
        "- **Theme**: light technical pages with deep-blue section rhythm",
        "- **Tone**: evidence-first, technical, restrained",
        "- **Chapter rhythm**: each report chapter starts with a section divider slide carrying the chapter number and name.",
        "",
        "| Role | HEX | Purpose |",
        "| ---- | --- | ------- |",
        "| **Background** | `#F7F9FC` | Page background |",
        "| **Secondary bg** | `#F1F6FA` | Technical bands and table surfaces |",
        "| **Primary** | `#01203C` | Chapter headers and emphasis |",
        "| **Accent** | `#4489C8` | Section labels and engineering highlights |",
        "| **Warning** | `#A33A2B` | High-risk values and limitations |",
        "| **Body text** | `#01203C` | Main body text |",
        "| **Secondary text** | `#566B7F` | Captions, source notes |",
        "| **Border/divider** | `#CCD8E2` | Table and figure frames |",
        "",
        "## IV. Typography System",
        "",
        "**Typography direction**: CJK-primary professional sans with serif fallback for report quotations.",
        "",
        "| Role | Chinese | English | Fallback tail |",
        "| ---- | ------- | ------- | ------------- |",
        '| **Title** | `"Microsoft YaHei"` | `Arial` | `sans-serif` |',
        '| **Body** | `"Microsoft YaHei"` | `Arial` | `sans-serif` |',
        "| **Emphasis** | `SimSun` | `Georgia` | `serif` |",
        '| **Code** | — | `Consolas, "Courier New"` | `monospace` |',
        "",
        "**Baseline**: Body font size = 18px.",
        "",
        "## V. Layout Principles",
        "",
        "- Header area: title + chapter + page number.",
        "- Content area prioritizes original figure/table readability.",
        "- Visible slides do not show bottom-left source footers; source notes stay in backend records and speaker notes.",
        "- Dense technical pages avoid decorative cards; tables and figures use framed source panels.",
        "- Visible panel headings must be content-specific, such as `图件重点：矿区强径流带分布` or `表格重点：涌水量预测结果`.",
        "- Generic workflow labels such as `报告对图件的说明`, `报告对表格的说明`, `报告计算口径`, and `报告阐述` are forbidden in visible slides.",
        "- Report section headings, figure captions, table captions and corresponding subsection titles take priority over agent-generated topic labels.",
        "- Visible ellipses are forbidden; long report text must be shortened to complete, source-faithful sentences or split across slides.",
        "- Horizontal report maps/profiles/charts should use top-figure/bottom-text layouts when side-by-side placement would make the figure too small.",
        "- Internal planning metadata, evidence IDs, asset filenames and row/column diagnostics remain only in backend contracts/QA.",
        "",
        "## VI. Icon Usage Specification",
        "",
        "- Icons are sparse; source figures and tables carry most visual meaning.",
        "",
        "## VII. Visualization Reference List",
        "",
        "The deck uses report-native tables, maps, figures and formulas because source fidelity is more important than restyling source data.",
        "",
    ]
    if template_profile.get("available"):
        lines.extend(
            [
                "## VII-A. Template Learning Profile",
                "",
                f"- Template file: `{template_profile.get('path')}`",
                f"- Template slides: {template_profile.get('slide_count')}",
                f"- Canvas inches: {template_profile.get('canvas_inches')}",
                f"- Dominant fonts: {template_profile.get('dominant_fonts')}",
                f"- Dominant text colors: {template_profile.get('dominant_text_colors')}",
                f"- Picture distribution: {template_profile.get('picture_count_distribution')}",
                f"- Layout keys: {template_profile.get('layout_keys')}",
                f"- Adopted rules: {'; '.join(template_profile.get('adopted_rules', []))}",
                "",
            ]
        )
    lines.extend(
        [
            "## VIII. Image Resource List",
            "",
            "| Filename | Dimensions | Ratio | Purpose | Type | Layout pattern | Acquire Via | Status | Reference | text_policy | page_role |",
            "| -------- | ---------- | ----- | ------- | ---- | -------------- | ----------- | ------ | --------- | ----------- | --------- |",
        ]
    )
    names = image_names or []
    for name in names:
        path = media_dir / name
        if path.exists():
            with Image.open(path) as im:
                w, h = im.size
            lines.append(
                f"| {name} | {w}x{h} | {w / h:.2f} | Key report source figure | Source figure | source-faithful figure panel | user | Existing | Extracted from source | none | local |"
            )
    lines.extend(["", "## IX. Content Outline", ""])
    for chapter, count in allocation.items():
        lines.append(f"### {chapter} ({count} pages)")
        for page in plan["slides"]:
            if page["chapter"] == chapter:
                lines.append(f"#### Slide {page['page']:02d} - {page['title']}")
                lines.append(f"- **Layout**: {page['layout_pattern']}")
                lines.append(f"- **Core message**: {page['title']}")
                lines.append(f"- **Source mode**: {page['source_mode']}")
                lines.append(f"- **Evidence**: {', '.join(page['evidence_ids'])}")
                lines.append(f"- **Source note**: {page['source_note']}")
        lines.append("")
    lines.extend(
        [
            "## X. Speaker Notes Requirements",
            "",
            "- One notes file per slide under `notes/`.",
            "- Notes use formal engineering review language and cite source notes.",
            "",
            "## XI. Technical Constraints Reminder",
            "",
            "- Final PPTX must remain natively editable.",
            "- Keep body text readable; do not shrink source tables below readable floor.",
            "- Source conflicts and limitations remain visible when relevant.",
        ]
    )
    return "\n".join(lines) + "\n"


def create_engineering_spec_lock(
    plan: dict[str, Any],
    media_dir: Path,
    important_anchor_pages: set[int] | None = None,
) -> str:
    anchor_pages = important_anchor_pages or {1, 2, len(plan.get("slides", []))}
    lines = [
        "## canvas",
        "- viewBox: 0 0 1280 720",
        "- format: PPT 16:9",
        "",
        "## colors",
        "- bg: #F7F9FC",
        "- paper: #FFFFFF",
        "- primary: #01203C",
        "- accent: #4489C8",
        "- warning: #A33A2B",
        "- secondary_accent: #EAF2F8",
        "- text: #01203C",
        "- text_secondary: #566B7F",
        "- border: #CCD8E2",
        "",
        "## typography",
        '- font_family: "Microsoft YaHei", Arial, sans-serif',
        '- title_family: "Microsoft YaHei", Arial, sans-serif',
        '- body_family: "Microsoft YaHei", Arial, sans-serif',
        "- emphasis_family: Georgia, SimSun, serif",
        '- code_family: Consolas, "Courier New", monospace',
        "- body: 18",
        "- title: 32",
        "- subtitle: 24",
        "- annotation: 14",
        "- footer: 10",
        "",
        "## images",
    ]
    for img in sorted(media_dir.glob("image_*.*")):
        if img.suffix.lower() in {".png", ".jpg", ".jpeg"}:
            lines.append(f"- {img.stem}: {img.as_posix()} | no-crop")
    lines.extend(["", "## page_rhythm"])
    for page in plan["slides"]:
        if page["page"] in anchor_pages:
            rhythm = "anchor"
        elif page["density"] == "dense":
            rhythm = "dense"
        else:
            rhythm = "breathing" if page["type"] in {"conclusion", "chapter_summary"} else "dense"
        lines.append(f"- P{page['page']:02d}: {rhythm}")
    lines.extend(
        [
            "",
            "## forbidden",
            "- rgba()",
            "- `<style>`, `class`, `<foreignObject>`, `textPath`, `@font-face`, `<animate*>`, `<script>`, `<iframe>`",
            "- `<g opacity>`",
            "- HTML named entities in visible text",
            "- Visible ellipses (`...`, `…`, `……`) as incomplete wording",
            "- Visible internal metadata labels such as 来源模式, 证据编号, 本页用于, 原始对象, 必讲内容, 保留理由, 原表定位, image_*.png, or E-* evidence IDs",
        ]
    )
    return "\n".join(lines) + "\n"
