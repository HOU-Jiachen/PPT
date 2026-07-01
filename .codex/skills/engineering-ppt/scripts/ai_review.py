#!/usr/bin/env python3
"""AI-assisted final review gate for engineering PPTX deliverables."""

from __future__ import annotations

import argparse
import json
import os
import re
import sys
import urllib.error
import urllib.request
from datetime import datetime
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


DEFAULT_MODEL = "gpt-4.1"


def compact(text: str, limit: int = 220) -> str:
    value = re.sub(r"\s+", " ", text or "").strip()
    return value if len(value) <= limit else value[:limit].rstrip() + "..."


def shape_payload(shape, slide_width: int, slide_height: int) -> dict[str, Any]:
    payload: dict[str, Any] = {
        "name": str(getattr(shape, "name", "") or ""),
        "type": str(getattr(shape, "shape_type", "")),
        "bbox_ratio": [
            round(int(getattr(shape, "left", 0) or 0) / max(slide_width, 1), 3),
            round(int(getattr(shape, "top", 0) or 0) / max(slide_height, 1), 3),
            round(int(getattr(shape, "width", 0) or 0) / max(slide_width, 1), 3),
            round(int(getattr(shape, "height", 0) or 0) / max(slide_height, 1), 3),
        ],
    }
    if getattr(shape, "has_text_frame", False):
        paragraphs = []
        for paragraph in shape.text_frame.paragraphs:
            text = re.sub(r"\s+", " ", paragraph.text or "").strip()
            if not text:
                continue
            sizes = []
            emphasized = 0
            for run in paragraph.runs:
                if not run.text.strip():
                    continue
                if run.font.size is not None:
                    sizes.append(round(run.font.size.pt, 1))
                if run.font.bold:
                    emphasized += 1
                else:
                    try:
                        if run.font.color.rgb is not None:
                            emphasized += 1
                    except Exception:
                        pass
            paragraphs.append(
                {
                    "text": compact(text),
                    "font_sizes": sizes,
                    "emphasized_runs": emphasized,
                }
            )
        payload["text"] = paragraphs
    if getattr(shape, "has_table", False):
        rows = []
        min_font = None
        for row in shape.table.rows:
            cells = []
            for cell in row.cells:
                cells.append(compact(cell.text, 80))
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.text.strip() and run.font.size is not None:
                            value = float(run.font.size.pt)
                            min_font = value if min_font is None else min(min_font, value)
            rows.append(cells)
        payload["table"] = {
            "rows": len(shape.table.rows),
            "cols": len(shape.table.columns),
            "sample_rows": rows[:4],
            "min_font_size": min_font,
        }
    if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
        payload["picture"] = True
    return payload


def pptx_payload(project: Path, pptx: Path) -> dict[str, Any]:
    prs = Presentation(pptx)
    slide_width = int(prs.slide_width)
    slide_height = int(prs.slide_height)
    slides = []
    for index, slide in enumerate(prs.slides, start=1):
        shapes = [shape_payload(shape, slide_width, slide_height) for shape in slide.shapes]
        visible_text = "\n".join(
            paragraph["text"]
            for shape in shapes
            for paragraph in shape.get("text", [])
            if paragraph.get("text")
        )
        slides.append(
            {
                "page": index,
                "visible_text": compact(visible_text, 900),
                "shape_count": len(shapes),
                "shapes": shapes,
            }
        )
    release_audit_path = project / "qa" / "release_audit.json"
    release_audit = {}
    if release_audit_path.exists():
        release_audit = json.loads(release_audit_path.read_text(encoding="utf-8-sig"))
    return {
        "project": project.name,
        "pptx": str(pptx),
        "created_at": datetime.now().isoformat(timespec="seconds"),
        "slide_count": len(prs.slides),
        "release_audit": release_audit.get("summary", {}),
        "release_findings": release_audit.get("findings", [])[:80],
        "slides": slides,
        "review_rules": [
            "Visible content must be source-faithful Chinese engineering report language.",
            "Do not allow internal IDs such as T036, T-036, evidence IDs, filenames, source_mode, or generation-process wording.",
            "Do not allow generic process sentences such as key values remain checkable against source tables.",
            "Tables must be readable; native table text below 10 pt or source-table crops marked TableImageNeedsSplit are blocking.",
            "Chapter divider names and chapter-introduction text must not use colored emphasis; emphasis is limited to substantive body content.",
            "Flag overlap, cramped layout, excessive density, incomplete sentences, poor table formatting, and missing evidence context.",
        ],
    }


def review_prompt(payload: dict[str, Any]) -> str:
    return (
        "You are a senior Chinese engineering-PPT reviewer. Review the following PPTX extraction. "
        "Return only JSON with keys passed:boolean, findings:[{severity,page,category,issue,fix}], "
        "and summary. Severity must be error or warning. Treat formatting, content faithfulness, "
        "layout, table readability, forbidden internal text, and chapter-emphasis misuse as review scope.\n\n"
        + json.dumps(payload, ensure_ascii=False, indent=2)
    )


def call_openai(prompt: str, model: str) -> str:
    api_key = os.environ.get("OPENAI_API_KEY")
    if not api_key:
        raise RuntimeError("OPENAI_API_KEY is not set.")
    body = {
        "model": model,
        "input": [
            {
                "role": "user",
                "content": [{"type": "input_text", "text": prompt}],
            }
        ],
        "temperature": 0,
    }
    request = urllib.request.Request(
        "https://api.openai.com/v1/responses",
        data=json.dumps(body).encode("utf-8"),
        headers={
            "Authorization": f"Bearer {api_key}",
            "Content-Type": "application/json",
        },
        method="POST",
    )
    try:
        with urllib.request.urlopen(request, timeout=120) as response:
            payload = json.loads(response.read().decode("utf-8"))
    except urllib.error.HTTPError as exc:
        detail = exc.read().decode("utf-8", errors="replace")
        raise RuntimeError(f"OpenAI review request failed: HTTP {exc.code}: {detail}") from exc
    return extract_response_text(payload)


def extract_response_text(payload: dict[str, Any]) -> str:
    if payload.get("output_text"):
        return str(payload["output_text"])
    parts: list[str] = []
    for item in payload.get("output", []) or []:
        for content in item.get("content", []) or []:
            text = content.get("text")
            if text:
                parts.append(str(text))
    return "\n".join(parts).strip()


def parse_review(text: str) -> dict[str, Any]:
    value = text.strip()
    if value.startswith("```"):
        value = re.sub(r"^```(?:json)?\s*", "", value)
        value = re.sub(r"\s*```$", "", value)
    try:
        return json.loads(value)
    except json.JSONDecodeError:
        return {"passed": False, "findings": [{"severity": "error", "category": "ai-review", "issue": "AI review did not return valid JSON.", "fix": "Inspect qa/ai_review_raw.txt and rerun the review."}], "raw": text}


def write_markdown(project: Path, review: dict[str, Any], blocked_reason: str | None = None) -> None:
    lines = ["# AI Review", ""]
    if blocked_reason:
        lines.append(f"- Status: blocked")
        lines.append(f"- Reason: {blocked_reason}")
    else:
        lines.append(f"- Passed: {bool(review.get('passed'))}")
    findings = review.get("findings", []) or []
    lines.append(f"- Findings: {len(findings)}")
    lines.append("")
    for item in findings:
        lines.append(
            f"- **{str(item.get('severity', 'warning')).upper()}** P{item.get('page', '-')}: "
            f"{item.get('category', 'review')} - {item.get('issue', '')} Fix: {item.get('fix', '')}"
        )
    lines.append("")
    (project / "qa" / "ai_review.md").write_text("\n".join(lines), encoding="utf-8")


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("project", type=Path)
    parser.add_argument("--pptx", type=Path, required=True)
    parser.add_argument("--model", default=os.environ.get("OPENAI_REVIEW_MODEL", DEFAULT_MODEL))
    parser.add_argument("--no-fail", action="store_true")
    args = parser.parse_args()

    project = args.project.resolve()
    pptx = args.pptx.resolve()
    qa = project / "qa"
    qa.mkdir(parents=True, exist_ok=True)
    payload = pptx_payload(project, pptx)
    prompt = review_prompt(payload)
    (qa / "ai_review_payload.json").write_text(json.dumps(payload, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    (qa / "ai_review_prompt.md").write_text(prompt, encoding="utf-8")

    try:
        raw = call_openai(prompt, args.model)
    except Exception as exc:
        blocked = str(exc)
        write_markdown(project, {"passed": False, "findings": []}, blocked)
        print(json.dumps({"passed": False, "blocked": True, "reason": blocked, "prompt": str(qa / "ai_review_prompt.md")}, ensure_ascii=False))
        if args.no_fail:
            return
        sys.exit(2)

    (qa / "ai_review_raw.txt").write_text(raw, encoding="utf-8")
    review = parse_review(raw)
    (qa / "ai_review.json").write_text(json.dumps(review, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    write_markdown(project, review)
    errors = [item for item in review.get("findings", []) or [] if str(item.get("severity", "")).lower() == "error"]
    print(json.dumps({"passed": bool(review.get("passed")) and not errors, "errors": len(errors), "findings": len(review.get("findings", []) or [])}, ensure_ascii=False))
    if errors and not args.no_fail:
        sys.exit(1)


if __name__ == "__main__":
    main()
